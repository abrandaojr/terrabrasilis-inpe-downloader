"""
04_generate_presentation.py
===========================
Generates a bilingual press PowerPoint on Brazilian deforestation data (PRODES/INPE).

All PRODES statistics (area by year, by biome) are calculated ON THE FLY from
the GeoParquet files produced by 02_convert_to_geoparquet.py.
Run script 02 before this one to generate the source data.

Reference data that is NOT in the PRODES pipeline
(forest cover % from MapBiomas, international comparison from GFW/FAO)
is kept as curated constants with explicit source labels.

Structure
---------
Slides 1–10  : PT-BR (Português Brasileiro)
Slides 11–20 : EN-US (American English)

Each section covers:
  1  Cover
  2  Lead stat  (% drop calculated from real data)
  3  Amazon historical series
  4  By biome
  5  Cerrado spotlight
  6  Forest cover remaining (MapBiomas reference)
  7  Target trajectory
  8  International context (GFW/FAO reference)
  9  What explains the decline
  10 Key takeaways

Usage
-----
    python 04_generate_presentation.py

Output
------
    PRODES_Press_Briefing.pptx  (same directory as this script)

Author
------
Amintas Brandão Jr. <abrandaojr@gmail.com>
Imazon — Instituto do Homem e Meio Ambiente da Amazônia

License
-------
MIT
"""

from __future__ import annotations

__version__ = "2.0.0"
__all__: list[str] = []

import importlib.util
import io
import re
import subprocess
import sys
from datetime import datetime
from pathlib import Path

HERE = Path(__file__).parent   # script directory — used for output paths

# ---------------------------------------------------------------------------
# Dependency bootstrap
# ---------------------------------------------------------------------------

def _bootstrap(*packages: tuple[str, str]) -> None:
    """Install missing packages into the current Python environment."""
    import importlib
    import shutil

    mod_by_pip = {pip: mod for pip, mod in packages}

    def _still_missing(pkgs: list[str]) -> list[str]:
        importlib.invalidate_caches()
        return [p for p in pkgs if not importlib.util.find_spec(mod_by_pip[p])]

    missing = _still_missing(list(mod_by_pip))
    if not missing:
        return

    if not shutil.which("uv"):
        subprocess.call(
            [sys.executable, "-m", "pip", "install", "--quiet", "uv"],
            stderr=subprocess.DEVNULL,
        )

    strategies = [
        [sys.executable, "-m", "pip", "install", "--quiet"],
        ["uv", "pip", "install", "--python", sys.executable, "--quiet"],
        [sys.executable, "-m", "uv", "pip", "install", "--python", sys.executable, "--quiet"],
        [sys.executable, "-m", "pip", "install", "--quiet", "--break-system-packages"],
    ]
    for base in strategies:
        if not missing:
            return
        try:
            subprocess.check_call(base + missing, stderr=subprocess.DEVNULL)
            missing = _still_missing(missing)
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass

    if missing:
        sys.exit(f"[FATAL] Could not install: {' '.join(missing)}")


_bootstrap(
    ("matplotlib",  "matplotlib"),
    ("numpy",       "numpy"),
    ("python-pptx", "pptx"),
    ("duckdb",      "duckdb"),
    ("pyarrow",     "pyarrow"),
)

import matplotlib.pyplot as plt       # noqa: E402
import numpy as np                    # noqa: E402
from pptx import Presentation         # noqa: E402
from pptx.util import Pt              # noqa: E402
from pptx.dml.color import RGBColor   # noqa: E402

# ---------------------------------------------------------------------------
# CONFIG  ← the only section that needs to be edited
# ---------------------------------------------------------------------------

CONFIG: dict[str, object] = {
    "output_path":    HERE / "PRODES_Press_Briefing.pptx",
    "chart_dpi":      220,
    # Directory with GeoParquet files from script 02.
    # Leave as None to auto-detect from the standard location.
    "geoparquet_dir": None,   # e.g. r"C:\Amintas\Prodes\geoparquet"
    # Policy targets (km²) — not data, not queried from GeoParquet.
    "target_2026_km2": 4_866,
    "target_2028_km2": 4_000,
}

SEP = "=" * 65
DIV = "-" * 65

# ---------------------------------------------------------------------------
# REFERENCE DATA  (sources outside the PRODES pipeline — kept as constants)
# ---------------------------------------------------------------------------

# Remaining native vegetation (% of original biome area) — MapBiomas 2023
# Source: https://mapbiomas.org
_COVER_PCT: dict[str, dict[str, float]] = {
    "pt": {
        "Pantanal":        86.0,
        "Amazônia":        81.0,
        "Caatinga":        60.0,
        "Cerrado":         53.0,
        "Pampa":           38.0,
        "Mata Atlântica":  13.0,
    },
    "en": {
        "Pantanal":         86.0,
        "Amazon":           81.0,
        "Caatinga":         60.0,
        "Cerrado Savanna":  53.0,
        "Pampa":            38.0,
        "Atlantic Forest":  13.0,
    },
}

# International tropical primary forest loss 2023 (km²) — GFW / FAO
# Different methodology from PRODES (tree-cover loss, not deforestation).
_INTL_LOSS_KM2: dict[str, int] = {
    "Brazil":      9_064,   # PRODES figure for comparability
    "D.R. Congo":  4_900,
    "Bolivia":     4_200,
    "Indonesia":   2_800,
    "Colombia":    1_450,
}

# ---------------------------------------------------------------------------
# STATS  — populated from GeoParquet in main(), never hardcoded
# ---------------------------------------------------------------------------

# All values computed on-the-fly from actual GeoParquet data.
STATS: dict = {}

# Biome directory name (from TerraBrasilis ZIP structure) → display names
_BIOME_TO_PT: dict[str, str] = {
    "Amazon Biome":    "Amazônia Legal",
    "Legal Amazon":    "Amazônia Legal",
    "Cerrado":         "Cerrado",
    "Caatinga":        "Caatinga",
    "Pantanal":        "Pantanal",
    "Mata Atlantica":  "Mata Atlântica",
    "Pampa":           "Pampa",
}
_BIOME_TO_EN: dict[str, str] = {
    "Amazon Biome":    "Legal Amazon",
    "Legal Amazon":    "Legal Amazon",
    "Cerrado":         "Cerrado Savanna",
    "Caatinga":        "Caatinga",
    "Pantanal":        "Pantanal",
    "Mata Atlantica":  "Atlantic Forest",
    "Pampa":           "Pampa",
}
_AMAZON_DIRS = {"Amazon Biome", "Legal Amazon"}
# Keywords that identify yearly deforestation categories
_DEFOR_KEYWORDS = ("yearly", "defor", "anual", "desmat")
# Candidate column names (area in km², checked in priority order)
_AREA_COLS = ("areakm", "area_km", "area_km2", "areakm2", "area")
# Candidate column names for year
_YEAR_COLS  = ("year", "ano", "yr", "data_year")

# ---------------------------------------------------------------------------
# DATA LOADING LAYER
# ---------------------------------------------------------------------------

def _auto_geoparquet_dir() -> Path | None:
    """Return the standard GeoParquet directory if it exists."""
    std = Path(r"C:\Amintas\Prodes\geoparquet")
    return std if std.exists() else None


def _detect_area_year(files: list[Path]) -> tuple[str | None, str | None]:
    """
    Inspect a sample parquet file's schema to find the area and year columns.
    Returns (area_col, year_col) — either may be None.
    """
    import pyarrow.parquet as pq

    for f in files[:5]:
        try:
            schema = pq.read_schema(str(f))
            low_to_orig = {n.lower(): n for n in schema.names}
            area_col = next((low_to_orig[c] for c in _AREA_COLS if c in low_to_orig), None)
            year_col = next((low_to_orig[c] for c in _YEAR_COLS if c in low_to_orig), None)
            if area_col is not None:
                return area_col, year_col
        except Exception:
            continue
    return None, None


def _infer_km2_factor(files: list[Path], area_col: str) -> float:
    """
    Read a few area values and guess the unit:
      km² → factor 1.0
      ha  → factor 0.01
      m²  → factor 0.000001
    """
    import pyarrow.parquet as pq

    sample: list[float] = []
    for f in files[:3]:
        try:
            col = pq.read_table(str(f), columns=[area_col]).column(area_col).to_pylist()
            sample += [v for v in col if v and v > 0][:20]
        except Exception:
            pass
    if not sample:
        return 1.0
    med = sorted(sample)[len(sample) // 2]
    if med > 500_000:    return 1 / 1_000_000   # m²
    if med > 5_000:      return 1 / 100          # ha
    return 1.0                                    # km²


def _query_series(files: list[Path], area_col: str, year_col: str, factor: float) -> dict[int, float]:
    """DuckDB: aggregate area by year, return {year: km²}."""
    import duckdb

    paths = [str(f).replace("\\", "/") for f in files]
    sql = f"""
        SELECT CAST("{year_col}" AS INTEGER)     AS yr,
               SUM(CAST("{area_col}" AS DOUBLE)) * {factor} AS km2
        FROM   read_parquet({paths!r})
        WHERE  "{year_col}"  IS NOT NULL
          AND  "{area_col}"  IS NOT NULL
          AND  CAST("{area_col}" AS DOUBLE) > 0
        GROUP  BY yr
        HAVING yr BETWEEN 2000 AND 2030
        ORDER  BY yr
    """
    try:
        rows = duckdb.connect().execute(sql).fetchall()
        return {int(r[0]): round(float(r[1]), 1) for r in rows if r[0] and r[1]}
    except Exception:
        return {}


def _query_total(files: list[Path], area_col: str, year_col: str | None,
                 target_year: int | None, factor: float) -> float | None:
    """DuckDB: sum area for a given year (or all years if year_col missing)."""
    import duckdb

    paths = [str(f).replace("\\", "/") for f in files]
    where = f'WHERE CAST("{year_col}" AS INTEGER) = {target_year}' if year_col and target_year else ""
    sql = f"""
        SELECT SUM(CAST("{area_col}" AS DOUBLE)) * {factor}
        FROM   read_parquet({paths!r})
        {where}
        HAVING SUM(CAST("{area_col}" AS DOUBLE)) > 0
    """
    try:
        row = duckdb.connect().execute(sql).fetchone()
        return round(float(row[0]), 1) if row and row[0] else None
    except Exception:
        return None


def _load_prodes_stats(geoparquet_dir: Path) -> dict:
    """
    Scan GeoParquet files and compute PRODES statistics on the fly.
    Returns a dict with keys: amazon_km2, biomes_km2_pt, biomes_km2_en, biome_year.
    Raises RuntimeError if no usable data found.
    """
    # Group parquet files by top-level biome directory + yearly category
    biome_files: dict[str, list[Path]] = {}
    for pf in sorted(geoparquet_dir.rglob("*.parquet")):
        try:
            parts = pf.relative_to(geoparquet_dir).parts
            if len(parts) < 2:
                continue
            biome_dir = parts[0]
            category  = parts[1] if len(parts) > 1 else ""
            if not any(k in category.lower() for k in _DEFOR_KEYWORDS):
                continue
            biome_files.setdefault(biome_dir, []).append(pf)
        except (ValueError, IndexError):
            continue

    if not biome_files:
        raise RuntimeError(
            "No yearly deforestation parquet files found in:\n"
            f"  {geoparquet_dir}\n"
            "Run  python 02_convert_to_geoparquet.py  first."
        )

    print(f"  [data] Biome folders found: {sorted(biome_files)}")

    result: dict = {
        "amazon_km2":    {},   # {year: km²}
        "biomes_km2_pt": {},   # {biome_name_pt: km²}
        "biomes_km2_en": {},   # {biome_name_en: km²}
        "biome_year":    {},   # {biome_dir: year used}
    }

    # ── Amazon annual series ─────────────────────────────────────────────
    amazon_files = [f for bd, fs in biome_files.items() if bd in _AMAZON_DIRS for f in fs]
    if amazon_files:
        area_col, year_col = _detect_area_year(amazon_files)
        if area_col and year_col:
            factor = _infer_km2_factor(amazon_files, area_col)
            series = _query_series(amazon_files, area_col, year_col, factor)
            if series:
                result["amazon_km2"] = series
                print(f"  [data] Amazon series: {min(series)}–{max(series)}  ({len(series)} yr)")
        else:
            print(f"  [data] WARN: could not detect columns in Amazon files (found: area_col={area_col})")

    # ── Per-biome total (most recent year with data) ─────────────────────
    ref_year = max(result["amazon_km2"]) if result["amazon_km2"] else None

    for biome_dir, files in biome_files.items():
        name_pt = _BIOME_TO_PT.get(biome_dir)
        name_en = _BIOME_TO_EN.get(biome_dir)
        if not name_pt:
            continue

        area_col, year_col = _detect_area_year(files)
        if not area_col:
            continue

        factor = _infer_km2_factor(files, area_col)

        # Find the most recent year available in this biome
        use_year = ref_year
        if year_col:
            try:
                import duckdb
                paths = [str(f).replace("\\", "/") for f in files]
                row = duckdb.connect().execute(
                    f'SELECT MAX(CAST("{year_col}" AS INTEGER)) FROM read_parquet({paths!r})'
                ).fetchone()
                if row and row[0]:
                    use_year = int(row[0])
            except Exception:
                pass

        total = _query_total(files, area_col, year_col, use_year, factor)
        if total and total > 0:
            result["biomes_km2_pt"][name_pt] = total
            result["biomes_km2_en"][name_en] = total
            result["biome_year"][biome_dir]   = use_year
            print(f"  [data] {biome_dir}: {total:,.0f} km² (year {use_year})")

    if not result["amazon_km2"] and not result["biomes_km2_pt"]:
        raise RuntimeError(
            "GeoParquet files found but no area/year data could be extracted.\n"
            "Check that column names include 'areakm' and 'year'."
        )

    return result


def _compute_derived_stats(raw: dict) -> dict:
    """Compute all derived statistics from raw query results."""
    s: dict = {}

    amazon = raw.get("amazon_km2", {})
    if amazon:
        s["amazon_km2"]    = amazon
        s["current_year"]  = max(amazon)
        s["current_km2"]   = amazon[s["current_year"]]
        s["peak_year"]     = max(amazon, key=lambda y: amazon[y])
        s["peak_km2"]      = amazon[s["peak_year"]]
        s["first_year"]    = min(amazon)
        s["first_km2"]     = amazon[s["first_year"]]
        s["n_years_decline"] = s["current_year"] - s["peak_year"]
        pct = (s["peak_km2"] - s["current_km2"]) / s["peak_km2"] * 100
        s["pct_decline"]   = round(pct, 1)
    else:
        s["amazon_km2"] = {}

    s["target_2026"]  = int(CONFIG["target_2026_km2"])
    s["target_2028"]  = int(CONFIG["target_2028_km2"])

    if "current_km2" in s and s["current_km2"] > 0:
        pct_to_target = (s["current_km2"] - s["target_2028"]) / s["current_km2"] * 100
        s["pct_to_target"] = round(pct_to_target, 1)
    else:
        s["pct_to_target"] = 0.0

    s["biomes_km2_pt"]   = raw.get("biomes_km2_pt", {})
    s["biomes_km2_en"]   = raw.get("biomes_km2_en", {})
    s["biome_year_by_dir"] = raw.get("biome_year", {})
    # Single display year: most recent year across all biomes
    by_dict = raw.get("biome_year", {})
    s["biome_year"] = max(by_dict.values()) if by_dict else (s.get("current_year", ""))

    # Cerrado-specific stats for spotlight slide
    cerrado_km2_pt = s["biomes_km2_pt"].get("Cerrado", 0)
    cerrado_km2_en = s["biomes_km2_en"].get("Cerrado Savanna", cerrado_km2_pt)
    s["cerrado_km2_pt"] = cerrado_km2_pt
    s["cerrado_km2_en"] = cerrado_km2_en

    # Reference data (not from GeoParquet)
    s["cover_pct"]  = _COVER_PCT
    s["intl_km2"]   = _INTL_LOSS_KM2

    # Update Brazil row in international comparison with real current year data
    if "current_km2" in s:
        s["intl_km2"] = {**_INTL_LOSS_KM2, "Brazil": int(round(s["current_km2"]))}

    return s


def _num(v: float, lang: str) -> str:
    """Format number: PT-BR uses . as thousands separator, EN uses ,"""
    formatted = f"{v:,.0f}"
    return formatted.replace(",", ".") if lang == "pt" else formatted

# ---------------------------------------------------------------------------
# STYLE
# ---------------------------------------------------------------------------

C_GREEN  = "#2E7D32"
C_RED    = "#C0392B"
C_BLUE   = "#1A5276"
C_ORANGE = "#E67E22"
C_GRAY   = "#CCCCCC"
C_DARK   = "#111111"
C_MED    = "#555555"
C_LIGHT  = "#999999"

_LANG_COLOR = {"pt": "#1B5E20", "en": "#0D47A1"}

_SW = 10.0
_SH = 5.625


def _in(x: float) -> int:
    return int(x * 914_400)


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

# ---------------------------------------------------------------------------
# BILINGUAL COPY  (templates filled at runtime from STATS)
# ---------------------------------------------------------------------------
# Keys beginning with "tmpl_" are format strings filled via _t(key, lang).
# Static keys are language dicts used directly.

COPY: dict[str, dict[str, str]] = {
    "sec_amazon":   {"pt": "AMAZÔNIA LEGAL",       "en": "LEGAL AMAZON"},
    "sec_biomes":   {"pt": "POR BIOMA",             "en": "BY BIOME"},
    "sec_cerrado":  {"pt": "CERRADO EM FOCO",       "en": "CERRADO SPOTLIGHT"},
    "sec_cover":    {"pt": "COBERTURA FLORESTAL",   "en": "FOREST COVER"},
    "sec_target":   {"pt": "META 2028",             "en": "2028 TARGET"},
    "sec_global":   {"pt": "CONTEXTO GLOBAL",       "en": "GLOBAL CONTEXT"},
    "sec_causes":   {"pt": "CAUSAS & RISCOS",       "en": "DRIVERS & RISKS"},
    "sec_takeaway": {"pt": "MENSAGENS-CHAVE",        "en": "KEY TAKEAWAYS"},

    # Cover
    "cover_title": {
        "pt": "Desmatamento no Brasil\nCaiu — Mas Ainda Não Chega.",
        "en": "Brazil's Deforestation\nIs Down — But Not Enough.",
    },
    "cover_sub": {
        "pt": "Dados PRODES · INPE  |  Apresentação para a Imprensa  ·  {year}",
        "en": "PRODES · INPE Data  |  Press Briefing  ·  {year}",
    },
    "cover_credit": {
        "pt": "Imazon — Instituto do Homem e Meio Ambiente da Amazônia",
        "en": "Imazon — Institute for People and the Environment in Amazonia",
    },

    # Slide 2: lead stat
    "stat_headline": {
        "pt": "O desmatamento na Amazônia Legal caiu {pct_decline:.0f}% em {n_years_decline} anos.",
        "en": "Amazon deforestation fell {pct_decline:.0f}% in {n_years_decline} years.",
    },
    "stat_body": {
        "pt": (
            "De {peak_km2_pt} km² em {peak_year} para {current_km2_pt} km² em {current_year} — "
            "menor área desmatada desde {first_year}.\n"
            "Ainda assim, {ratio_to_target:.1f}× a meta estabelecida para 2028."
        ),
        "en": (
            "From {peak_km2_en} km² in {peak_year} to {current_km2_en} km² in {current_year} — "
            "lowest figure since {first_year}.\n"
            "Yet still {ratio_to_target:.1f}× the target set for 2028."
        ),
    },

    # Slide 3: historical
    "hist_headline": {
        "pt": "Pico em {peak_year}, queda constante desde então — meta ainda distante.",
        "en": "Peak in {peak_year}, steady decline since — but the 2028 target remains out of reach.",
    },
    "hist_sub": {
        "pt": "Desmatamento anual na Amazônia Legal (km²)  ·  {first_year}–2028",
        "en": "Annual deforestation in Brazil's Legal Amazon (km²)  ·  {first_year}–2028",
    },

    # Slide 4: by biome
    "biome_headline": {
        "pt": "O Cerrado perde tanto quanto a Amazônia — e recebe menos atenção.",
        "en": "The Cerrado loses as much as the Amazon — and gets far less attention.",
    },
    "biome_sub": {
        "pt": "Desmatamento por bioma brasileiro (km²)  ·  {biome_year}",
        "en": "Deforestation by Brazilian biome (km²)  ·  {biome_year}",
    },

    # Slide 5: Cerrado
    "cerrado_headline": {
        "pt": "O Cerrado é o bioma brasileiro mais ameaçado proporcionalmente.",
        "en": "The Cerrado is Brazil's most proportionally threatened biome.",
    },
    "cerrado_note": {
        "pt": (
            "O Cerrado abriga 5% da biodiversidade mundial e regula o ciclo hídrico "
            "das principais bacias hidrográficas do Brasil. "
            "Recebe menos de 10% dos recursos do Fundo Amazônia."
        ),
        "en": (
            "The Cerrado harbors 5% of the world's biodiversity and regulates "
            "the water cycle of Brazil's main river basins. "
            "It receives less than 10% of Amazon Fund resources."
        ),
    },

    # Slide 6: forest cover
    "cover_headline": {
        "pt": "Mata Atlântica: 13% restam. Pantanal: 86%. Cerrado: 53%.",
        "en": "Atlantic Forest: 13% remains. Pantanal: 86%. Cerrado: 53%.",
    },
    "cover_sub_text": {
        "pt": "Vegetação nativa remanescente por bioma (% da área original)  ·  MapBiomas 2023",
        "en": "Remaining native vegetation by biome (% of original area)  ·  MapBiomas 2023",
    },

    # Slide 7: target
    "target_headline": {
        "pt": "A meta de {target_2028} km² para 2028 exige reduzir mais {pct_to_target:.0f}%.",
        "en": "The {target_2028} km² target for 2028 requires a further {pct_to_target:.0f}% cut.",
    },
    "target_sub": {
        "pt": "Trajetória observada e meta de redução (km²)  ·  {first_year}–2028",
        "en": "Observed trend and reduction target (km²)  ·  {first_year}–2028",
    },

    # Slide 8: international
    "intl_headline": {
        "pt": "Brasil lidera a queda — mas ainda é o maior desmatador tropical.",
        "en": "Brazil leads the decline — but remains the world's largest tropical deforester.",
    },
    "intl_sub": {
        "pt": "Perda de floresta tropical por país (km²)  ·  2023  ·  Fonte: GFW / FAO",
        "en": "Tropical forest loss by country (km²)  ·  2023  ·  Source: GFW / FAO",
    },
    "intl_note": {
        "pt": "* GFW mede perda de cobertura arbórea; PRODES mede desmatamento. Metodologias distintas.",
        "en": "* GFW measures tree-cover loss; PRODES measures deforestation. Different methodologies.",
    },

    # Slide 9: causes
    "causes_headline": {
        "pt": "O que explica a queda — e o que pode revertê-la.",
        "en": "What drove the decline — and what could reverse it.",
    },

    # Slide 10: takeaways
    "takeaway_headline": {
        "pt": "Três mensagens desta apresentação.",
        "en": "Three messages from this briefing.",
    },

    # Sources
    "src_prodes": {
        "pt": "Fonte: INPE/PRODES · Calculado a partir dos dados GeoParquet",
        "en": "Source: INPE/PRODES · Calculated from GeoParquet data",
    },
    "src_mapbiomas": {
        "pt": "Fonte: MapBiomas 2023 · Mapeamento Anual da Cobertura e Uso da Terra",
        "en": "Source: MapBiomas 2023 · Annual Land Cover and Use Mapping Project",
    },
    "src_gfw": {
        "pt": "Fonte: Global Forest Watch / FAO 2023 · Dados aproximados — metodologias distintas",
        "en": "Source: Global Forest Watch / FAO 2023 · Approximate data — methodologies differ",
    },
}


def _t(key: str, lang: str) -> str:
    """Return COPY[key][lang] formatted with STATS values."""
    template = COPY[key][lang]
    try:
        return template.format(**STATS)
    except KeyError:
        return template   # if STATS not yet populated, return raw template


def _causes_items(lang: str) -> list[tuple[str, str, str, str]]:
    return [
        ("▲", C_GREEN, {
            "pt": "Fiscalização reforçada",
            "en": "Enforcement strengthened",
        }[lang], {
            "pt": "Ibama e PF multiplicaram autuações e embargos desde 2023. "
                  "Operações coordenadas reduziram o desmatamento ilegal.",
            "en": "Ibama and Federal Police multiplied fines and embargoes from 2023. "
                  "Coordinated operations reduced illegal deforestation.",
        }[lang]),
        ("▲", C_GREEN, {
            "pt": "Financiamento climático",
            "en": "Climate finance increased",
        }[lang], {
            "pt": "Fundo Amazônia recebeu +R$ 3 bi em 2023–24 (Noruega, Alemanha, EUA). "
                  "Primeira fase do REDD+ Amazon operacional.",
            "en": "Amazon Fund received BRL 3 bn+ in 2023–24 (Norway, Germany, USA). "
                  "First phase of REDD+ Amazon operational.",
        }[lang]),
        ("▼", C_RED, {
            "pt": "Risco: anistia fundiária",
            "en": "Risk: land regularization bills",
        }[lang], {
            "pt": "Projetos de lei que regularizam desmatamentos ilegais antes de 2008 "
                  "ameaçam criar incentivos para novos crimes ambientais.",
            "en": "Legislation that would amnesty illegal deforestation before 2008 "
                  "risks creating incentives for new environmental crimes.",
        }[lang]),
    ]


def _takeaway_items(lang: str) -> list[tuple[str, str]]:
    pct = STATS.get("pct_decline", 0)
    yr  = STATS.get("peak_year", "")
    tgt = STATS.get("target_2028", 4_000)
    nt  = STATS.get("pct_to_target", 0)
    return {
        "pt": [
            (f"A queda de {pct:.0f}% é real — mas frágil.",
             f"A redução desde {yr} é histórica. "
             "Qualquer mudança na política de fiscalização pode revertê-la rapidamente."),
            ("O Cerrado está em crise silenciosa.",
             "Perde tanto quanto a Amazônia, tem 53% de cobertura original "
             "e recebe proporcionalmente muito menos recursos e atenção."),
            (f"A meta de {_num(tgt,'pt')} km² em 2028 exige ação agora.",
             f"Para chegar à meta, o Brasil precisa reduzir mais {nt:.0f}% "
             "nos próximos três anos. O tempo está curto."),
        ],
        "en": [
            (f"The {pct:.0f}% decline is real — but fragile.",
             f"The drop since {yr} is historic. "
             "Any shift in enforcement policy could quickly reverse those gains."),
            ("The Cerrado is in silent crisis.",
             "It loses as much as the Amazon, retains only 53% of original cover, "
             "and receives far less resources and political attention."),
            (f"The {_num(tgt,'en')} km² target requires action now.",
             f"To reach the target, Brazil must cut a further {nt:.0f}% "
             "in the next three years. Time is running short."),
        ],
    }[lang]

# ---------------------------------------------------------------------------
# MATPLOTLIB CHART ENGINE
# ---------------------------------------------------------------------------

def _style() -> None:
    plt.rcParams.update({
        "font.family":        "sans-serif",
        "font.sans-serif":    ["Arial", "Helvetica Neue", "Liberation Sans", "DejaVu Sans"],
        "figure.facecolor":   "white",
        "axes.facecolor":     "white",
        "axes.grid":          True,
        "axes.grid.axis":     "y",
        "grid.color":         "#F0F0F0",
        "grid.linewidth":     0.8,
        "axes.spines.top":    False,
        "axes.spines.right":  False,
        "axes.spines.left":   False,
        "axes.spines.bottom": True,
        "axes.axisbelow":     True,
        "xtick.bottom":       False,
        "ytick.left":         False,
    })


def _fig(w: float = 9.0, h: float = 3.8) -> tuple:
    _style()
    fig, ax = plt.subplots(figsize=(w, h))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    ax.spines["bottom"].set_color("#E0E0E0")
    ax.yaxis.set_visible(False)
    ax.xaxis.set_tick_params(length=0)
    return fig, ax


def _buf(fig) -> io.BytesIO:
    b = io.BytesIO()
    fig.savefig(b, format="png", dpi=int(CONFIG["chart_dpi"]),
                bbox_inches="tight", facecolor="white")
    b.seek(0)
    plt.close(fig)
    return b


def _lbl(ax, x, y, text, color=C_MED, size=8.5, bold=False, ha="center", va="bottom"):
    ax.text(x, y, text, ha=ha, va=va, fontsize=size, color=color,
            fontweight="bold" if bold else "normal")

# ---------------------------------------------------------------------------
# CHART FUNCTIONS  (all read from STATS)
# ---------------------------------------------------------------------------

def chart_amazon_historical(lang: str) -> io.BytesIO:
    amazon  = STATS["amazon_km2"]
    t26     = STATS["target_2026"]
    t28     = STATS["target_2028"]
    pk_yr   = STATS["peak_year"]
    cur_yr  = STATS["current_year"]
    cur_km2 = STATS["current_km2"]

    hist_yrs = sorted(amazon)
    all_yrs  = hist_yrs + [2026, 2028]
    all_vals = [amazon[y] for y in hist_yrs] + [t26, t28]
    colors   = []
    for y in hist_yrs:
        if y == pk_yr:   colors.append(C_RED)
        elif y == cur_yr: colors.append(C_BLUE)
        else:             colors.append(C_GRAY)
    colors += [C_GREEN, C_GREEN]

    pos = np.arange(len(all_yrs), dtype=float)
    fig, ax = _fig(9.2, 3.9)
    ax.bar(pos, all_vals, width=0.72, color=colors, zorder=3, linewidth=0)
    ax.set_ylim(0, max(all_vals) * 1.22)
    ax.set_xticks(pos)
    ax.set_xticklabels([str(y) for y in all_yrs], fontsize=9, color=C_MED)

    for i, (y, v) in enumerate(zip(all_yrs, all_vals)):
        is_pk  = y == pk_yr
        is_cur = y == cur_yr
        is_tgt = y in (2026, 2028)
        c  = C_RED if is_pk else (C_GREEN if is_tgt else (C_BLUE if is_cur else C_MED))
        fs = 9.5 if (is_pk or is_cur) else 8.5
        _lbl(ax, pos[i], v + max(all_vals)*0.015, _num(v, lang), c, fs, is_pk or is_cur)

    i_cur = all_yrs.index(cur_yr)
    i_26  = all_yrs.index(2026)
    divx  = (pos[i_cur] + pos[i_26]) / 2
    ax.axvline(divx, color="#DDDDDD", lw=0.9, ls="--", zorder=2)
    obs = "Observado" if lang == "pt" else "Observed"
    prj = "Projetado" if lang == "pt" else "Projected"
    _lbl(ax, divx - 0.12, max(all_vals)*1.17, obs, C_LIGHT, 7.5, ha="right")
    _lbl(ax, divx + 0.12, max(all_vals)*1.17, prj, C_GREEN, 7.5, ha="left")

    i_28 = all_yrs.index(2028)
    ax.hlines(cur_km2, pos[i_cur], pos[i_28]+0.45, colors="#BBBBBB", lw=1.0, ls=":", zorder=4)
    bx = pos[i_28] + 0.72
    ax.annotate("", xy=(bx, t28), xytext=(bx, cur_km2),
                arrowprops=dict(arrowstyle="<->", color=C_DARK, lw=1.0))
    pct = (cur_km2 - t28) / cur_km2 * 100
    ax.text(bx+0.15, (t28+cur_km2)/2, f"−{pct:.0f}%", va="center",
            fontsize=9.5, color=C_DARK, fontweight="bold")
    tgt = "Meta" if lang == "pt" else "Target"
    _lbl(ax, pos[i_28], t28 + max(all_vals)*0.03, tgt, C_GREEN, 8)

    plt.tight_layout(pad=0.2)
    return _buf(fig)


def chart_by_biome(lang: str) -> io.BytesIO:
    key   = "biomes_km2_pt" if lang == "pt" else "biomes_km2_en"
    data  = STATS.get(key, {})
    if not data:
        raise RuntimeError("No biome data in STATS")

    pairs  = sorted(data.items(), key=lambda x: x[1], reverse=True)
    biomes = [p[0] for p in pairs]
    values = [p[1] for p in pairs]

    amazon_keywords = {"amazônia", "amazon", "legal"}
    cerrado_keywords = {"cerrado"}
    bar_colors = []
    for b in biomes:
        bl = b.lower()
        if any(k in bl for k in cerrado_keywords):
            bar_colors.append(C_RED)
        elif any(k in bl for k in amazon_keywords):
            bar_colors.append(C_BLUE)
        else:
            bar_colors.append(C_GRAY)

    _style()
    fig, ax = plt.subplots(figsize=(8.8, 3.5))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    for sp in ax.spines.values():
        sp.set_visible(False)
    ax.xaxis.set_visible(False)
    ax.yaxis.set_tick_params(length=0)
    ax.set_axisbelow(True)

    pos = np.arange(len(biomes))
    ax.barh(pos, values, height=0.6, color=bar_colors, zorder=3, linewidth=0)
    ax.set_yticks(pos)
    ax.set_yticklabels(biomes, fontsize=10.5, color=C_DARK)
    ax.invert_yaxis()
    ax.set_xlim(0, max(values) * 1.28)

    for i, (v, c) in enumerate(zip(values, bar_colors)):
        ax.text(v + max(values)*0.015, i, _num(v, lang),
                va="center", ha="left", fontsize=9.5,
                color=c if c != C_GRAY else C_MED,
                fontweight="bold" if c != C_GRAY else "normal")

    plt.tight_layout(pad=0.2)
    return _buf(fig)


def chart_forest_cover(lang: str) -> io.BytesIO:
    data   = STATS["cover_pct"][lang]
    pairs  = sorted(data.items(), key=lambda x: x[1])
    biomes = [p[0] for p in pairs]
    values = [p[1] for p in pairs]
    bar_colors = [C_RED if v < 40 else (C_ORANGE if v < 65 else C_GREEN) for v in values]

    _style()
    fig, ax = plt.subplots(figsize=(8.8, 3.5))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    for sp in ax.spines.values():
        sp.set_visible(False)
    ax.xaxis.set_visible(False)
    ax.yaxis.set_tick_params(length=0)

    pos = np.arange(len(biomes))
    ax.barh(pos, values, height=0.6, color=bar_colors, alpha=0.88, zorder=3, linewidth=0)
    ax.set_yticks(pos)
    ax.set_yticklabels(biomes, fontsize=10.5, color=C_DARK)
    ax.set_xlim(0, 118)
    ax.axvline(100, color="#DDDDDD", lw=0.8, ls="--", zorder=2)
    ref = "Cobertura original" if lang == "pt" else "Original cover"
    ax.text(99, len(biomes)-0.6, ref, ha="right", va="top", fontsize=7, color=C_LIGHT)
    for i, (v, c) in enumerate(zip(values, bar_colors)):
        ax.text(v+1.5, i, f"{v:.0f}%", va="center", ha="left",
                fontsize=9.5, color=c, fontweight="bold")
    plt.tight_layout(pad=0.2)
    return _buf(fig)


def chart_target_trajectory(lang: str) -> io.BytesIO:
    amazon  = STATS["amazon_km2"]
    t26     = STATS["target_2026"]
    t28     = STATS["target_2028"]
    cur_yr  = STATS["current_year"]
    cur_km2 = STATS["current_km2"]

    hist_yrs = sorted(amazon)
    all_yrs  = hist_yrs + [2026, 2028]
    all_vals = [amazon[y] for y in hist_yrs] + [t26, t28]
    colors   = [C_GRAY] * len(hist_yrs) + [C_GREEN, C_GREEN]

    pos = np.arange(len(all_yrs), dtype=float)
    fig, ax = _fig(9.2, 3.9)
    ax.bar(pos, all_vals, width=0.72, color=colors, zorder=3, linewidth=0)
    ax.set_ylim(0, max(all_vals) * 1.22)
    ax.set_xticks(pos)
    ax.set_xticklabels([str(y) for y in all_yrs], fontsize=9, color=C_MED)

    i_cur = all_yrs.index(cur_yr)
    i_26  = all_yrs.index(2026)
    i_28  = all_yrs.index(2028)
    tx = [pos[i_cur], pos[i_26], pos[i_28]]
    ty = [cur_km2, t26, t28]
    ax.plot(tx, ty, color=C_GREEN, lw=2.2, ls="--", zorder=5,
            marker="o", markersize=5, markerfacecolor=C_GREEN)
    ax.fill_between(tx, ty, 0, alpha=0.05, color=C_GREEN, zorder=1)
    for x, y in zip(tx, ty):
        _lbl(ax, x, y + max(all_vals)*0.025, _num(y, lang), C_GREEN, 9, True)
    tgt_lbl = "Trajetória da meta" if lang == "pt" else "Target trajectory"
    ax.text(pos[i_28]+0.12, t28 - max(all_vals)*0.07, f"← {tgt_lbl}",
            ha="left", va="top", fontsize=8, color=C_GREEN, style="italic")
    plt.tight_layout(pad=0.2)
    return _buf(fig)


def chart_international(lang: str) -> io.BytesIO:
    data  = STATS["intl_km2"]
    pairs = sorted(data.items(), key=lambda x: x[1], reverse=True)
    countries = [p[0] for p in pairs]
    values    = [p[1] for p in pairs]
    bar_colors = [C_BLUE if c == "Brazil" else C_GRAY for c in countries]

    _style()
    fig, ax = plt.subplots(figsize=(8.8, 3.2))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    for sp in ax.spines.values():
        sp.set_visible(False)
    ax.xaxis.set_visible(False)
    ax.yaxis.set_tick_params(length=0)

    pos = np.arange(len(countries))
    ax.barh(pos, values, height=0.55, color=bar_colors, zorder=3, linewidth=0)
    ax.set_yticks(pos)
    ax.set_yticklabels(countries, fontsize=10.5, color=C_DARK)
    ax.invert_yaxis()
    ax.set_xlim(0, max(values) * 1.3)
    for i, (v, c) in enumerate(zip(values, bar_colors)):
        ax.text(v + max(values)*0.015, i, _num(v, lang),
                va="center", ha="left", fontsize=9.5,
                color=C_BLUE if c == C_BLUE else C_MED,
                fontweight="bold" if c == C_BLUE else "normal")
    plt.tight_layout(pad=0.2)
    return _buf(fig)

# ---------------------------------------------------------------------------
# SLIDE HELPERS
# ---------------------------------------------------------------------------

def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, l, t, w, h, fill_hex, line_hex=None):
    sh = slide.shapes.add_shape(1, _in(l), _in(t), _in(w), _in(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex:
        sh.line.color.rgb = _rgb(line_hex)
        sh.line.width = int(0.5 * 12_700)
    else:
        sh.line.fill.background()
    return sh


def _tb(slide, text: str, l, t, w, h, *,
        size=10, bold=False, italic=False, color=C_DARK,
        align=1, wrap=True, font="Arial"):
    box = slide.shapes.add_textbox(_in(l), _in(t), _in(w), _in(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(text.split("\n")):
        p   = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text       = line
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color)
        run.font.name  = font
    return box


def _pic(slide, buf: io.BytesIO, l, t, w):
    buf.seek(0)
    slide.shapes.add_picture(buf, _in(l), _in(t), width=_in(w))


def _src(slide, text: str, color=C_LIGHT):
    _tb(slide, text, 0.35, 5.32, 9.3, 0.28, size=6.5, color=color)


def _section(slide, text: str, lang: str):
    _tb(slide, text, 0.35, 0.18, 9.3, 0.28, size=7, bold=True, color=_LANG_COLOR[lang])


def _divider(slide, lang: str):
    _rect(slide, 0.35, 0.47, 9.3, 0.005, "#E8E8E8")


def _headline(slide, text: str, t=0.52, size=13.5):
    _tb(slide, text, 0.35, t, 9.3, 0.85, size=size, bold=True, color=C_DARK)


def _sub(slide, text: str, t=1.35):
    _tb(slide, text, 0.35, t, 9.3, 0.38, size=8.5, color=C_MED)


def _lang_bar(slide, lang: str):
    lc = _LANG_COLOR[lang]
    _rect(slide, 0, 0, _SW, 0.055, lc)
    _rect(slide, 0, _SH - 0.055, _SW, 0.055, lc)

# ---------------------------------------------------------------------------
# SLIDE BUILDERS
# ---------------------------------------------------------------------------

def s_cover(prs: Presentation, lang: str) -> None:
    sl = _blank(prs)
    lc = _LANG_COLOR[lang]
    _lang_bar(sl, lang)
    ll = "PT-BR" if lang == "pt" else "EN-US"
    _tb(sl, ll, 0.35, 0.12, 1.5, 0.25, size=7, bold=True, color=lc)
    _tb(sl, COPY["cover_title"][lang], 0.55, 0.7, 8.5, 1.9,
        size=30, bold=True, color=C_DARK, font="Georgia")
    _rect(sl, 0.55, 2.75, 2.8, 0.055, lc)
    _tb(sl, _t("cover_sub", lang),
        0.55, 2.92, 8.5, 0.42, size=10, color=C_MED)
    _tb(sl, COPY["cover_credit"][lang], 0.55, 3.48, 8.5, 0.38,
        size=9, italic=True, color=C_LIGHT)


def s_lead_stat(prs: Presentation, lang: str) -> None:
    sl  = _blank(prs)
    lc  = _LANG_COLOR[lang]
    pct = STATS.get("pct_decline", 0)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_amazon"][lang], lang)
    _divider(sl, lang)
    _tb(sl, f"{pct:.0f}%", 0.35, 0.55, 3.8, 1.5,
        size=88, bold=True, color=C_GREEN, wrap=False)
    _headline(sl, _t("stat_headline", lang), t=1.9, size=15)
    _tb(sl, _t("stat_body", lang), 0.35, 2.68, 9.3, 1.0, size=10.5, color=C_MED)
    _src(sl, COPY["src_prodes"][lang])


def s_amazon_historical(prs: Presentation, lang: str) -> None:
    sl  = _blank(prs)
    buf = chart_amazon_historical(lang)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_amazon"][lang], lang)
    _divider(sl, lang)
    _headline(sl, _t("hist_headline", lang))
    _sub(sl, _t("hist_sub", lang))
    _pic(sl, buf, 0.2, 1.62, 9.6)
    _src(sl, COPY["src_prodes"][lang])


def s_by_biome(prs: Presentation, lang: str) -> None:
    sl  = _blank(prs)
    buf = chart_by_biome(lang)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_biomes"][lang], lang)
    _divider(sl, lang)
    _headline(sl, COPY["biome_headline"][lang])
    _sub(sl, _t("biome_sub", lang))
    _pic(sl, buf, 0.4, 1.62, 9.2)
    _src(sl, COPY["src_prodes"][lang])


def s_cerrado_spotlight(prs: Presentation, lang: str) -> None:
    sl  = _blank(prs)
    lc  = _LANG_COLOR[lang]
    km2 = STATS.get("cerrado_km2_pt" if lang == "pt" else "cerrado_km2_en", 0)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_cerrado"][lang], lang)
    _divider(sl, lang)
    _headline(sl, COPY["cerrado_headline"][lang])

    cards = [
        (_num(km2, lang),
         "km² desmatados no Cerrado\n— quase o mesmo que a Amazônia" if lang == "pt"
         else "km² deforested in the Cerrado\n— almost as much as the Amazon"),
        ("53%",
         "de cobertura original remanescente\n— Mata Atlântica já perdeu 87%" if lang == "pt"
         else "of original cover remaining\n— Atlantic Forest already lost 87%"),
    ]
    for i, (num, lbl) in enumerate(cards):
        x = 0.4 + i * 4.75
        _rect(sl, x, 1.48, 4.35, 1.95, "#FFF8E1", C_RED)
        _tb(sl, num, x+0.25, 1.58, 3.85, 0.9, size=38, bold=True, color=C_RED)
        _tb(sl, lbl, x+0.25, 2.38, 3.85, 0.9, size=9.5, color=C_MED)

    _tb(sl, COPY["cerrado_note"][lang], 0.35, 3.55, 9.3, 1.0, size=9.5, color=C_MED)
    _src(sl, COPY["src_prodes"][lang] + "  ·  MapBiomas 2023")


def s_forest_cover(prs: Presentation, lang: str) -> None:
    sl  = _blank(prs)
    buf = chart_forest_cover(lang)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_cover"][lang], lang)
    _divider(sl, lang)
    _headline(sl, COPY["cover_headline"][lang])
    _sub(sl, COPY["cover_sub_text"][lang])
    _pic(sl, buf, 0.4, 1.62, 9.2)
    _src(sl, COPY["src_mapbiomas"][lang])


def s_target_2028(prs: Presentation, lang: str) -> None:
    sl  = _blank(prs)
    buf = chart_target_trajectory(lang)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_target"][lang], lang)
    _divider(sl, lang)
    _headline(sl, _t("target_headline", lang))
    _sub(sl, _t("target_sub", lang))
    _pic(sl, buf, 0.2, 1.62, 9.6)
    _src(sl, COPY["src_prodes"][lang])


def s_international(prs: Presentation, lang: str) -> None:
    sl  = _blank(prs)
    buf = chart_international(lang)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_global"][lang], lang)
    _divider(sl, lang)
    _headline(sl, COPY["intl_headline"][lang])
    _sub(sl, COPY["intl_sub"][lang])
    _pic(sl, buf, 0.4, 1.62, 9.2)
    _tb(sl, COPY["intl_note"][lang], 0.35, 4.85, 9.3, 0.35,
        size=6.5, italic=True, color=C_LIGHT)
    _src(sl, COPY["src_gfw"][lang])


def s_causes(prs: Presentation, lang: str) -> None:
    sl = _blank(prs)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_causes"][lang], lang)
    _divider(sl, lang)
    _headline(sl, COPY["causes_headline"][lang])
    for i, (icon, color, title, body) in enumerate(_causes_items(lang)):
        y = 1.48 + i * 1.25
        _tb(sl, f"{icon}  {title}", 0.35, y, 9.3, 0.38, size=11, bold=True, color=color)
        _tb(sl, body, 0.35, y+0.37, 9.3, 0.75, size=9.5, color=C_MED)


def s_takeaways(prs: Presentation, lang: str) -> None:
    sl = _blank(prs)
    lc = _LANG_COLOR[lang]
    _rect(sl, 0, 0, _SW, 1.05, lc)
    _tb(sl, COPY["sec_takeaway"][lang], 0.35, 0.08, 9.3, 0.28,
        size=7, bold=True, color="#FFFFFF")
    _tb(sl, COPY["takeaway_headline"][lang], 0.35, 0.33, 9.3, 0.62,
        size=14, bold=True, color="#FFFFFF")
    _rect(sl, 0, _SH-0.055, _SW, 0.055, lc)
    for i, (title, body) in enumerate(_takeaway_items(lang)):
        y = 1.12 + i * 1.42
        _rect(sl, 0.28, y, 9.44, 1.28, "#F5F5F5", "#E0E0E0")
        _rect(sl, 0.28, y, 0.38, 1.28, lc)
        _tb(sl, str(i+1), 0.29, y+0.38, 0.36, 0.5,
            size=18, bold=True, color="#FFFFFF", align=2)
        _tb(sl, title, 0.75, y+0.1,  8.85, 0.42, size=11,  bold=True, color=C_DARK)
        _tb(sl, body,  0.75, y+0.52, 8.85, 0.65, size=9.5, color=C_MED)

# ---------------------------------------------------------------------------
# SLIDE SEQUENCE
# ---------------------------------------------------------------------------

BUILDERS = [
    s_cover,
    s_lead_stat,
    s_amazon_historical,
    s_by_biome,
    s_cerrado_spotlight,
    s_forest_cover,
    s_target_2028,
    s_international,
    s_causes,
    s_takeaways,
]

# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------

def main() -> None:
    now = datetime.now().strftime("%Y-%m-%d %H:%M")
    print(f"\n{SEP}")
    print(f"  PRODES Press Briefing Generator  v{__version__}  |  {now}")
    print(f"{SEP}\n")

    # ── 1. Locate GeoParquet directory ──────────────────────────────────
    gpq_dir = (
        Path(str(CONFIG["geoparquet_dir"])) if CONFIG.get("geoparquet_dir")
        else _auto_geoparquet_dir()
    )
    if gpq_dir is None or not gpq_dir.exists():
        sys.exit(
            "[FATAL] GeoParquet directory not found.\n"
            "Run  python 02_convert_to_geoparquet.py  first, then retry."
        )

    print(f"  GeoParquet dir: {gpq_dir}\n")

    # ── 2. Load & compute real stats ────────────────────────────────────
    print(f"  Loading statistics from GeoParquet files...")
    raw   = _load_prodes_stats(gpq_dir)
    stats = _compute_derived_stats(raw)

    # Populate format-string variables expected by COPY templates
    stats["peak_km2_pt"]    = _num(stats["peak_km2"],    "pt")
    stats["peak_km2_en"]    = _num(stats["peak_km2"],    "en")
    stats["current_km2_pt"] = _num(stats["current_km2"], "pt")
    stats["current_km2_en"] = _num(stats["current_km2"], "en")
    stats["target_2028"]    = int(CONFIG["target_2028_km2"])
    stats["target_2026"]    = int(CONFIG["target_2026_km2"])
    stats["ratio_to_target"] = stats["current_km2"] / stats["target_2028"]
    stats["year"]            = now[:4]

    # Make STATS globally available for chart functions and slide builders
    STATS.update(stats)

    print(f"\n  Key statistics calculated from real data:")
    print(f"    Peak:    {STATS['peak_km2']:,.0f} km² ({STATS['peak_year']})")
    print(f"    Current: {STATS['current_km2']:,.0f} km² ({STATS['current_year']})")
    print(f"    Decline: {STATS['pct_decline']:.1f}%  over {STATS['n_years_decline']} years")
    print(f"    To 2028 target: need {STATS['pct_to_target']:.1f}% more reduction")

    # ── 3. Build presentation ────────────────────────────────────────────
    print(f"\n  Building slides...")
    prs = Presentation()
    prs.slide_width  = _in(_SW)
    prs.slide_height = _in(_SH)

    total = len(BUILDERS) * 2
    n = 0

    for lang in ("pt", "en"):
        label = "PT-BR" if lang == "pt" else "EN-US"
        print(f"  [{label}]")
        for fn in BUILDERS:
            n += 1
            name = fn.__name__.replace("s_", "")
            print(f"    [{n:2d}/{total}]  {name}", end=" ", flush=True)
            try:
                fn(prs, lang)
                print("✓")
            except Exception as exc:
                print(f"⚠  {exc} — blank slide inserted")
                _blank(prs)

    # ── 4. Save ─────────────────────────────────────────────────────────
    out = Path(str(CONFIG["output_path"]))
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    print(f"\n{DIV}")
    print(f"  Saved  : {out.resolve()}")
    print(f"  Slides : {len(prs.slides)}  ({len(BUILDERS)} PT-BR + {len(BUILDERS)} EN-US)")
    print(f"{SEP}\n")


if __name__ == "__main__":
    main()
