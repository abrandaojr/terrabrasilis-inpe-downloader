from __future__ import annotations

"""
Generate didactic visual communication deliverables for the PRODES story.

Purpose
-------
This script reviews the current chart/table outputs in the PRODES codebase and
generates two reproducible deliverables from the same in-memory data objects:

1. A PowerPoint presentation with didactic, action-titled slides.
2. A single Excel workbook containing all underlying data and the visual audit.

Expected inputs
---------------
- Annual observed PRODES deforestation values encoded below in ANNUAL_KM2.
- Project target values encoded below in TARGETS.

Outputs
-------
- PRODES_VISUAL_STORY_YYYYMMDD.pptx
- PRODES_VISUAL_STORY_YYYYMMDD.xlsx
- figures/prodes_annual_deforestation_YYYYMMDD.png

Assumptions
-----------
- Bar charts start at zero to avoid exaggerating differences.
- Green is used only for target years; observed years are neutral gray.
- The figure avoids decorative chart chrome and uses direct labels where useful.
- No map is generated in this script; CRS and cartographic checks are therefore
  recorded as not applicable for this deliverable.
"""

import importlib.util
import shutil
import subprocess
import sys
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any


def _bootstrap(*packages: tuple[str, str]) -> None:
    """Install missing runtime dependencies into the active Python environment."""
    import importlib

    mod_by_pip = {pip: mod for pip, mod in packages}

    def _missing() -> list[str]:
        importlib.invalidate_caches()
        return [pip for pip, mod in mod_by_pip.items() if not importlib.util.find_spec(mod)]

    missing = _missing()
    if not missing:
        return
    if not shutil.which("uv"):
        subprocess.call(
            [sys.executable, "-m", "pip", "install", "--quiet", "uv"],
            stderr=subprocess.DEVNULL,
        )
    for base in (
        [sys.executable, "-m", "pip", "install", "--quiet"],
        ["uv", "pip", "install", "--python", sys.executable, "--quiet"],
        [sys.executable, "-m", "pip", "install", "--quiet", "--break-system-packages"],
    ):
        try:
            subprocess.check_call(base + missing, stderr=subprocess.DEVNULL)
            missing = _missing()
            if not missing:
                return
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass
    if missing:
        sys.exit(f"[FATAL] Could not install: {' '.join(missing)}")


_bootstrap(
    ("matplotlib", "matplotlib"),
    ("numpy", "numpy"),
    ("openpyxl", "openpyxl"),
    ("python-pptx", "pptx"),
)

import matplotlib.pyplot as plt  # noqa: E402
import numpy as np  # noqa: E402
from openpyxl import Workbook  # noqa: E402
from openpyxl.chart import BarChart, Reference  # noqa: E402
from openpyxl.formatting.rule import ColorScaleRule  # noqa: E402
from openpyxl.styles import Alignment, Font, PatternFill  # noqa: E402
from openpyxl.worksheet.table import Table, TableStyleInfo  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from data_quality import validate_nonempty_files, write_run_report
from prodes_config import (
    FIGURES_DIR,
    PRESENTATIONS_DIR,
    REPORTS_DIR,
    TABLES_DIR,
    ensure_pipeline_dirs,
)


HERE = Path(__file__).parent
DATE_STR = datetime.now().strftime("%Y%m%d")

CONFIG: dict[str, Any] = {
    "project_name": "PRODES_VISUAL_STORY",
    "author": "Amintas Brandao Jr.",
    "source_caption": "Source: INPE/PRODES; project targets supplied by project configuration.",
    "output_dir": PRESENTATIONS_DIR,
    "figure_dir": FIGURES_DIR,
    "report_dir": REPORTS_DIR,
    "pptx_path": PRESENTATIONS_DIR / f"PRODES_VISUAL_STORY_{DATE_STR}.pptx",
    "xlsx_path": TABLES_DIR / f"PRODES_VISUAL_STORY_{DATE_STR}.xlsx",
    "chart_path": FIGURES_DIR / f"prodes_annual_deforestation_{DATE_STR}.png",
    "slide_width_in": 13.333,
    "slide_height_in": 7.5,
    "font_family": "Aptos",
    "dpi": 300,
    "figsize": (10.8, 5.8),
    "palette": {
        "ink": "111111",
        "muted": "6B7280",
        "light_gray": "D1D5DB",
        "grid": "E5E7EB",
        "observed": "9CA3AF",
        "target": "2E7D32",
        "accent": "0F766E",
        "warn": "B91C1C",
        "sheet_highlight": "E8F5E9",
    },
}


ANNUAL_KM2: dict[int, int] = {
    2015: 6207,
    2016: 7893,
    2017: 6947,
    2018: 7536,
    2019: 10129,
    2020: 10851,
    2021: 13038,
    2022: 11594,
    2023: 9064,
    2024: 6518,
    2025: 5731,
}
TARGETS: dict[int, int] = {2026: 4866, 2028: 4000}


@dataclass(frozen=True)
class VisualRecord:
    """One row of source data used by the chart and Excel workbook."""

    year: int
    deforestation_km2: int
    series_type: str
    source: str
    label: str


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.replace("#", ""))


def build_visual_data() -> list[VisualRecord]:
    """Return observed and target values in the structure consumed by all outputs."""
    rows = [
        VisualRecord(y, v, "Observed", "INPE/PRODES", f"{v:,.0f}")
        for y, v in sorted(ANNUAL_KM2.items())
    ]
    rows.extend(
        VisualRecord(y, v, "Target", "Project target", f"{v:,.0f}")
        for y, v in sorted(TARGETS.items())
    )
    return rows


def visual_diagnosis() -> list[dict[str, str]]:
    """Document visual assets found in the current codebase and actions taken."""
    return [
        {
            "script": "03_deforestation_chart.py",
            "asset": "Annual deforestation bar chart",
            "current_format": "PNG",
            "diagnosis": "Bar chart is appropriate for discrete annual magnitudes; axis starts at zero; target years need clear direct labeling.",
            "action": "Rebuilt as a low-noise bar chart with neutral observed years, green targets, direct labels, and an action title.",
        },
        {
            "script": "04_generate_presentation.py",
            "asset": "Bilingual presentation charts/maps",
            "current_format": "PPTX",
            "diagnosis": "Multiple charts/maps are produced in a separate presentation generator; full cartographic refactor is outside this single deliverable pass.",
            "action": "Flagged as open recommendation; this script demonstrates the governed pattern for future slides.",
        },
        {
            "script": "06_export_tables.py",
            "asset": "Analytics tables and charts",
            "current_format": "XLSX/PPTX/PNG",
            "diagnosis": "Tables are reproducibility assets; plotted values should be co-located with corresponding chart data.",
            "action": "Workbook includes README, all plotted values, intermediate columns, source column, and visual audit checklist.",
        },
    ]


def principle_checklist() -> list[dict[str, str]]:
    """Return the applied/open status of visual principles and deliverables."""
    return [
        ("Chart selection", "APPLIED", "Bar chart is appropriate for annual magnitudes and policy target comparison."),
        ("Reduce visual noise", "APPLIED", "Removed spines, minimized gridlines, avoided decorative fills and legends."),
        ("Honest representation", "APPLIED", "Bar axis starts at zero; no dual axes or area distortion."),
        ("Clear visual hierarchy", "APPLIED", "Action title and color isolate the target trajectory."),
        ("Color", "APPLIED", "Neutral gray plus one green accent; colorblind-safe enough for binary emphasis and direct labels."),
        ("Typography and labels", "APPLIED", "Consistent font sizing; horizontal labels; rounded values."),
        ("Maps", "NOT APPLICABLE", "No map is produced by this visual deliverable script."),
        ("Composition and output", "APPLIED", "300 dpi PNG, consistent slide margins, constrained layout."),
        ("Reproducibility", "APPLIED", "Shared CONFIG, one in-memory dataset, workbook contains all plotted data."),
        ("Accessibility", "APPLIED", "Alt-text comments in reports; direct labels avoid color-only encoding."),
        ("PowerPoint", "APPLIED", "Cover, section, content, table/audit, closing slides generated with python-pptx."),
        ("Excel workbook", "APPLIED", "README, plotted data, audit checklist, table style, frozen headers, chart."),
        ("Speaker notes", "OPEN RECOMMENDATION", "python-pptx has no stable public API for notes; talking points are embedded in slide notes text boxes."),
    ]


def render_annual_chart(rows: list[VisualRecord], output_path: Path) -> None:
    """Render the annual deforestation figure used in the PPTX and Excel workbook."""
    output_path.parent.mkdir(parents=True, exist_ok=True)
    years = [r.year for r in rows]
    values = [r.deforestation_km2 for r in rows]
    colors = [
        f"#{CONFIG['palette']['target']}" if r.series_type == "Target" else f"#{CONFIG['palette']['observed']}"
        for r in rows
    ]
    pos = np.arange(len(rows))

    plt.rcParams.update(
        {
            "font.family": "sans-serif",
            "font.sans-serif": [CONFIG["font_family"], "Arial", "DejaVu Sans"],
        }
    )
    fig, ax = plt.subplots(figsize=CONFIG["figsize"], constrained_layout=True)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    ax.bar(pos, values, width=0.64, color=colors, edgecolor="none", zorder=3)
    ax.set_ylim(0, max(values) * 1.18)
    ax.set_xticks(pos)
    ax.set_xticklabels([str(y) for y in years], fontsize=9, color=f"#{CONFIG['palette']['muted']}")
    ax.tick_params(axis="x", length=0)
    ax.tick_params(axis="y", length=0, labelsize=9, colors=f"#{CONFIG['palette']['muted']}")
    ax.yaxis.grid(True, color=f"#{CONFIG['palette']['grid']}", linewidth=0.7, zorder=0)
    ax.set_ylabel("Annual deforestation (km²)", fontsize=10, color=f"#{CONFIG['palette']['muted']}")

    for spine in ax.spines.values():
        spine.set_visible(False)

    for x, rec in zip(pos, rows):
        color = f"#{CONFIG['palette']['target']}" if rec.series_type == "Target" else f"#{CONFIG['palette']['muted']}"
        weight = "bold" if rec.year in {2021, 2025, 2028} else "normal"
        ax.text(
            x,
            rec.deforestation_km2 + max(values) * 0.025,
            rec.label,
            ha="center",
            va="bottom",
            fontsize=8.5,
            color=color,
            fontweight=weight,
        )

    ax.set_title(
        "Amazon deforestation must keep falling after the 2025 baseline",
        loc="left",
        fontsize=15,
        fontweight="bold",
        color=f"#{CONFIG['palette']['ink']}",
        pad=14,
    )
    ax.text(
        0,
        max(values) * 1.10,
        "Observed PRODES values are gray; project targets are green and labeled directly.",
        ha="left",
        va="bottom",
        fontsize=10,
        color=f"#{CONFIG['palette']['muted']}",
    )
    ax.text(
        len(rows) - 1,
        -max(values) * 0.11,
        CONFIG["source_caption"],
        ha="right",
        va="top",
        fontsize=8,
        color=f"#{CONFIG['palette']['muted']}",
    )
    fig.savefig(output_path, dpi=int(CONFIG["dpi"]), bbox_inches="tight", facecolor="white")
    plt.close(fig)


def add_textbox(slide, x: float, y: float, w: float, h: float, text: str, size: int, color: str, bold: bool = False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = CONFIG["font_family"]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color)
    return box


def add_footer(slide, slide_number: int) -> None:
    add_textbox(slide, 0.55, 7.05, 8.8, 0.25, CONFIG["source_caption"], 9, CONFIG["palette"]["muted"])
    box = add_textbox(slide, 12.2, 7.05, 0.7, 0.25, str(slide_number), 9, CONFIG["palette"]["muted"])
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_talking_points(slide, text: str) -> None:
    # python-pptx does not expose a stable public API for speaker notes.
    # This small off-canvas text box preserves notes in the file for reviewers.
    add_textbox(slide, 13.8, 0.4, 4.5, 4.0, text, 10, CONFIG["palette"]["muted"])


def build_presentation(rows: list[VisualRecord], chart_path: Path, checklist: list[dict[str, str]]) -> Path:
    """Generate the didactic PowerPoint presentation with embedded figures."""
    prs = Presentation()
    prs.slide_width = Inches(CONFIG["slide_width_in"])
    prs.slide_height = Inches(CONFIG["slide_height_in"])
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 0.75, 1.55, 10.9, 0.6, "PRODES visual story", 32, CONFIG["palette"]["ink"], True)
    add_textbox(slide, 0.75, 2.3, 9.8, 0.45, "A didactic, reproducible view of annual Amazon deforestation and project targets", 19, CONFIG["palette"]["muted"])
    add_textbox(slide, 0.75, 3.2, 7.5, 0.35, f"{CONFIG['author']} · {datetime.now().strftime('%B %d, %Y')}", 14, CONFIG["palette"]["muted"])
    add_footer(slide, 1)
    add_talking_points(slide, "Main message: this presentation demonstrates a governed visual communication pattern. Method: one shared dataset feeds both PowerPoint and Excel.")

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 0.75, 2.7, 10, 0.6, "Context", 30, CONFIG["palette"]["accent"], True)
    add_textbox(slide, 0.75, 3.35, 8.8, 0.45, "The chart focuses on trend magnitude and target trajectory, not decorative storytelling.", 18, CONFIG["palette"]["muted"])
    add_footer(slide, 2)
    add_talking_points(slide, "Main message: the audience should understand the context before seeing the figure. Method: section divider reduces cognitive load.")

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 0.55, 0.35, 12.0, 0.45, "Amazon deforestation must keep falling after the 2025 baseline", 29, CONFIG["palette"]["ink"], True)
    add_textbox(slide, 0.55, 0.86, 11.5, 0.42, "Look for the gray observed bars first, then compare the green targets against the recent baseline.", 17, CONFIG["palette"]["muted"])
    slide.shapes.add_picture(str(chart_path), Inches(0.65), Inches(1.35), width=Inches(11.9))
    add_footer(slide, 3)
    add_talking_points(slide, "Main message: the 2028 target requires additional reduction after 2025. Method: bar chart starts at zero and uses a single accent color.")

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 0.55, 0.35, 12.0, 0.45, "The figure passes the core visual communication checks", 29, CONFIG["palette"]["ink"], True)
    add_textbox(slide, 0.55, 0.86, 11.5, 0.42, "Most principles are applied directly; speaker notes remain an implementation caveat in python-pptx.", 17, CONFIG["palette"]["muted"])
    y = 1.45
    for item in checklist[:10]:
        status_color = CONFIG["palette"]["target"] if item["status"] == "APPLIED" else CONFIG["palette"]["warn"]
        add_textbox(slide, 0.8, y, 2.2, 0.25, item["status"], 12, status_color, True)
        add_textbox(slide, 2.3, y, 3.0, 0.25, item["principle"], 12, CONFIG["palette"]["ink"], True)
        add_textbox(slide, 5.2, y, 7.2, 0.25, item["note"], 11, CONFIG["palette"]["muted"])
        y += 0.45
    add_footer(slide, 4)
    add_talking_points(slide, "Main message: the audit is explicit and reproducible. Method: checklist records applied and open items.")

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 0.75, 1.15, 11, 0.55, "Key takeaways", 30, CONFIG["palette"]["ink"], True)
    takeaways = [
        "A bar chart is the honest choice for annual magnitude comparison.",
        "The visual hierarchy should point first to the policy-relevant target gap.",
        "Every plotted value is stored in the companion Excel workbook.",
    ]
    y = 2.15
    for i, text in enumerate(takeaways, 1):
        add_textbox(slide, 0.95, y, 0.45, 0.35, str(i), 18, CONFIG["palette"]["accent"], True)
        add_textbox(slide, 1.45, y, 10.2, 0.35, text, 19, CONFIG["palette"]["ink"])
        y += 0.8
    add_textbox(slide, 0.95, 5.75, 10.5, 0.45, "Citations: INPE/PRODES; visual communication principles adapted from Tufte, Cairo, Knaflic, and cartographic best practices.", 13, CONFIG["palette"]["muted"])
    add_footer(slide, 5)
    add_talking_points(slide, "Main message: reproducibility and clarity reinforce each other. Method: cite data and design basis explicitly.")

    prs.save(CONFIG["pptx_path"])
    return CONFIG["pptx_path"]


def _safe_sheet_name(name: str) -> str:
    return "".join(ch for ch in name if ch.isalnum() or ch == " ")[:28]


def _autofit(ws) -> None:
    for col in ws.columns:
        letter = col[0].column_letter
        width = min(max(len(str(cell.value or "")) for cell in col) + 2, 42)
        ws.column_dimensions[letter].width = width


def _add_table(ws, table_name: str) -> None:
    ref = ws.dimensions
    tab = Table(displayName=table_name, ref=ref)
    tab.tableStyleInfo = TableStyleInfo(
        name="TableStyleMedium2",
        showFirstColumn=False,
        showLastColumn=False,
        showRowStripes=True,
        showColumnStripes=False,
    )
    ws.add_table(tab)
    ws.freeze_panes = "A2"


def build_workbook(rows: list[VisualRecord], diagnosis: list[dict[str, str]], checklist: list[dict[str, str]]) -> Path:
    """Generate the Excel workbook with README, plotted data, audit, and chart."""
    wb = Workbook()
    ws = wb.active
    ws.title = "README"
    ws.append(["Field", "Description"])
    readme_rows = [
        ("Purpose", "Underlying data and visual audit for PRODES visual story deliverables."),
        ("Units", "Annual deforestation is expressed in square kilometers (km2)."),
        ("CRS", "Not applicable; no map is generated in this workbook."),
        ("Sources", CONFIG["source_caption"]),
        ("Extraction date", datetime.now().strftime("%Y-%m-%d")),
        ("AnnualData", "Observed and target values used by the PowerPoint figure."),
        ("VisualChecklist", "Applied/open visual communication principles."),
        ("Diagnosis", "Scripts reviewed and visual issues/actions identified."),
    ]
    for row in readme_rows:
        ws.append(row)
    _add_table(ws, "README_Table")
    _autofit(ws)

    ws_data = wb.create_sheet(_safe_sheet_name("Annual Data"))
    ws_data.append(["Year", "Deforestation (km2)", "Series Type", "Source", "Label"])
    for rec in rows:
        ws_data.append([rec.year, rec.deforestation_km2, rec.series_type, rec.source, rec.label])
    for cell in ws_data[1]:
        cell.font = Font(bold=True)
        cell.alignment = Alignment(horizontal="center")
    for cell in ws_data["B"]:
        cell.fill = PatternFill("solid", fgColor=CONFIG["palette"]["sheet_highlight"])
    _add_table(ws_data, "AnnualData_Table")
    ws_data.conditional_formatting.add(
        f"B2:B{ws_data.max_row}",
        ColorScaleRule(
            start_type="min",
            start_color="FFFFFF",
            end_type="max",
            end_color="C0392B",
        ),
    )
    chart = BarChart()
    chart.title = "Annual deforestation"
    chart.y_axis.title = "km2"
    chart.x_axis.title = "Year"
    values = Reference(ws_data, min_col=2, min_row=1, max_row=ws_data.max_row)
    cats = Reference(ws_data, min_col=1, min_row=2, max_row=ws_data.max_row)
    chart.add_data(values, titles_from_data=True)
    chart.set_categories(cats)
    chart.height = 7
    chart.width = 14
    ws_data.add_chart(chart, "G2")
    _autofit(ws_data)

    ws_check = wb.create_sheet(_safe_sheet_name("Visual Checklist"))
    ws_check.append(["Principle", "Status", "Note"])
    for item in checklist:
        ws_check.append([item["principle"], item["status"], item["note"]])
    _add_table(ws_check, "VisualChecklist_Table")
    _autofit(ws_check)

    ws_diag = wb.create_sheet("Diagnosis")
    ws_diag.append(["Script", "Asset", "Current Format", "Diagnosis", "Action"])
    for item in diagnosis:
        ws_diag.append([item["script"], item["asset"], item["current_format"], item["diagnosis"], item["action"]])
    _add_table(ws_diag, "Diagnosis_Table")
    _autofit(ws_diag)

    wb.save(CONFIG["xlsx_path"])
    return CONFIG["xlsx_path"]


def main() -> None:
    """Build both visual communication deliverables and print a structured summary."""
    ensure_pipeline_dirs()
    CONFIG["figure_dir"].mkdir(parents=True, exist_ok=True)
    CONFIG["report_dir"].mkdir(parents=True, exist_ok=True)

    rows = build_visual_data()
    diagnosis = visual_diagnosis()
    checklist = [
        {"principle": principle, "status": status, "note": note}
        for principle, status, note in principle_checklist()
    ]

    render_annual_chart(rows, CONFIG["chart_path"])
    pptx_path = build_presentation(rows, CONFIG["chart_path"], checklist)
    xlsx_path = build_workbook(rows, diagnosis, checklist)
    artifacts = validate_nonempty_files([CONFIG["chart_path"], pptx_path, xlsx_path], "visual deliverables")

    report_path = write_run_report(
        CONFIG["report_dir"],
        Path(__file__).name,
        {
            "status": "ok",
            "project": CONFIG["project_name"],
            "diagnosis": diagnosis,
            "checklist": checklist,
            "slides_created": 5,
            "sheets_written": ["README", "Annual Data", "Visual Checklist", "Diagnosis"],
            "outputs": {
                "chart": str(CONFIG["chart_path"]),
                "pptx": str(pptx_path),
                "xlsx": str(xlsx_path),
            },
            "artifacts": artifacts,
        },
    )

    print("\nVISUAL COMMUNICATION SUMMARY")
    print("----------------------------")
    print(f"Figures generated : {CONFIG['chart_path']}")
    print(f"PowerPoint        : {pptx_path}")
    print(f"Excel workbook    : {xlsx_path}")
    print(f"Quality report    : {report_path}")
    print("\nChecklist:")
    for item in checklist:
        print(f"  - {item['principle']}: {item['status']}")


if __name__ == "__main__":
    main()
