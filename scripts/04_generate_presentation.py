from __future__ import annotations

import importlib.util
import io
import shutil
import subprocess
import sys
from datetime import datetime
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

__version__ = "2.0.0"
__all__: list[str] = []

HERE = Path(__file__).parent  # script directory â€” used for output paths


# ---------------------------------------------------------------------------
# Dependency bootstrap
# ---------------------------------------------------------------------------


def _bootstrap(*packages: tuple[str, str]) -> None:
    """Install missing packages into the current Python environment."""
    import importlib

    mod_by_pip = {pip: mod for pip, mod in packages}

    def _still_missing(pkgs: list[str]) -> list[str]:
        importlib.invalidate_caches()
        return [p for p in pkgs if not importlib.util.find_spec(mod_by_pip[p])]

    missing = _still_missing(list(mod_by_pip))
    if not missing:
        return

    if not shutil.which("uv"):
        subprocess.call(
            [sys.executable, "-m", "pip", "install", "--quiet", "uv"],
            stderr=subprocess.DEVNULL,
        )

    strategies = [
        [sys.executable, "-m", "pip", "install", "--quiet"],
        ["uv", "pip", "install", "--python", sys.executable, "--quiet"],
        [
            sys.executable,
            "-m",
            "uv",
            "pip",
            "install",
            "--python",
            sys.executable,
            "--quiet",
        ],
        [sys.executable, "-m", "pip", "install",
         "--quiet", "--break-system-packages"],
    ]
    for base in strategies:
        if not missing:
            return
        try:
            subprocess.check_call(base + missing, stderr=subprocess.DEVNULL)
            missing = _still_missing(missing)
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass

    if missing:
        sys.exit(f"[FATAL] Could not install: {' '.join(missing)}")


_bootstrap(
    ("matplotlib", "matplotlib"),
    ("numpy", "numpy"),
    ("python-pptx", "pptx"),
    ("duckdb", "duckdb"),
    ("pyarrow", "pyarrow"),
    ("geopandas", "geopandas"),
    ("pyogrio", "pyogrio"),
)

# Third-party imports after bootstrap
import matplotlib.patches as mpatches  # noqa: E402
import matplotlib.pyplot as plt  # noqa: E402
import numpy as np  # noqa: E402
import pandas as pd  # noqa: E402
import geopandas as gpd  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN as MsoParagraphAlignment  # noqa: E402
from pptx.util import Pt  # noqa: E402

from prodes_pipeline.data_quality import (
    LineageRecord,
    StageTimer,
    configure_json_logging,
    file_inventory,
    freshness_metrics,
    parquet_quality_profile,
    require_existing_dir,
    to_jsonable,
    validate_nonempty_files,
    write_run_report,
)
from prodes_pipeline.pipeline_contracts import GEOPARQUET_CONTRACT
from prodes_pipeline.config import (
    GEOPARQUET_DIR,
    PRESENTATIONS_DIR,
    REPORTS_DIR,
    ensure_pipeline_dirs,
)


# ---------------------------------------------------------------------------
# CONFIG  â† the only section that needs to be edited
# ---------------------------------------------------------------------------


CONFIG: dict[str, object] = {
    "output_path": PRESENTATIONS_DIR / "PRODES_Press_Briefing.pptx",
    "chart_dpi": 220,
    # Directory with GeoParquet files from script 02.
    # Leave as None to auto-detect from the standard location.
    "geoparquet_dir": GEOPARQUET_DIR,
    # Policy targets (kmÂ²) â€” not data, not queried from GeoParquet.
    "target_2026_km2": 4_866,
    "target_2028_km2": 4_000,
}

SEP = "=" * 65
DIV = "-" * 65
REPORT_DIR = REPORTS_DIR
OBS_LOG = configure_json_logging(REPORT_DIR / "observability.jsonl")
_DUCKDB_CONN = None


# ---------------------------------------------------------------------------
# REFERENCE DATA  (sources outside the PRODES pipeline â€” kept as constants)
# ---------------------------------------------------------------------------


# Remaining native vegetation (% of original biome area) â€” MapBiomas 2023
# Source: https://mapbiomas.org
_COVER_PCT: dict[str, dict[str, float]] = {
    "pt": {
        "Pantanal": 86.0,
        "AmazÃ´nia": 81.0,
        "Caatinga": 60.0,
        "Cerrado": 53.0,
        "Pampa": 38.0,
        "Mata AtlÃ¢ntica": 13.0,
    },
    "en": {
        "Pantanal": 86.0,
        "Amazon": 81.0,
        "Caatinga": 60.0,
        "Cerrado Savanna": 53.0,
        "Pampa": 38.0,
        "Atlantic Forest": 13.0,
    },
}

# International tropical primary forest loss 2023 (kmÂ²) â€” GFW / FAO
# Different methodology from PRODES (tree-cover loss, not deforestation).
_INTL_LOSS_KM2: dict[str, int] = {
    "Brazil": 9_064,  # PRODES figure for comparability
    "D.R. Congo": 4_900,
    "Bolivia": 4_200,
    "Indonesia": 2_800,
    "Colombia": 1_450,
}


# ---------------------------------------------------------------------------
# STATS  â€” populated from GeoParquet in main(), never hardcoded
# ---------------------------------------------------------------------------


# All values computed on-the-fly from actual GeoParquet data.
_STATS: dict = {}

# Biome directory name (from TerraBrasilis ZIP structure) â†’ display names
_BIOME_TO_PT: dict[str, str] = {
    "Amazon Biome": "AmazÃ´nia Legal",
    "Legal Amazon": "AmazÃ´nia Legal",
    "Cerrado": "Cerrado",
    "Caatinga": "Caatinga",
    "Pantanal": "Pantanal",
    "Mata Atlantica": "Mata AtlÃ¢ntica",
    "Pampa": "Pampa",
}
_BIOME_TO_EN: dict[str, str] = {
    "Amazon Biome": "Legal Amazon",
    "Legal Amazon": "Legal Amazon",
    "Cerrado": "Cerrado Savanna",
    "Caatinga": "Caatinga",
    "Pantanal": "Pantanal",
    "Mata Atlantica": "Atlantic Forest",
    "Pampa": "Pampa",
}
_BIOME_TO_PT.update({
    "Amazon_Biome": _BIOME_TO_PT["Amazon Biome"],
    "Legal_Amazon": _BIOME_TO_PT["Legal Amazon"],
    "Mata_Atlantica": _BIOME_TO_PT["Mata Atlantica"],
})
_BIOME_TO_EN.update({
    "Amazon_Biome": _BIOME_TO_EN["Amazon Biome"],
    "Legal_Amazon": _BIOME_TO_EN["Legal Amazon"],
    "Mata_Atlantica": _BIOME_TO_EN["Mata Atlantica"],
})
_AMAZON_DIRS = {"Amazon Biome", "Amazon_Biome", "Legal Amazon", "Legal_Amazon"}
# Keywords in the ZIP stem or layer name that identify deforestation layers
_DEFOR_KEYWORDS = ("deforestation", "desmatamento", "desmat")
# Keywords that identify auxiliary layers to skip (not deforestation measurements)
_AUX_KEYWORDS = (
    "border",
    "boundary",
    "hydrography",
    "hydro",
    "indigenous_area",
    "conservation_units",
    "settlement",
    "quilombola",
    "terra_indigena",
    "unidade_conservacao",
    "uc_",
    "biome_border",
    "limite",
    "carbon",
    "biomass",
)
# Candidate column names (area in kmÂ², checked in priority order)
_AREA_COLS = ("areakm", "area_km", "area_km2", "areakm2", "area")
# Candidate column names for year
_YEAR_COLS = ("year", "ano", "yr", "data_year")


# Map biome display name (from _COVER_PCT) to its corresponding
# keyword used in _find_border_file (derived from _BIOME_TO_PT keys)
_BIOME_LABEL_TO_FILE_KEYWORD: dict[str, str] = {
    "Pantanal": "pantanal",
    "AmazÃ´nia": "amazon_biome_border",
    "Amazon": "amazon_biome_border",
    "Caatinga": "caatinga",
    "Cerrado": "cerrado",
    "Cerrado Savanna": "cerrado",
    "Pampa": "pampa",
    "Mata AtlÃ¢ntica": "mata_atlantica",
    "Atlantic Forest": "mata_atlantica",
}

_BIOME_KEYWORD_TO_ORG_DIR: dict[str, str] = {
    "amazon_biome_border": "Amazon_Biome",
    "caatinga": "Caatinga",
    "cerrado": "Cerrado",
    "mata_atlantica": "Mata_Atlantica",
    "pampa": "Pampa",
    "pantanal": "Pantanal",
}
_BORDER_FILE_CACHE: dict[tuple[str, str], Path | None] = {}


# ---------------------------------------------------------------------------
# DATA LOADING LAYER
# ---------------------------------------------------------------------------


def _duckdb():
    """Return a process-local DuckDB connection reused across stats queries."""
    global _DUCKDB_CONN
    if _DUCKDB_CONN is None:
        import duckdb

        _DUCKDB_CONN = duckdb.connect(":memory:")
    return _DUCKDB_CONN


def _auto_geoparquet_dir() -> Path | None:
    """Return the standard GeoParquet directory if it exists."""
    std = GEOPARQUET_DIR
    return std if std.exists() else None


def _detect_area_year(files: list[Path]) -> tuple[str | None, str | None]:
    """
    Inspect a sample parquet file's schema to find the area and year columns.
    Returns (area_col, year_col) â€” either may be None.
    """
    import pyarrow.parquet as pq

    for f in files[:5]:
        try:
            schema = pq.read_schema(str(f))
            low_to_orig = {n.lower(): n for n in schema.names}
            area_col = next(
                (low_to_orig[c] for c in _AREA_COLS if c in low_to_orig),
                None,
            )
            year_col = next(
                (low_to_orig[c] for c in _YEAR_COLS if c in low_to_orig),
                None,
            )
            if area_col is not None:
                return area_col, year_col
        except Exception:
            continue
    return None, None


def _infer_km2_factor(files: list[Path], area_col: str) -> float:
    """
    Read a few area values and guess the unit:
      kmÂ² â†’ factor 1.0
      ha  â†’ factor 0.01
      mÂ²  â†’ factor 0.000001
    """
    import pyarrow.parquet as pq

    sample: list[float] = []
    for f in files[:3]:
        try:
            col = (
                pq.read_table(str(f), columns=[area_col])
                .column(area_col)
                .to_pylist()
            )
            sample += [v for v in col if v and v > 0][:20]
        except Exception:
            pass
    if not sample:
        return 1.0
    med = sorted(sample)[len(sample) // 2]
    if med > 500_000:
        return 1 / 1_000_000  # mÂ²
    if med > 5_000:
        return 1 / 100  # ha
    return 1.0  # kmÂ²


def _query_series(
    files: list[Path], area_col: str, year_col: str, factor: float
) -> dict[int, float]:
    """DuckDB: aggregate area by year, return {year: kmÂ²}."""
    paths = [str(f).replace("\\", "/") for f in files]
    sql = f"""
        SELECT CAST("{year_col}" AS INTEGER)     AS yr,
               SUM(CAST("{area_col}" AS DOUBLE)) * {factor} AS km2
        FROM   read_parquet({paths!r})
        WHERE  "{year_col}"  IS NOT NULL
          AND  "{area_col}"  IS NOT NULL
          AND  CAST("{area_col}" AS DOUBLE) > 0
        GROUP  BY yr
        HAVING yr BETWEEN 2000 AND 2030
        ORDER  BY yr
    """
    try:
        rows = _duckdb().execute(sql).fetchall()
        return {int(r[0]): round(float(r[1]), 1) for r in rows if r[0] and r[1]}
    except Exception:
        return {}


def _query_total(
    files: list[Path],
    area_col: str,
    year_col: str | None,
    target_year: int | None,
    factor: float,
) -> float | None:
    """DuckDB: sum area for a given year (or all years if year_col missing)."""
    paths = [str(f).replace("\\", "/") for f in files]
    where = (
        f'WHERE CAST("{year_col}" AS INTEGER) = {target_year}'
        if year_col and target_year
        else ""
    )
    sql = f"""
        SELECT SUM(CAST("{area_col}" AS DOUBLE)) * {factor}
        FROM   read_parquet({paths!r})
        {where}
        HAVING SUM(CAST("{area_col}" AS DOUBLE)) > 0
    """
    try:
        row = _duckdb().execute(sql).fetchone()
        return round(float(row[0]), 1) if row and row[0] else None
    except Exception:
        return None


def _load_prodes_stats(geoparquet_dir: Path) -> dict:
    """
    Scan GeoParquet files and compute PRODES statistics on the fly.
    Returns a dict with keys: amazon_km2, biomes_km2_pt, biomes_km2_en,
    biome_year.
    Raises RuntimeError if no usable data found.

    The GeoParquet directory structure produced by script 02 is:
        <geoparquet_dir>/<date_or_biome>/.../<biome_name>/<category>/
        <zip>/<layer>.parquet
    The biome name is detected by matching parts against _BIOME_TO_PT.
    The category (e.g. "Yearly_deforestation") is the part immediately after.
    This is robust to any number of date-folder prefixes.
    """
    # Diagnostic: print everything found for easier debugging
    all_parquets = sorted(geoparquet_dir.rglob("*.parquet"))
    if not all_parquets:
        raise RuntimeError(
            "No parquet files found at all in:\n"
            f"  {geoparquet_dir}\n"
            "Run  python 02_convert_to_geoparquet.py  first to convert ZIPs."
        )

    print(f"  [data] Total parquet files found: {len(all_parquets)}")

    # Group parquet files by biome name.
    # Strategy: scan the full relative path for a known biome name AND a
    # deforestation keyword (anywhere in zip-stem or layer name).
    # Exclude paths that contain auxiliary-data keywords.
    biome_files: dict[str, list[Path]] = {}
    for pf in all_parquets:
        try:
            parts = pf.relative_to(geoparquet_dir).parts
            path_low = "/".join(p.lower() for p in parts)

            # Must contain at least one deforestation keyword
            if not any(k in path_low for k in _DEFOR_KEYWORDS):
                continue

            # Must NOT be an auxiliary layer
            if any(k in path_low for k in _AUX_KEYWORDS):
                continue

            # Find the known biome directory
            biome_dir = next(
                (p for p in parts[:-1] if p in _BIOME_TO_PT), None
            )
            if not biome_dir:
                continue

            biome_files.setdefault(biome_dir, []).append(pf)
        except (ValueError, IndexError):
            continue

    if not biome_files:
        sample = "\n    ".join(
            str(p.relative_to(geoparquet_dir)) for p in all_parquets[:8]
        )
        raise RuntimeError(
            "Parquet files exist but none matched a known biome + "
            "deforestation keyword.\n"
            f"  directory: {geoparquet_dir}\n"
            f"  sample paths:\n    {sample}\n"
            f"  known biome names: {sorted(_BIOME_TO_PT)}\n"
            f"  deforestation keywords searched: {_DEFOR_KEYWORDS}\n"
            f"  auxiliary keywords excluded:     {_AUX_KEYWORDS}"
        )

    print(f"  [data] Biome folders found: {sorted(biome_files)}")

    result: dict = {
        "amazon_km2": {},  # {year: kmÂ²}
        "biomes_km2_pt": {},  # {biome_name_pt: kmÂ²}
        "biomes_km2_en": {},  # {biome_name_en: kmÂ²}
        "biome_year": {},  # {biome_dir: year used}
    }

    # â”€â”€ Amazon annual series â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    amazon_files = [
        f for bd, fs in biome_files.items() if bd in _AMAZON_DIRS for f in fs
    ]
    if amazon_files:
        area_col, year_col = _detect_area_year(amazon_files)
        if area_col and year_col:
            factor = _infer_km2_factor(amazon_files, area_col)
            series = _query_series(amazon_files, area_col, year_col, factor)
            if series:
                result["amazon_km2"] = series
                print(
                    f"  [data] Amazon series: {min(series)}â€“{max(series)}  "
                    f"({len(series)} yr)"
                )
        else:
            print(
                "  [data] WARN: could not detect columns in Amazon files "
                f"(found: area_col={area_col})"
            )

    # â”€â”€ Per-biome total (most recent year with data) â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    ref_year = max(result["amazon_km2"]) if result["amazon_km2"] else None

    for biome_dir, files in biome_files.items():
        name_pt = _BIOME_TO_PT.get(biome_dir)
        name_en = _BIOME_TO_EN.get(biome_dir)
        if not name_pt:
            continue

        area_col, year_col = _detect_area_year(files)
        if not area_col:
            continue

        factor = _infer_km2_factor(files, area_col)

        # Find the most recent year available in this biome
        use_year = ref_year
        if year_col:
            try:
                paths = [str(f).replace("\\", "/") for f in files]
                row = _duckdb().execute(
                    f'SELECT MAX(CAST("{year_col}" AS INTEGER)) '
                    f"FROM read_parquet({paths!r})"
                ).fetchone()
                if row and row[0]:
                    use_year = int(row[0])
            except Exception:
                pass

        total = _query_total(files, area_col, year_col, use_year, factor)
        if total and total > 0:
            result["biomes_km2_pt"][name_pt] = total
            result["biomes_km2_en"][name_en] = total
            result["biome_year"][biome_dir] = use_year
            print(f"  [data] {biome_dir}: {total:,.0f} kmÂ² (year {use_year})")

    if not result["amazon_km2"] and not result["biomes_km2_pt"]:
        raise RuntimeError(
            "GeoParquet files found but no area/year data could be extracted.\n"
            "Check that column names include 'areakm' and 'year'."
        )

    return result


def _compute_derived_stats(raw: dict) -> dict:
    """Compute all derived statistics from raw query results."""
    s: dict = {}

    amazon = raw.get("amazon_km2", {})
    if amazon:
        s["amazon_km2"] = amazon
        s["current_year"] = max(amazon)
        s["current_km2"] = amazon[s["current_year"]]
        s["peak_year"] = max(amazon, key=lambda y: amazon[y])
        s["peak_km2"] = amazon[s["peak_year"]]
        s["first_year"] = min(amazon)
        s["first_km2"] = amazon[s["first_year"]]
        s["n_years_decline"] = s["current_year"] - s["peak_year"]
        pct = (s["peak_km2"] - s["current_km2"]) / s["peak_km2"] * 100
        s["pct_decline"] = round(pct, 1)
    else:
        s["amazon_km2"] = {}

    s["target_2026"] = int(CONFIG["target_2026_km2"])
    s["target_2028"] = int(CONFIG["target_2028_km2"])

    if "current_km2" in s and s["current_km2"] > 0:
        pct_to_target = (
            (s["current_km2"] - s["target_2028"]) / s["current_km2"] * 100
        )
        s["pct_to_target"] = round(pct_to_target, 1)
        s["ratio_to_target"] = s["current_km2"] / s["target_2028"]
    else:
        s["pct_to_target"] = 0.0
        s["ratio_to_target"] = 0.0

    s["biomes_km2_pt"] = raw.get("biomes_km2_pt", {})
    s["biomes_km2_en"] = raw.get("biomes_km2_en", {})
    s["biome_year_by_dir"] = raw.get("biome_year", {})
    # Single display year: most recent year across all biomes
    by_dict = raw.get("biome_year", {})
    s["biome_year"] = max(by_dict.values()) if by_dict else (s.get("current_year", ""))

    # Cerrado-specific stats for spotlight slide
    cerrado_km2_pt = s["biomes_km2_pt"].get("Cerrado", 0)
    cerrado_km2_en = s["biomes_km2_en"].get("Cerrado Savanna", cerrado_km2_pt)
    s["cerrado_km2_pt"] = cerrado_km2_pt
    s["cerrado_km2_en"] = cerrado_km2_en

    # Reference data (not from GeoParquet)
    s["cover_pct"] = _COVER_PCT
    s["intl_km2"] = _INTL_LOSS_KM2

    # Update Brazil row in international comparison with real current year data
    if "current_km2" in s:
        s["intl_km2"] = {**_INTL_LOSS_KM2, "Brazil": int(round(s["current_km2"]))}

    # Populate format-string variables expected by COPY templates
    s["peak_km2_pt"] = _num(s["peak_km2"], "pt")
    s["peak_km2_en"] = _num(s["peak_km2"], "en")
    s["current_km2_pt"] = _num(s["current_km2"], "pt")
    s["current_km2_en"] = _num(s["current_km2"], "en")
    s["year"] = datetime.now().strftime("%Y")  # Current year for cover slide

    return s


def _num(v: float, lang: str) -> str:
    """Format number: PT-BR uses . as thousands separator, EN uses ,"""
    formatted = f"{v:,.0f}"
    return formatted.replace(",", ".") if lang == "pt" else formatted


# ---------------------------------------------------------------------------
# STYLE
# ---------------------------------------------------------------------------


C_GREEN = "#2E7D32"
C_RED = "#C0392B"
C_BLUE = "#1A5276"
C_ORANGE = "#E67E22"
C_GRAY = "#CCCCCC"
C_DARK = "#111111"
C_MED = "#555555"
C_LIGHT = "#999999"

_LANG_COLOR = {"pt": "#1B5E20", "en": "#0D47A1"}

_SW = 10.0
_SH = 5.625


def _in(x: float) -> int:
    return int(x * 914_400)


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# BILINGUAL COPY  (templates filled at runtime from _STATS)
# ---------------------------------------------------------------------------
# Keys beginning with "tmpl_" are format strings filled via _t(key, lang).
# Static keys are language dicts used directly.


COPY: dict[str, dict[str, str]] = {
    "sec_amazon": {"pt": "AMAZÃ”NIA LEGAL", "en": "LEGAL AMAZON"},
    "sec_biomes": {"pt": "POR BIOMA", "en": "BY BIOME"},
    "sec_cerrado": {"pt": "CERRADO EM FOCO", "en": "CERRADO SPOTLIGHT"},
    "sec_cover": {"pt": "COBERTURA FLORESTAL", "en": "FOREST COVER"},
    "sec_target": {"pt": "META 2028", "en": "2028 TARGET"},
    "sec_global": {"pt": "CONTEXTO GLOBAL", "en": "GLOBAL CONTEXT"},
    "sec_causes": {"pt": "CAUSAS & RISCOS", "en": "DRIVERS & RISKS"},
    "sec_takeaway": {"pt": "MENSAGENS-CHAVE", "en": "KEY TAKEAWAYS"},
    # Cover
    "cover_title": {
        "pt": "Desmatamento no Brasil\nCaiu â€” Mas Ainda NÃ£o Chega.",
        "en": "Brazil's Deforestation\nIs Down â€” But Not Enough.",
    },
    "cover_sub": {
        "pt": "Dados PRODES Â· INPE  |  ApresentaÃ§Ã£o para a Imprensa  Â·  {year}",
        "en": "PRODES Â· INPE Data  |  Press Briefing  Â·  {year}",
    },
    "cover_credit": {
        "pt": "Imazon â€” Instituto do Homem e Meio Ambiente da AmazÃ´nia",
        "en": "Imazon â€” Institute for People and the Environment in Amazonia",
    },
    # Slide 2: lead stat
    "stat_headline": {
        "pt": (
            "O desmatamento na AmazÃ´nia Legal caiu {pct_decline:.0f}% em "
            "{n_years_decline} anos."
        ),
        "en": (
            "Amazon deforestation fell {pct_decline:.0f}% in "
            "{n_years_decline} years."
        ),
    },
    "stat_body": {
        "pt": (
            "De {peak_km2_pt} kmÂ² em {peak_year} para {current_km2_pt} kmÂ² "
            "em {current_year} â€” menor Ã¡rea desmatada desde {first_year}.\n"
            "Ainda assim, {ratio_to_target:.1f}Ã— a meta estabelecida para 2028."
        ),
        "en": (
            "From {peak_km2_en} kmÂ² in {peak_year} to {current_km2_en} kmÂ² "
            "in {current_year} â€” lowest figure since {first_year}.\n"
            "Yet still {ratio_to_target:.1f}Ã— the target set for 2028."
        ),
    },
    # Slide 3: historical
    "hist_headline": {
        "pt": (
            "Pico em {peak_year}, queda constante desde entÃ£o â€” meta ainda "
            "distante."
        ),
        "en": (
            "Peak in {peak_year}, steady decline since â€” but the 2028 "
            "target remains out of reach."
        ),
    },
    "hist_sub": {
        "pt": (
            "Desmatamento anual na AmazÃ´nia Legal (kmÂ²)  Â·  "
            "{first_year}â€“2028"
        ),
        "en": (
            "Annual deforestation in Brazil's Legal Amazon (kmÂ²)  Â·  "
            "{first_year}â€“2028"
        ),
    },
    # Slide 4: by biome
    "biome_headline": {
        "pt": (
            "O Cerrado perde tanto quanto a AmazÃ´nia â€” e recebe menos atenÃ§Ã£o."
        ),
        "en": (
            "The Cerrado loses as much as the Amazon â€” and gets far less "
            "attention."
        ),
    },
    "biome_sub": {
        "pt": "Desmatamento por bioma brasileiro (kmÂ²)  Â·  {biome_year}",
        "en": "Deforestation by Brazilian biome (kmÂ²)  Â·  {biome_year}",
    },
    # Slide 5: Cerrado
    "cerrado_headline": {
        "pt": "O Cerrado Ã© o bioma brasileiro mais ameaÃ§ado proporcionalmente.",
        "en": "The Cerrado is Brazil's most proportionally threatened biome.",
    },
    "cerrado_note": {
        "pt": (
            "O Cerrado abriga 5% da biodiversidade mundial e regula o ciclo "
            "hÃ­drico das principais bacias hidrogrÃ¡ficas do Brasil. "
            "Recebe menos de 10% dos recursos do Fundo AmazÃ´nia."
        ),
        "en": (
            "The Cerrado harbors 5% of the world's biodiversity and regulates "
            "the water cycle of Brazil's main river basins. "
            "It receives less than 10% of Amazon Fund resources."
        ),
    },
    # Slide 6: forest cover
    "cover_headline": {
        "pt": "Mata AtlÃ¢ntica: 13% restam. Pantanal: 86%. Cerrado: 53%.",
        "en": "Atlantic Forest: 13% remains. Pantanal: 86%. Cerrado: 53%.",
    },
    "cover_sub_text": {
        "pt": (
            "VegetaÃ§Ã£o nativa remanescente por bioma (% da Ã¡rea original)  "
            "Â·  MapBiomas 2023"
        ),
        "en": (
            "Remaining native vegetation by biome (% of original area)  "
            "Â·  MapBiomas 2023"
        ),
    },
    # Slide 7: target
    "target_headline": {
        "pt": (
            "A meta de {target_2028} kmÂ² para 2028 exige reduzir mais "
            "{pct_to_target:.0f}%."
        ),
        "en": (
            "The {target_2028} kmÂ² target for 2028 requires a further "
            "{pct_to_target:.0f}% cut."
        ),
    },
    "target_sub": {
        "pt": (
            "TrajetÃ³ria observada e meta de reduÃ§Ã£o (kmÂ²)  Â·  "
            "{first_year}â€“2028"
        ),
        "en": (
            "Observed trend and reduction target (kmÂ²)  Â·  "
            "{first_year}â€“2028"
        ),
    },
    # Slide 8: international
    "intl_headline": {
        "pt": (
            "Brasil lidera a queda â€” mas ainda Ã© o maior desmatador tropical."
        ),
        "en": (
            "Brazil leads the decline â€” but remains the world's largest "
            "tropical deforester."
        ),
    },
    "intl_sub": {
        "pt": (
            "Perda de floresta tropical por paÃ­s (kmÂ²)  Â·  2023  Â·  "
            "Fonte: GFW / FAO"
        ),
        "en": (
            "Tropical forest loss by country (kmÂ²)  Â·  2023  Â·  "
            "Source: GFW / FAO"
        ),
    },
    "intl_note": {
        "pt": (
            "* GFW mede perda de cobertura arbÃ³rea; PRODES mede desmatamento. "
            "Metodologias distintas."
        ),
        "en": (
            "* GFW measures tree-cover loss; PRODES measures deforestation. "
            "Different methodologies."
        ),
    },
    # Slide 9: causes
    "causes_headline": {
        "pt": "O que explica a queda â€” e o que pode revertÃª-la.",
        "en": "What drove the decline â€” and what could reverse it.",
    },
    # Slide 10: takeaways
    "takeaway_headline": {
        "pt": "TrÃªs mensagens desta apresentaÃ§Ã£o.",
        "en": "Three messages from this briefing.",
    },
    # Sources
    "src_prodes": {
        "pt": "Fonte: INPE/PRODES Â· Calculado a partir dos dados GeoParquet",
        "en": "Source: INPE/PRODES Â· Calculated from GeoParquet data",
    },
    "src_mapbiomas": {
        "pt": "Fonte: MapBiomas 2023 Â· Mapeamento Anual da Cobertura e Uso da Terra",
        "en": "Source: MapBiomas 2023 Â· Annual Land Cover and Use Mapping Project",
    },
    "src_gfw": {
        "pt": (
            "Fonte: Global Forest Watch / FAO 2023 Â· Dados aproximados â€” "
            "metodologias distintas"
        ),
        "en": (
            "Source: Global Forest Watch / FAO 2023 Â· Approximate data â€” "
            "methodologies differ"
        ),
    },
}


def _t(key: str, lang: str) -> str:
    """Return COPY[key][lang] formatted with _STATS values."""
    template = COPY[key][lang]
    try:
        return template.format(**_STATS)
    except KeyError:
        return template  # if _STATS not yet populated, return raw template


def _causes_items(lang: str) -> list[tuple[str, str, str, str]]:
    return [
        (
            "â–²",
            C_GREEN,
            {"pt": "FiscalizaÃ§Ã£o reforÃ§ada", "en": "Enforcement strengthened"}[
                lang
            ],
            (
                "Ibama e PF multiplicaram autuaÃ§Ãµes e embargos desde 2023. "
                "OperaÃ§Ãµes coordenadas reduziram o desmatamento ilegal."
            )
            if lang == "pt"
            else (
                "Ibama and Federal Police multiplied fines and embargoes "
                "from 2023. Coordinated operations reduced illegal "
                "deforestation."
            ),
        ),
        (
            "â–²",
            C_GREEN,
            {"pt": "Financiamento climÃ¡tico", "en": "Climate finance increased"}[
                lang
            ],
            (
                "Fundo AmazÃ´nia recebeu +R$ 3 bi em 2023â€“24 (Noruega, Alemanha, EUA). "
                "Primeira fase do REDD+ Amazon operacional."
            )
            if lang == "pt"
            else (
                "Amazon Fund received BRL 3 bn+ in 2023â€“24 (Norway, Germany, USA). "
                "First phase of REDD+ Amazon operational."
            ),
        ),
        (
            "â–¼",
            C_RED,
            {"pt": "Risco: anistia fundiÃ¡ria", "en": "Risk: land regularization bills"}[
                lang
            ],
            (
                "Projetos de lei que regularizam desmatamentos ilegais antes de 2008 "
                "ameaÃ§am criar incentivos para novos crimes ambientais."
            )
            if lang == "pt"
            else (
                "Legislation that would amnesty illegal deforestation before 2008 "
                "risks creating incentives for new environmental crimes."
            ),
        ),
    ]


def _takeaway_items(lang: str) -> list[tuple[str, str]]:
    pct = _STATS.get("pct_decline", 0)
    yr = _STATS.get("peak_year", "")
    tgt = _STATS.get("target_2028", 4_000)
    nt = _STATS.get("pct_to_target", 0)
    return {
        "pt": [
            (
                f"A queda de {pct:.0f}% Ã© real â€” mas frÃ¡gil.",
                f"A reduÃ§Ã£o desde {yr} Ã© histÃ³rica. "
                "Qualquer mudanÃ§a na polÃ­tica de fiscalizaÃ§Ã£o pode revertÃª-la "
                "rapidamente.",
            ),
            (
                "O Cerrado estÃ¡ em crise silenciosa.",
                "Perde tanto quanto a AmazÃ´nia, tem 53% de cobertura original "
                "e recebe proporcionalmente muito menos recursos e atenÃ§Ã£o.",
            ),
            (
                f"A meta de {_num(tgt,'pt')} kmÂ² em 2028 exige aÃ§Ã£o agora.",
                f"Para chegar Ã  meta, o Brasil precisa reduzir mais {nt:.0f}% "
                "nos prÃ³ximos trÃªs anos. O tempo estÃ¡ curto.",
            ),
        ],
        "en": [
            (
                f"The {pct:.0f}% decline is real â€” but fragile.",
                f"The drop since {yr} is historic. "
                "Any shift in enforcement policy could quickly reverse "
                "those gains.",
            ),
            (
                "The Cerrado is in silent crisis.",
                "It loses as much as the Amazon, retains only 53% of "
                "original cover, and receives far less resources and "
                "political attention.",
            ),
            (
                f"The {_num(tgt,'en')} kmÂ² target requires action now.",
                f"To reach the target, Brazil must cut a further {nt:.0f}% "
                "in the next three years. Time is running short.",
            ),
        ],
    }[lang]


# ---------------------------------------------------------------------------
# MATPLOTLIB CHART ENGINE
# ---------------------------------------------------------------------------


def _style() -> None:
    plt.rcParams.update(
        {
            "font.family": "sans-serif",
            "font.sans-serif": [
                "Arial",
                "Helvetica Neue",
                "Liberation Sans",
                "DejaVu Sans",
            ],
            "figure.facecolor": "white",
            "axes.facecolor": "white",
            "axes.grid": True,
            "axes.grid.axis": "y",
            "grid.color": "#F0F0F0",
            "grid.linewidth": 0.8,
            "axes.spines.top": False,
            "axes.spines.right": False,
            "axes.spines.left": False,
            "axes.spines.bottom": True,
            "axes.axisbelow": True,
            "xtick.bottom": False,
            "ytick.left": False,
        }
    )


def _fig(w: float = 9.0, h: float = 3.8) -> tuple:
    _style()
    fig, ax = plt.subplots(figsize=(w, h))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    ax.spines["bottom"].set_color("#E0E0E0")
    ax.yaxis.set_visible(False)
    ax.xaxis.set_tick_params(length=0)
    return fig, ax


def _buf(fig) -> io.BytesIO:
    b = io.BytesIO()
    fig.savefig(
        b,
        format="png",
        dpi=int(CONFIG["chart_dpi"]),
        bbox_inches="tight",
        facecolor="white",
    )
    b.seek(0)
    plt.close(fig)
    return b


def _lbl(
    ax,
    x,
    y,
    text,
    color=C_MED,
    size=8.5,
    bold=False,
    ha="center",
    va="bottom",
):
    ax.text(
        x,
        y,
        text,
        ha=ha,
        va=va,
        fontsize=size,
        color=color,
        fontweight="bold" if bold else "normal",
    )


# ---------------------------------------------------------------------------
# CHART FUNCTIONS  (all read from _STATS)
# ---------------------------------------------------------------------------


def chart_amazon_historical(lang: str) -> io.BytesIO:
    amazon = _STATS["amazon_km2"]
    t26 = _STATS["target_2026"]
    t28 = _STATS["target_2028"]
    pk_yr = _STATS["peak_year"]
    cur_yr = _STATS["current_year"]
    cur_km2 = _STATS["current_km2"]

    hist_yrs = sorted(amazon)
    all_yrs = hist_yrs + [2026, 2028]
    all_vals = [amazon[y] for y in hist_yrs] + [t26, t28]
    colors = []
    for y in hist_yrs:
        if y == pk_yr:
            colors.append(C_RED)
        elif y == cur_yr:
            colors.append(C_BLUE)
        else:
            colors.append(C_GRAY)
    colors += [C_GREEN, C_GREEN]

    pos = np.arange(len(all_yrs), dtype=float)
    fig, ax = _fig(9.2, 3.9)
    ax.bar(pos, all_vals, width=0.72, color=colors, zorder=3, linewidth=0)
    ax.set_ylim(0, max(all_vals) * 1.22)
    ax.set_xticks(pos)
    ax.set_xticklabels([str(y) for y in all_yrs], fontsize=9, color=C_MED)

    for i, (y, v) in enumerate(zip(all_yrs, all_vals)):
        is_pk = y == pk_yr
        is_cur = y == cur_yr
        is_tgt = y in (2026, 2028)
        c = C_RED if is_pk else (C_GREEN if is_tgt else (C_BLUE if is_cur else C_MED))
        fs = 9.5 if (is_pk or is_cur) else 8.5
        _lbl(
            ax,
            pos[i],
            v + max(all_vals) * 0.015,
            _num(v, lang),
            c,
            fs,
            is_pk or is_cur,
        )

    i_cur = all_yrs.index(cur_yr)
    i_26 = all_yrs.index(2026)
    divx = (pos[i_cur] + pos[i_26]) / 2
    ax.axvline(divx, color="#DDDDDD", lw=0.9, ls="--", zorder=2)
    obs = "Observado" if lang == "pt" else "Observed"
    prj = "Projetado" if lang == "pt" else "Projected"
    _lbl(ax, divx - 0.12, max(all_vals) * 1.17, obs, C_LIGHT, 7.5, ha="right")
    _lbl(ax, divx + 0.12, max(all_vals) * 1.17, prj, C_GREEN, 7.5, ha="left")

    i_28 = all_yrs.index(2028)
    ax.hlines(
        cur_km2, pos[i_cur], pos[i_28] + 0.45, colors="#BBBBBB", lw=1.0, ls=":", zorder=4
    )
    bx = pos[i_28] + 0.72
    ax.annotate(
        "",
        xy=(bx, t28),
        xytext=(bx, cur_km2),
        arrowprops=dict(arrowstyle="<->", color=C_DARK, lw=1.0),
    )
    pct = (cur_km2 - t28) / cur_km2 * 100
    ax.text(
        bx + 0.15,
        (t28 + cur_km2) / 2,
        f"âˆ’{pct:.0f}%",
        va="center",
        fontsize=9.5,
        color=C_DARK,
        fontweight="bold",
    )
    tgt = "Meta" if lang == "pt" else "Target"
    _lbl(ax, pos[i_28], t28 + max(all_vals) * 0.03, tgt, C_GREEN, 8)

    plt.tight_layout(pad=0.2)
    return _buf(fig)


def chart_by_biome(lang: str) -> io.BytesIO:
    key = "biomes_km2_pt" if lang == "pt" else "biomes_km2_en"
    data = _STATS.get(key, {})
    if not data:
        raise RuntimeError("No biome data in _STATS")

    pairs = sorted(data.items(), key=lambda x: x[1], reverse=True)
    biomes = [p[0] for p in pairs]
    values = [p[1] for p in pairs]

    amazon_keywords = {"amazÃ´nia", "amazon", "legal"}
    cerrado_keywords = {"cerrado"}
    bar_colors = []
    for b in biomes:
        bl = b.lower()
        if any(k in bl for k in cerrado_keywords):
            bar_colors.append(C_RED)
        elif any(k in bl for k in amazon_keywords):
            bar_colors.append(C_BLUE)
        else:
            bar_colors.append(C_GRAY)

    _style()
    fig, ax = plt.subplots(figsize=(8.8, 3.5))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    for sp in ax.spines.values():
        sp.set_visible(False)
    ax.xaxis.set_visible(False)
    ax.yaxis.set_tick_params(length=0)
    ax.set_axisbelow(True)

    pos = np.arange(len(biomes))
    ax.barh(pos, values, height=0.6, color=bar_colors, zorder=3, linewidth=0)
    ax.set_yticks(pos)
    ax.set_yticklabels(biomes, fontsize=10.5, color=C_DARK)
    ax.invert_yaxis()
    ax.set_xlim(0, max(values) * 1.28)

    for i, (v, c) in enumerate(zip(values, bar_colors)):
        ax.text(
            v + max(values) * 0.015,
            i,
            _num(v, lang),
            va="center",
            ha="left",
            fontsize=9.5,
            color=c if c != C_GRAY else C_MED,
            fontweight="bold" if c != C_GRAY else "normal",
        )

    plt.tight_layout(pad=0.2)
    return _buf(fig)


def chart_forest_cover(lang: str) -> io.BytesIO:
    data = _STATS["cover_pct"][lang]
    pairs = sorted(data.items(), key=lambda x: x[1])
    biomes = [p[0] for p in pairs]
    values = [p[1] for p in pairs]
    bar_colors = [
        C_RED if v < 40 else (C_ORANGE if v < 65 else C_GREEN) for v in values
    ]

    _style()
    fig, ax = plt.subplots(figsize=(8.8, 3.5))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    for sp in ax.spines.values():
        sp.set_visible(False)
    ax.xaxis.set_visible(False)
    ax.yaxis.set_tick_params(length=0)

    pos = np.arange(len(biomes))
    ax.barh(
        pos, values, height=0.6, color=bar_colors, alpha=0.88, zorder=3, linewidth=0
    )
    ax.set_yticks(pos)
    ax.set_yticklabels(biomes, fontsize=10.5, color=C_DARK)
    ax.set_xlim(0, 118)
    ax.axvline(100, color="#DDDDDD", lw=0.8, ls="--", zorder=2)
    ref = "Cobertura original" if lang == "pt" else "Original cover"
    ax.text(99, len(biomes) - 0.6, ref, ha="right", va="top", fontsize=7, color=C_LIGHT)
    for i, (v, c) in enumerate(zip(values, bar_colors)):
        ax.text(
            v + 1.5,
            i,
            f"{v:.0f}%",
            va="center",
            ha="left",
            fontsize=9.5,
            color=c,
            fontweight="bold",
        )
    plt.tight_layout(pad=0.2)
    return _buf(fig)


def chart_target_trajectory(lang: str) -> io.BytesIO:
    amazon = _STATS["amazon_km2"]
    t26 = _STATS["target_2026"]
    t28 = _STATS["target_2028"]
    cur_yr = _STATS["current_year"]
    cur_km2 = _STATS["current_km2"]

    hist_yrs = sorted(amazon)
    all_yrs = hist_yrs + [2026, 2028]
    all_vals = [amazon[y] for y in hist_yrs] + [t26, t28]
    colors = [C_GRAY] * len(hist_yrs) + [C_GREEN, C_GREEN]

    pos = np.arange(len(all_yrs), dtype=float)
    fig, ax = _fig(9.2, 3.9)
    ax.bar(pos, all_vals, width=0.72, color=colors, zorder=3, linewidth=0)
    ax.set_ylim(0, max(all_vals) * 1.22)
    ax.set_xticks(pos)
    ax.set_xticklabels([str(y) for y in all_yrs], fontsize=9, color=C_MED)

    i_cur = all_yrs.index(cur_yr)
    i_26 = all_yrs.index(2026)
    i_28 = all_yrs.index(2028)
    tx = [pos[i_cur], pos[i_26], pos[i_28]]
    ty = [cur_km2, t26, t28]
    ax.plot(
        tx,
        ty,
        color=C_GREEN,
        lw=2.2,
        ls="--",
        zorder=5,
        marker="o",
        markersize=5,
        markerfacecolor=C_GREEN,
    )
    ax.fill_between(tx, ty, 0, alpha=0.05, color=C_GREEN, zorder=1)
    for x, y in zip(tx, ty):
        _lbl(
            ax,
            x,
            y + max(all_vals) * 0.025,
            _num(y, lang),
            C_GREEN,
            9,
            True,
        )
    tgt_lbl = "TrajetÃ³ria da meta" if lang == "pt" else "Target trajectory"
    ax.text(
        pos[i_28] + 0.12,
        t28 - max(all_vals) * 0.07,
        f"â† {tgt_lbl}",
        ha="left",
        va="top",
        fontsize=8,
        color=C_GREEN,
        style="italic",
    )
    plt.tight_layout(pad=0.2)
    return _buf(fig)


def chart_international(lang: str) -> io.BytesIO:
    data = _STATS["intl_km2"]
    pairs = sorted(data.items(), key=lambda x: x[1], reverse=True)
    countries = [p[0] for p in pairs]
    values = [p[1] for p in pairs]
    bar_colors = [C_BLUE if c == "Brazil" else C_GRAY for c in countries]

    _style()
    fig, ax = plt.subplots(figsize=(8.8, 3.2))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    for sp in ax.spines.values():
        sp.set_visible(False)
    ax.xaxis.set_visible(False)
    ax.yaxis.set_tick_params(length=0)

    pos = np.arange(len(countries))
    ax.barh(
        pos, values, height=0.55, color=bar_colors, zorder=3, linewidth=0
    )
    ax.set_yticks(pos)
    ax.set_yticklabels(countries, fontsize=10.5, color=C_DARK)
    ax.invert_yaxis()
    ax.set_xlim(0, max(values) * 1.3)
    for i, (v, c) in enumerate(zip(values, bar_colors)):
        ax.text(
            v + max(values) * 0.015,
            i,
            _num(v, lang),
            va="center",
            ha="left",
            fontsize=9.5,
            color=C_BLUE if c == C_BLUE else C_MED,
            fontweight="bold" if c == C_BLUE else "normal",
        )
    plt.tight_layout(pad=0.2)
    return _buf(fig)


# ---------------------------------------------------------------------------
# SLIDE HELPERS
# ---------------------------------------------------------------------------


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, l, t, w, h, fill_hex, line_hex=None):
    sh = slide.shapes.add_shape(1, _in(l), _in(t), _in(w), _in(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex:
        sh.line.color.rgb = _rgb(line_hex)
        sh.line.width = int(0.5 * 12_700)
    else:
        sh.line.fill.background()
    return sh


def _tb(
    slide,
    text: str,
    l,
    t,
    w,
    h,
    *,
    size=10,
    bold=False,
    italic=False,
    color=C_DARK,
    align=MsoParagraphAlignment.LEFT,
    wrap=True,
    font="Arial",
):
    box = slide.shapes.add_textbox(_in(l), _in(t), _in(w), _in(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color)
        run.font.name = font
    return box


def _pic(slide, buf: io.BytesIO, l, t, w):
    buf.seek(0)
    slide.shapes.add_picture(buf, _in(l), _in(t), width=_in(w))


def _src(slide, text: str, color=C_LIGHT):
    _tb(slide, text, 0.35, 5.32, 9.3, 0.28, size=6.5, color=color)


def _section(slide, text: str, lang: str):
    _tb(slide, text, 0.35, 0.18, 9.3, 0.28, size=7, bold=True, color=_LANG_COLOR[lang])


def _divider(slide, lang: str):
    _rect(slide, 0.35, 0.47, 9.3, 0.005, "#E8E8E8")


def _headline(slide, text: str, t=0.52, size=13.5):
    _tb(slide, text, 0.35, t, 9.3, 0.85, size=size, bold=True, color=C_DARK)


def _sub(slide, text: str, t=1.35):
    _tb(slide, text, 0.35, t, 9.3, 0.38, size=8.5, color=C_MED)


def _lang_bar(slide, lang: str):
    lc = _LANG_COLOR[lang]
    _rect(slide, 0, 0, _SW, 0.055, lc)
    _rect(slide, 0, _SH - 0.055, _SW, 0.055, lc)


# ---------------------------------------------------------------------------
# SLIDE BUILDERS
# ---------------------------------------------------------------------------


def s_cover(prs: Presentation, lang: str) -> None:
    sl = _blank(prs)
    lc = _LANG_COLOR[lang]
    _lang_bar(sl, lang)
    ll = "PT-BR" if lang == "pt" else "EN-US"
    _tb(sl, ll, 0.35, 0.12, 1.5, 0.25, size=7, bold=True, color=lc)
    _tb(
        sl,
        COPY["cover_title"][lang],
        0.55,
        0.7,
        8.5,
        1.9,
        size=30,
        bold=True,
        color=C_DARK,
        font="Georgia",
    )
    _rect(sl, 0.55, 2.75, 2.8, 0.055, lc)
    _tb(sl, _t("cover_sub", lang), 0.55, 2.92, 8.5, 0.42, size=10, color=C_MED)
    _tb(
        sl,
        COPY["cover_credit"][lang],
        0.55,
        3.48,
        8.5,
        0.38,
        size=9,
        italic=True,
        color=C_LIGHT,
    )


def s_lead_stat(prs: Presentation, lang: str) -> None:
    sl = _blank(prs)
    pct = _STATS.get("pct_decline", 0)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_amazon"][lang], lang)
    _divider(sl, lang)
    _tb(sl, f"{pct:.0f}%", 0.35, 0.55, 3.8, 1.5,
        size=88, bold=True, color=C_GREEN, wrap=False)
    _headline(sl, _t("stat_headline", lang), t=1.9, size=15)
    _tb(sl, _t("stat_body", lang), 0.35, 2.68, 9.3, 1.0, size=10.5, color=C_MED)
    _src(sl, COPY["src_prodes"][lang])


def s_amazon_historical(prs: Presentation, lang: str) -> None:
    sl = _blank(prs)
    buf = chart_amazon_historical(lang)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_amazon"][lang], lang)
    _divider(sl, lang)
    _headline(sl, _t("hist_headline", lang))
    _sub(sl, _t("hist_sub", lang))
    _pic(sl, buf, 0.2, 1.62, 9.6)
    _src(sl, COPY["src_prodes"][lang])


def s_by_biome(prs: Presentation, lang: str) -> None:
    sl = _blank(prs)
    buf = chart_by_biome(lang)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_biomes"][lang], lang)
    _divider(sl, lang)
    _headline(sl, COPY["biome_headline"][lang])
    _sub(sl, _t("biome_sub", lang))
    _pic(sl, buf, 0.4, 1.62, 9.2)
    _src(sl, COPY["src_prodes"][lang])


def s_cerrado_spotlight(prs: Presentation, lang: str) -> None:
    sl = _blank(prs)
    km2 = _STATS.get("cerrado_km2_pt" if lang == "pt" else "cerrado_km2_en", 0)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_cerrado"][lang], lang)
    _divider(sl, lang)
    _headline(sl, COPY["cerrado_headline"][lang])

    cards = [
        (
            _num(km2, lang),
            "kmÂ² desmatados no Cerrado\nâ€” quase o mesmo que a AmazÃ´nia"
            if lang == "pt"
            else "kmÂ² deforested in the Cerrado\nâ€” almost as much as the Amazon",
        ),
        (
            "53%",
            "de cobertura original remanescente\nâ€” Mata AtlÃ¢ntica jÃ¡ perdeu 87%"
            if lang == "pt"
            else "of original cover remaining\nâ€” Atlantic Forest already lost 87%",
        ),
    ]
    for i, (num, lbl) in enumerate(cards):
        x = 0.4 + i * 4.75
        _rect(sl, x, 1.48, 4.35, 1.95, "#FFF8E1", C_RED)
        _tb(sl, num, x + 0.25, 1.58, 3.85, 0.9, size=38, bold=True, color=C_RED)
        _tb(sl, lbl, x + 0.25, 2.38, 3.85, 0.9, size=9.5, color=C_MED)

    _tb(sl, COPY["cerrado_note"][lang], 0.35, 3.55, 9.3, 1.0, size=9.5, color=C_MED)
    _src(sl, COPY["src_prodes"][lang] + "  Â·  MapBiomas 2023")


def s_forest_cover(prs: Presentation, lang: str) -> None:
    sl = _blank(prs)
    buf = chart_forest_cover(lang)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_cover"][lang], lang)
    _divider(sl, lang)
    _headline(sl, COPY["cover_headline"][lang])
    _sub(sl, COPY["cover_sub_text"][lang])
    _pic(sl, buf, 0.4, 1.62, 9.2)
    _src(sl, COPY["src_mapbiomas"][lang])


def s_target_2028(prs: Presentation, lang: str) -> None:
    sl = _blank(prs)
    buf = chart_target_trajectory(lang)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_target"][lang], lang)
    _divider(sl, lang)
    _headline(sl, _t("target_headline", lang))
    _sub(sl, _t("target_sub", lang))
    _pic(sl, buf, 0.2, 1.62, 9.6)
    _src(sl, COPY["src_prodes"][lang])


def s_international(prs: Presentation, lang: str) -> None:
    sl = _blank(prs)
    buf = chart_international(lang)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_global"][lang], lang)
    _divider(sl, lang)
    _headline(sl, COPY["intl_headline"][lang])
    _sub(sl, COPY["intl_sub"][lang])
    _tb(
        sl,
        COPY["intl_note"][lang],
        0.35,
        4.85,
        9.3,
        0.35,
        size=6.5,
        italic=True,
        color=C_LIGHT,
    )
    _src(sl, COPY["src_gfw"][lang])


def s_causes(prs: Presentation, lang: str) -> None:
    sl = _blank(prs)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_causes"][lang], lang)
    _divider(sl, lang)
    _headline(sl, COPY["causes_headline"][lang])
    for i, (icon, color, title, body) in enumerate(_causes_items(lang)):
        y = 1.48 + i * 1.25
        _tb(
            sl,
            f"{icon}  {title}",
            0.35,
            y,
            9.3,
            0.38,
            size=11,
            bold=True,
            color=color,
        )
        _tb(sl, body, 0.35, y + 0.37, 9.3, 0.75, size=9.5, color=C_MED)


def s_takeaways(prs: Presentation, lang: str) -> None:
    sl = _blank(prs)
    lc = _LANG_COLOR[lang]
    _rect(sl, 0, 0, _SW, 1.05, lc)
    _tb(
        sl,
        COPY["sec_takeaway"][lang],
        0.35,
        0.08,
        9.3,
        0.28,
        size=7,
        bold=True,
        color="#FFFFFF",
    )
    _tb(
        sl,
        COPY["takeaway_headline"][lang],
        0.35,
        0.33,
        9.3,
        0.62,
        size=14,
        bold=True,
        color="#FFFFFF",
    )
    _rect(sl, 0, _SH - 0.055, _SW, 0.055, lc)
    for i, (title, body) in enumerate(_takeaway_items(lang)):
        y = 1.12 + i * 1.42
        _rect(sl, 0.28, y, 9.44, 1.28, "#F5F5F5", "#E0E0E0")
        _rect(sl, 0.28, y, 0.38, 1.28, lc)
        _tb(
            sl,
            str(i + 1),
            0.29,
            y + 0.38,
            0.36,
            0.5,
            size=18,
            bold=True,
            color="#FFFFFF",
            align=MsoParagraphAlignment.RIGHT,
        )
        _tb(sl, title, 0.75, y + 0.1, 8.85, 0.42, size=11, bold=True, color=C_DARK)
        _tb(sl, body, 0.75, y + 0.52, 8.85, 0.65, size=9.5, color=C_MED)


# ---------------------------------------------------------------------------
# SLIDE SEQUENCE
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# MAP HELPERS
# ---------------------------------------------------------------------------

# Color ramps for NYT/Economist map style
_MAP_BG = "#F8F5F0"  # off-white parchment background
_MAP_OCEAN = "#D6E8F0"  # muted blue for ocean/water
_MAP_COUNTRY = "#E8E4DF"  # light gray for neighboring countries
_MAP_BORDER = "#BBBBBB"  # thin country borders
_MAP_STATE = "#CCCCCC"  # state borders
_MAP_AMAZON = "#E0ECD8"  # light green â€” Amazon biome base

# Year â†’ color for deforestation trend (older = lighter, recent = darker red)
_DEFOR_CMAP = "YlOrRd"


def _find_border_file(gpq_dir: Path, keyword: str) -> Path | None:
    """Search _organized or raw geoparquet for a border/boundary parquet."""
    cache_key = (str(gpq_dir.resolve()), keyword)
    if cache_key in _BORDER_FILE_CACHE:
        return _BORDER_FILE_CACHE[cache_key]

    org_dir_name = _BIOME_KEYWORD_TO_ORG_DIR.get(keyword)
    if org_dir_name:
        aux_dir = gpq_dir / "_organized" / org_dir_name / "auxiliary"
        candidates = (
            [aux_dir / "amazon_biome_border.parquet"]
            if keyword == "amazon_biome_border"
            else [aux_dir / "biome_border.parquet"]
        )
        for candidate in candidates:
            if candidate.exists():
                _BORDER_FILE_CACHE[cache_key] = candidate
                return candidate

    search_keyword = keyword.lower().replace("_", "")
    for base in (gpq_dir / "_organized", gpq_dir):
        if not base.exists():
            continue
        for pf in sorted(base.rglob("*.parquet")):
            stem = pf.stem.lower().replace("_", "")
            parent = pf.parent.name.lower()
            if parent != "auxiliary":
                continue
            if search_keyword in stem and "border" in stem:
                _BORDER_FILE_CACHE[cache_key] = pf
                return pf
    _BORDER_FILE_CACHE[cache_key] = None
    return None


def _load_geodf(
    files: list[Path],
    columns: list[str] | None = None,
    year_filter: int | None = None,
    year_col: str = "year",
) -> gpd.GeoDataFrame | None:
    """Load parquet files into a GeoDataFrame, optionally filtered by year."""
    dfs = []
    for f in files[:8]:  # limit files to keep memory manageable
        try:
            gdf = gpd.read_parquet(str(f), columns=columns)
            if year_filter is not None and year_col in gdf.columns:
                gdf = gdf[gdf[year_col].astype(int) == year_filter]
            if not gdf.empty:
                dfs.append(gdf)
        except Exception:
            continue
    if not dfs:
        return None
    try:
        return gpd.GeoDataFrame(
            pd.concat(dfs, ignore_index=True),
            crs=dfs[0].crs,
        )
    except Exception:
        return dfs[0] if dfs else None


def _apply_map_style(fig, ax) -> None:
    """Apply publication-quality map styling."""
    fig.patch.set_facecolor(_MAP_BG)
    ax.set_facecolor(_MAP_OCEAN)
    ax.set_axis_off()
    for spine in ax.spines.values():
        spine.set_visible(False)


def _map_source(ax, text: str) -> None:
    ax.text(
        0.01,
        0.01,
        text,
        transform=ax.transAxes,
        fontsize=6.5,
        color="#999999",
        style="italic",
        va="bottom",
        ha="left",
    )


# â”€â”€ MAP 1: Amazon deforestation polygons colored by year â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€


def chart_map_amazon_deforestation(lang: str) -> io.BytesIO:
    """
    Map showing Amazon deforestation polygons colored by year (2015-current).
    NYT style: off-white background, YlOrRd color ramp, clean annotation.
    Falls back to a styled infographic if geodata unavailable.
    """
    gpq_dir = (
        Path(str(CONFIG["geoparquet_dir"]))
        if CONFIG.get("geoparquet_dir")
        else GEOPARQUET_DIR
    )

    # Try to load deforestation polygons
    amazon_files = [
        pf
        for pf in sorted(gpq_dir.rglob("*.parquet"))
        if any(
            k in pf.name.lower()
            for k in ("accumulated_deforestation", "desmatamento")
        )
        and "border" not in pf.name.lower()
        and any(p in pf.parts for p in ("Amazon Biome", "Legal Amazon"))
    ]

    # Try to load biome/border for context
    border_file = _find_border_file(gpq_dir, "amazon_biome_border")

    fig, ax = plt.subplots(figsize=(6.5, 6.0))
    _apply_map_style(fig, ax)

    drawn_map = False

    if border_file:
        try:
            border = gpd.read_parquet(str(border_file))
            if not border.empty:
                border = border.to_crs("EPSG:4326") if border.crs else border
                border.plot(
                    ax=ax, color=_MAP_AMAZON, edgecolor=_MAP_BORDER,
                    linewidth=0.5, zorder=1
                )
                drawn_map = True
        except Exception:
            pass

    current_year = _STATS.get("current_year", datetime.now().year)
    years_to_show = list(range(max(2015, current_year - 9), current_year + 1))

    if amazon_files:
        try:
            dfs = []
            for f in amazon_files[:4]:
                gdf = gpd.read_parquet(str(f))
                if "year" in gdf.columns and "geometry" in gdf.columns:
                    gdf = gdf[gdf["year"].astype(int).isin(years_to_show)]
                    if not gdf.empty:
                        dfs.append(gdf[["year", "geometry"]])
            if dfs:
                combined = gpd.GeoDataFrame(
                    pd.concat(dfs, ignore_index=True),
                    crs=dfs[0].crs,
                )
                if combined.crs and combined.crs.to_epsg() != 4326:
                    combined = combined.to_crs("EPSG:4326")
                combined["year"] = combined["year"].astype(int)
                combined.plot(
                    ax=ax,
                    column="year",
                    cmap=_DEFOR_CMAP,
                    vmin=min(years_to_show),
                    vmax=max(years_to_show),
                    alpha=0.85,
                    linewidth=0,
                    zorder=2,
                )
                drawn_map = True

                # Colorbar
                import matplotlib.cm as cm_m
                import matplotlib.colors as mcolors_m

                norm = mcolors_m.Normalize(
                    vmin=min(years_to_show), vmax=max(years_to_show)
                )
                sm = cm_m.ScalarMappable(cmap=_DEFOR_CMAP, norm=norm)
                sm.set_array([])
                cbar = fig.colorbar(
                    sm, ax=ax, orientation="horizontal",
                    fraction=0.04, pad=0.01, shrink=0.6
                )
                cbar.set_label(
                    "Ano de detecÃ§Ã£o" if lang == "pt" else "Year detected",
                    fontsize=8,
                    color="#555555",
                )
                cbar.ax.tick_params(labelsize=7, color="#555555")
                for spine in cbar.ax.spines.values():
                    spine.set_visible(False)
        except Exception:
            pass

    if not drawn_map:
        # Fallback: styled text placeholder
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.text(
            0.5,
            0.6,
            "ðŸ—º",
            ha="center",
            va="center",
            fontsize=48,
            color="#CCCCCC",
            transform=ax.transAxes,
        )
        msg = (
            "Mapa indisponÃ­vel â€” execute o script 02\n"
            "para gerar os dados GeoParquet"
            if lang == "pt"
            else "Map unavailable â€” run script 02\n" "to generate GeoParquet data"
        )
        ax.text(
            0.5,
            0.35,
            msg,
            ha="center",
            va="center",
            fontsize=9,
            color="#AAAAAA",
            transform=ax.transAxes,
            style="italic",
        )
    else:
        # Annotation: current year kmÂ²
        current_km2 = _STATS.get("current_km2", 0)
        txt = (
            f"{_num(current_km2,'pt')} kmÂ²\ndesmatados em {current_year}"
            if lang == "pt"
            else f"{_num(current_km2,'en')} kmÂ²\ndeforested in {current_year}"
        )
        ax.text(
            0.97,
            0.97,
            txt,
            transform=ax.transAxes,
            ha="right",
            va="top",
            fontsize=9,
            color="#C0392B",
            fontweight="bold",
            bbox=dict(
                boxstyle="round,pad=0.3", fc="white", ec="#CCCCCC", alpha=0.9
            ),
        )

    src = "Fonte: INPE/PRODES" if lang == "pt" else "Source: INPE/PRODES"
    _map_source(ax, src)

    plt.tight_layout(pad=0.1)
    return _buf(fig)


# â”€â”€ MAP 2: Brazil biomes colored by remaining forest cover â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€


def chart_map_biomes_coverage(lang: str) -> io.BytesIO:
    """
    Map of Brazil biomes colored by remaining native vegetation cover.
    Economist style: muted palette, direct labels, source note.
    Falls back to a styled bar visual if border data unavailable.
    """
    gpq_dir = (
        Path(str(CONFIG["geoparquet_dir"]))
        if CONFIG.get("geoparquet_dir")
        else GEOPARQUET_DIR
    )

    # Forest cover reference (MapBiomas)
    cover_data = _COVER_PCT[lang]  # [(biome_name, pct), ...]

    import matplotlib.cm as cm_m
    import matplotlib.colors as mcolors_m

    # Color scale: red (<40%) â†’ orange (<65%) â†’ green (>65%)
    _COV_CMAP = mcolors_m.LinearSegmentedColormap.from_list(
        "cover", ["#C0392B", "#E67E22", "#27AE60"], N=256
    )
    norm = mcolors_m.Normalize(vmin=0, vmax=100)

    fig, ax = plt.subplots(figsize=(6.5, 6.0))
    _apply_map_style(fig, ax)

    drawn_any = False

    for biome_label, pct in sorted(cover_data.items(), key=lambda x: x[1], reverse=True):
        color = mcolors_m.to_hex(_COV_CMAP(norm(pct)))

        # Find file keyword for this biome_label
        file_keyword = _BIOME_LABEL_TO_FILE_KEYWORD.get(biome_label)
        border_file = (
            _find_border_file(gpq_dir, file_keyword) if file_keyword else None
        )

        if border_file:
            try:
                gdf = gpd.read_parquet(str(border_file), columns=["geometry"])
                if not gdf.empty:
                    if gdf.crs and gdf.crs.to_epsg() != 4326:
                        gdf = gdf.to_crs("EPSG:4326")
                    gdf = gdf.copy()
                    gdf["geometry"] = gdf.geometry.simplify(
                        0.03, preserve_topology=True
                    )
                    gdf.plot(
                        ax=ax,
                        color=color,
                        edgecolor="#FFFFFF",
                        linewidth=0.6,
                        alpha=0.9,
                        zorder=2,
                    )
                    drawn_any = True

                    # Centroid label
                    try:
                        geometry_union = (
                            gdf.geometry.union_all()
                            if hasattr(gdf.geometry, "union_all")
                            else gdf.geometry.unary_union
                        )
                        cx = geometry_union.centroid.x
                        cy = geometry_union.centroid.y
                        short = biome_label.split()[0]  # first word
                        ax.text(
                            cx,
                            cy,
                            f"{short}\n{pct:.0f}%",
                            ha="center",
                            va="center",
                            fontsize=6.5,
                            color="white" if pct < 50 else "#333333",
                            fontweight="bold",
                            zorder=5,
                        )
                    except Exception:
                        pass
            except Exception:
                pass

    if not drawn_any:
        # Fallback: horizontal bar chart styled as a map legend
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.set_facecolor(_MAP_BG)
        title = (
            "Cobertura Vegetal Remanescente por Bioma (%)"
            if lang == "pt"
            else "Remaining Vegetation Cover by Biome (%)"
        )
        ax.text(
            0.5,
            0.95,
            title,
            ha="center",
            va="top",
            fontsize=9,
            fontweight="bold",
            color="#111111",
            transform=ax.transAxes,
        )

        sorted_cover = sorted(cover_data.items(), key=lambda x: x[1])
        bar_h = 0.10
        for i, (bname, pct) in enumerate(sorted_cover):
            y = 0.12 + i * 0.13
            color = mcolors_m.to_hex(_COV_CMAP(norm(pct)))
            bar = plt.Rectangle(
                (0.08, y),
                pct / 100 * 0.60,
                bar_h,
                color=color,
                transform=ax.transAxes,
                zorder=2,
            )
            ax.add_patch(bar)
            ax.text(
                0.07,
                y + bar_h / 2,
                bname,
                ha="right",
                va="center",
                fontsize=8,
                color="#333333",
                transform=ax.transAxes,
            )
            ax.text(
                0.08 + pct / 100 * 0.60 + 0.01,
                y + bar_h / 2,
                f"{pct:.0f}%",
                ha="left",
                va="center",
                fontsize=8,
                fontweight="bold",
                color=color,
                transform=ax.transAxes,
            )
        drawn_any = True

    # Colorbar
    if drawn_any:
        sm = cm_m.ScalarMappable(cmap=_COV_CMAP, norm=norm)
        sm.set_array([])
        cbar = fig.colorbar(
            sm, ax=ax, orientation="horizontal", fraction=0.03, pad=0.01, shrink=0.55
        )
        lbl = (
            "Cobertura remanescente (%)" if lang == "pt" else "Remaining cover (%)"
        )
        cbar.set_label(lbl, fontsize=8, color="#555555")
        cbar.ax.tick_params(labelsize=7, color="#555555")
        for spine in cbar.ax.spines.values():
            spine.set_visible(False)

    src = "Fonte: MapBiomas 2023" if lang == "pt" else "Source: MapBiomas 2023"
    _map_source(ax, src)

    plt.tight_layout(pad=0.1)
    return _buf(fig)


# â”€â”€ Map slide builders â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€


def s_map_amazon(prs: Presentation, lang: str) -> None:
    """Map slide: Amazon deforestation polygons colored by detection year."""
    sl = _blank(prs)
    buf = chart_map_amazon_deforestation(lang)
    cur_yr = _STATS.get("current_year", "")
    cur_km2 = _STATS.get("current_km2", 0)
    pct = _STATS.get("pct_decline", 0)

    _lang_bar(sl, lang)
    _section(sl, COPY["sec_amazon"][lang], lang)
    _divider(sl, lang)

    headline_pt = f"Onde estÃ£o as {_num(cur_km2,'pt')} kmÂ² desmatadas em {cur_yr}?"
    headline_en = f"Where are the {_num(cur_km2,'en')} kmÂ² deforested in {cur_yr}?"
    _headline(sl, headline_pt if lang == "pt" else headline_en)

    sub_pt = (
        "PolÃ­gonos de desmatamento na AmazÃ´nia Legal Â· 2015â€“"
        f"{cur_yr} Â· cores por ano de detecÃ§Ã£o"
    )
    sub_en = (
        "Deforestation polygons in Brazil's Legal Amazon Â· 2015â€“"
        f"{cur_yr} Â· colored by detection year"
    )
    _sub(sl, sub_pt if lang == "pt" else sub_en)

    # Map on left (60% width), stat panel on right
    _pic(sl, buf, 0.2, 1.6, 6.0)

    # Right panel: key stats
    rx = 6.4
    _rect(sl, rx, 1.62, 3.35, 3.7, "#F8F5F0", "#E0E0E0")

    stats_items_pt = [
        (f"{_num(cur_km2,'pt')} kmÂ²", f"desmatados em {cur_yr}"),
        (
            f"âˆ’{pct:.0f}%",
            f"desde o pico de {_STATS.get('peak_year','')}",
        ),
        (f"{_num(_STATS.get('target_2028',4000),'pt')} kmÂ²", "meta para 2028"),
    ]
    stats_items_en = [
        (f"{_num(cur_km2,'en')} kmÂ²", f"deforested in {cur_yr}"),
        (f"âˆ’{pct:.0f}%", f"since {_STATS.get('peak_year','')} peak"),
        (f"{_num(_STATS.get('target_2028',4000),'en')} kmÂ²", "2028 target"),
    ]
    items = stats_items_pt if lang == "pt" else stats_items_en

    for i, (num, lbl) in enumerate(items):
        y = 1.75 + i * 1.15
        _tb(
            sl,
            num,
            rx + 0.2,
            y,
            3.0,
            0.55,
            size=22,
            bold=True,
            color=C_RED if i == 0 else (C_GREEN if i == 2 else C_BLUE),
        )
        _tb(sl, lbl, rx + 0.2, y + 0.52, 3.0, 0.42, size=9, color=C_MED)

    _src(sl, COPY["src_prodes"][lang])


def s_map_biomes(prs: Presentation, lang: str) -> None:
    """Map slide: Brazil biomes colored by remaining forest cover (MapBiomas)."""
    sl = _blank(prs)
    buf = chart_map_biomes_coverage(lang)
    _lang_bar(sl, lang)
    _section(sl, COPY["sec_cover"][lang], lang)
    _divider(sl, lang)
    _headline(sl, COPY["cover_headline"][lang])
    _sub(sl, COPY["cover_sub_text"][lang])

    # Map on left, legend on right
    _pic(sl, buf, 0.2, 1.6, 6.2)

    # Right: compact cover table
    rx = 6.55
    _rect(sl, rx, 1.62, 3.2, 3.7, "#F8F5F0", "#E0E0E0")

    cover = sorted(_STATS["cover_pct"][lang].items(), key=lambda x: x[1])
    _tb(
        sl,
        "Cobertura\nremanescente" if lang == "pt" else "Remaining\ncover",
        rx + 0.15,
        1.68,
        2.9,
        0.6,
        size=8,
        bold=True,
        color=C_MED,
    )

    for i, (biome, pct) in enumerate(cover):
        y = 2.28 + i * 0.56
        color = C_RED if pct < 40 else (C_ORANGE if pct < 65 else C_GREEN)
        _tb(sl, biome[:18], rx + 0.15, y, 2.0, 0.45, size=8.5, color=C_DARK)
        _tb(
            sl,
            f"{pct:.0f}%",
            rx + 2.3,
            y,
            0.7,
            0.45,
            size=9,
            bold=True,
            color=color,
        )

    _src(sl, COPY["src_mapbiomas"][lang])


_BUILDERS = [
    s_cover,
    s_lead_stat,
    s_amazon_historical,
    s_map_amazon,  # map: deforestation polygons colored by year
    s_by_biome,
    s_cerrado_spotlight,
    s_map_biomes,  # map: biomes colored by remaining forest cover
    s_forest_cover,
    s_target_2028,
    s_international,
    s_causes,
    s_takeaways,
]


# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------


def main() -> None:
    ensure_pipeline_dirs()
    now = datetime.now().strftime("%Y-%m-%d %H:%M")
    print(f"\n{SEP}")
    print(f"  PRODES Press Briefing Generator  v{__version__}  |  {now}")
    print(f"{SEP}\n")

    # 1. Locate GeoParquet directory
    gpq_dir = (
        Path(str(CONFIG["geoparquet_dir"]))
        if CONFIG.get("geoparquet_dir")
        else _auto_geoparquet_dir()
    )
    if gpq_dir is None:
        sys.exit(
            "[FATAL] GeoParquet directory not found.\n"
            "Run  python 02_convert_to_geoparquet.py  first, then retry."
        )
    require_existing_dir(gpq_dir, "GeoParquet")

    print(f"  GeoParquet dir: {gpq_dir}\n")
    input_files = list(gpq_dir.rglob("*.parquet"))
    input_quality = {
        "contract": to_jsonable(GEOPARQUET_CONTRACT),
        "inventory": file_inventory(input_files),
        "freshness": freshness_metrics(input_files, GEOPARQUET_CONTRACT.freshness),
        "parquet_profile": parquet_quality_profile(input_files, GEOPARQUET_CONTRACT),
    }
    OBS_LOG.emit(
        "data_contract",
        stage_name="04_geoparquet_input_contract",
        contract=to_jsonable(GEOPARQUET_CONTRACT),
        metrics=input_quality,
    )

    # 2. Load & compute real stats
    print("  Loading statistics from GeoParquet files...")
    stats_timer = StageTimer("04_load_compute_stats")
    raw = _load_prodes_stats(gpq_dir)
    stats = _compute_derived_stats(raw)
    OBS_LOG.emit(
        "stage_metrics",
        **to_jsonable(
            stats_timer.finish(
                "ok",
                input_row_count=input_quality["parquet_profile"].get("row_count"),
                output_row_count=len(stats),
                anomalies={
                    "schema": input_quality["parquet_profile"].get(
                        "schema_anomalies", []
                    ),
                    "distribution": input_quality["parquet_profile"].get(
                        "distribution_anomalies", []
                    ),
                    "freshness": input_quality["freshness"].get("stale", []),
                },
            )
        ),
    )

    # Make _STATS globally available for chart functions and slide builders
    _STATS.update(stats)

    print("\n  Key statistics calculated from real data:")
    print(f"    Peak:    {_STATS['peak_km2']:,.0f} kmÂ² ({_STATS['peak_year']})")
    print(f"    Current: {_STATS['current_km2']:,.0f} kmÂ² ({_STATS['current_year']})")
    print(
        f"    Decline: {_STATS['pct_decline']:.1f}%  "
        f"over {_STATS['n_years_decline']} years"
    )
    print(
        f"    To 2028 target: need {_STATS['pct_to_target']:.1f}% more reduction"
    )

    # 3. Build presentation
    print("\n  Building slides...")
    prs = Presentation()
    build_timer = StageTimer("04_build_presentation")
    prs.slide_width = _in(_SW)
    prs.slide_height = _in(_SH)

    total = len(_BUILDERS) * 2
    n = 0

    for lang in ("pt", "en"):
        label = "PT-BR" if lang == "pt" else "EN-US"
        print(f"  [{label}]")
        for fn in _BUILDERS:
            n += 1
            name = fn.__name__.replace("s_", "")
            print(f"    [{n:2d}/{total}]  {name}", end=" ", flush=True)
            try:
                fn(prs, lang)
                print("OK")
            except Exception as exc:
                print(f"WARN  {exc} - blank slide inserted")
                _blank(prs)

    # 4. Save
    out = Path(str(CONFIG["output_path"]))
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    artifacts = validate_nonempty_files([out], "presentation")
    OBS_LOG.emit(
        "stage_metrics",
        **to_jsonable(
            build_timer.finish(
                "ok",
                input_row_count=len(_STATS),
                output_row_count=len(prs.slides),
            )
        ),
    )
    report_path = write_run_report(
        REPORT_DIR,
        Path(__file__).name,
        {
            "status": "ok",
            "version": __version__,
            "geoparquet_dir": str(gpq_dir),
            "input_quality": input_quality,
            "stats_keys": sorted(_STATS),
            "slide_count": len(prs.slides),
            "expected_slide_count": len(_BUILDERS) * 2,
            "artifacts": artifacts,
            "lineage": LineageRecord(
                stage_name="04_generate_presentation",
                upstream_sources=[str(gpq_dir)],
                transformation="Compute PRODES summary statistics from GeoParquet files and render bilingual PowerPoint slides.",
                downstream_outputs=[str(out)],
                contracts=[GEOPARQUET_CONTRACT.name],
            ),
        },
    )

    print(f"\n{DIV}")
    print(f"  Saved  : {out.resolve()}")
    print(
        f"  Slides : {len(prs.slides)}  "
        f"({len(_BUILDERS)} PT-BR + {len(_BUILDERS)} EN-US)"
    )
    print(f"  Quality report: {report_path}")
    print(f"{SEP}\n")


if __name__ == "__main__":
    main()

