"""
06_export_tables.py
===================
Senior Environmental Data Science Analytics Pipeline â€” v2.0

Spatial-Temporal Dynamics of Natural Vegetation Suppression
and Secondary Succession Trends in Brazilian Biomes

MATHEMATICAL FRAMEWORK
----------------------
All parameters are computed exclusively from primary GeoParquet files
produced by scripts 02 and 05, via DuckDB SQL engines. No external
datasets are used.

Let P = {páµ¢} be the universe of suppression polygons with attributes
(t=year, s=state, m=municipality, c=class, Î±=area_km2), and
V = {vâ±¼} the universe of secondary vegetation polygons with additional
attribute Î´=age_class.

P1  Annual Suppression Series
    S(t,s,m) = Î£_{p: t,s,m} Î±_p        [GROUP BY year, state, muni]

P2  Cumulative Suppression
    C(t,s,m) = Î£_{Ï„â‰¤t} S(Ï„,s,m)        [SUM() OVER (PARTITION BY s,m ORDER BY t)]

P3  Natural Vegetation Remaining â€” Parameter A  (stock estimate)
    NV_A(t)  = Ã‚â‚€ âˆ’ C(t)               [Ã‚â‚€ = class-aggregate at tâ‚€]

P4  Natural Vegetation â€” Parameter B  (class partition)
    NV_B(t,c) = Î£_{p: t,c} Î±_p         [GROUP BY year, classname]

P5  Secondary Vegetation Annual Extent
    VS(t,s)  = Î£_{v: t,s} Î±_v           [GROUP BY year, state]

P6  Net Annual Increment of Secondary Vegetation
    Î”VS(t,s) = VS(t,s) âˆ’ VS(tâˆ’1,s)     [LAG() window function]

P7  SV Parameter A â€” Age-Class Partition
    VS_A(t,Î´) = Î£_{v: t,Î´} Î±_v         [Î´ âˆˆ {young, intermediate, mature}]

P8  SV Parameter B â€” Land-Use-History Partition
    VS_B(t,h) = Î£_{v: t,h} Î±_v         [h from classname / land-use field]

P9  Administrative Cross-Tabulation
    For all k âˆˆ {P1..P8}: X_k(t,s,m) at municipality and state level

OUTPUTS
-------
  configured tables workspace
      PRODES_Analytics_PT_<date>.xlsx
      PRODES_Analytics_EN_<date>.xlsx
      PRODES_Analytics_PT_<date>.pptx
      PRODES_Analytics_EN_<date>.pptx
  charts/
      suppression_trend_<date>.png
      cumulative_suppression_<date>.png
      sv_dynamics_<date>.png
      sv_increment_<date>.png
      state_ranking_<date>.png
      sv_subclass_<date>.png

Author
------
Amintas BrandÃ£o Jr. <abrandaojr@gmail.com>
Imazon â€” Instituto do Homem e Meio Ambiente da AmazÃ´nia

License
-------
MIT
"""

from __future__ import annotations

__version__ = "2.0.0"
__all__: list[str] = []

import importlib.util
import io
import subprocess
import sys
from collections import defaultdict
from datetime import datetime, timezone
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))
from typing import NamedTuple

HERE = Path(__file__).parent

from prodes_pipeline.config import (  # noqa: E402
    GEOPARQUET_DIR,
    REPORTS_DIR,
    TABLES_DIR,
    ensure_pipeline_dirs,
)

GPQ_DIR = GEOPARQUET_DIR
CHART_DIR = TABLES_DIR / "charts"

# ============================================================================
# DEPENDENCY BOOTSTRAP
# ============================================================================

def _bootstrap(*packages: tuple[str, str]) -> None:
    """Install missing packages into the active Python environment."""
    import importlib, shutil
    mod_by_pip = {pip: mod for pip, mod in packages}

    def _still_missing(pkgs: list[str]) -> list[str]:
        importlib.invalidate_caches()
        return [p for p in pkgs if not importlib.util.find_spec(mod_by_pip[p])]

    missing = _still_missing(list(mod_by_pip))
    if not missing:
        return
    if not shutil.which("uv"):
        subprocess.call([sys.executable, "-m", "pip", "install", "--quiet", "uv"],
                        stderr=subprocess.DEVNULL)
    for base in [
        [sys.executable, "-m", "pip", "install", "--quiet"],
        ["uv", "pip", "install", "--python", sys.executable, "--quiet"],
        [sys.executable, "-m", "pip", "install", "--quiet", "--break-system-packages"],
    ]:
        if not missing:
            return
        try:
            subprocess.check_call(base + missing, stderr=subprocess.DEVNULL)
            missing = _still_missing(missing)
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass
    if missing:
        sys.exit(f"[FATAL] Could not install: {' '.join(missing)}")


_bootstrap(
    ("duckdb",     "duckdb"),
    ("pyarrow",    "pyarrow"),
    ("openpyxl",   "openpyxl"),
    ("matplotlib", "matplotlib"),
    ("numpy",      "numpy"),
    ("python-pptx", "pptx"),
)

import duckdb                           # noqa: E402
import matplotlib.pyplot as plt        # noqa: E402
import matplotlib.ticker as mticker    # noqa: E402
import numpy as np                     # noqa: E402
import pyarrow.parquet as pq           # noqa: E402
from openpyxl import Workbook          # noqa: E402
from openpyxl.styles import (Alignment, Border, Font, PatternFill, Side)  # noqa
from openpyxl.utils import get_column_letter  # noqa: E402
from pptx import Presentation          # noqa: E402
from pptx.dml.color import RGBColor    # noqa: E402
from pptx.util import Inches, Pt       # noqa: E402

from prodes_pipeline.data_quality import (
    LineageRecord,
    StageTimer,
    configure_json_logging,
    file_inventory,
    freshness_metrics,
    parquet_quality_profile,
    require_existing_dir,
    to_jsonable,
    validate_nonempty_files,
    write_run_report,
)
from prodes_pipeline.pipeline_contracts import ANALYTICS_EXPORT_CONTRACT, GEOPARQUET_CONTRACT

SEP = "=" * 70
DIV = "-" * 70
REPORT_DIR = REPORTS_DIR
OBS_LOG = configure_json_logging(REPORT_DIR / "observability.jsonl")
_DUCKDB = duckdb.connect(":memory:")

# ============================================================================
# CONFIG
# ============================================================================

CONFIG: dict[str, object] = {
    "geoparquet_dir": str(GPQ_DIR),
    "tables_dir":     str(TABLES_DIR),
    "chart_dpi":      300,         # journal-quality minimum
    "chart_fmt":      "png",
}

# ============================================================================
# BILINGUAL COPY  (zero language mixing â€” every string fully localized)
# ============================================================================

T: dict[str, dict[str, str]] = {
    # â”€â”€ Section labels â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    "title_suppression": {
        "pt": "SÃ©rie HistÃ³rica de SupressÃ£o de VegetaÃ§Ã£o Nativa",
        "en": "Historical Series of Natural Vegetation Suppression",
    },
    "title_cumulative": {
        "pt": "SupressÃ£o Acumulada de VegetaÃ§Ã£o Nativa",
        "en": "Cumulative Natural Vegetation Suppression",
    },
    "title_nv_a": {
        "pt": "VegetaÃ§Ã£o Nativa Remanescente â€” ParÃ¢metro A",
        "en": "Remaining Natural Vegetation â€” Parameter A",
    },
    "title_nv_b": {
        "pt": "VegetaÃ§Ã£o Nativa por Classe â€” ParÃ¢metro B",
        "en": "Natural Vegetation by Class â€” Parameter B",
    },
    "title_sv_extent": {
        "pt": "ExtensÃ£o Espacial Anual da VegetaÃ§Ã£o SecundÃ¡ria",
        "en": "Annual Spatial Extent of Secondary Vegetation",
    },
    "title_sv_increment": {
        "pt": "Incremento LÃ­quido Anual de VegetaÃ§Ã£o SecundÃ¡ria",
        "en": "Annual Net Increment of Secondary Vegetation",
    },
    "title_sv_a": {
        "pt": "VegetaÃ§Ã£o SecundÃ¡ria por Classe de Idade â€” ParÃ¢metro A",
        "en": "Secondary Vegetation by Age Class â€” Parameter A",
    },
    "title_sv_b": {
        "pt": "VegetaÃ§Ã£o SecundÃ¡ria por HistÃ³rico de Uso â€” ParÃ¢metro B",
        "en": "Secondary Vegetation by Land-Use History â€” Parameter B",
    },
    "title_muni": {
        "pt": "Matriz MunicÃ­pio Ã— Estado â€” Todos os ParÃ¢metros",
        "en": "Municipality Ã— State Matrix â€” All Parameters",
    },
    "title_methodology": {
        "pt": "Notas MetodolÃ³gicas",
        "en": "Methodological Notes",
    },
    # â”€â”€ Column headers â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    "col_year":         {"pt": "Ano",              "en": "Year"},
    "col_state":        {"pt": "Estado",           "en": "State"},
    "col_muni":         {"pt": "MunicÃ­pio",        "en": "Municipality"},
    "col_class":        {"pt": "Classe",           "en": "Class"},
    "col_age_class":    {"pt": "Classe de Idade",  "en": "Age Class"},
    "col_suppression":  {"pt": "SupressÃ£o (kmÂ²)",  "en": "Suppression (kmÂ²)"},
    "col_cumulative":   {"pt": "Acumulado (kmÂ²)",  "en": "Cumulative (kmÂ²)"},
    "col_nv_remaining": {"pt": "VN Remanescente (kmÂ²)", "en": "NV Remaining (kmÂ²)"},
    "col_nv_pct":       {"pt": "VN Remanescente (%)",   "en": "NV Remaining (%)"},
    "col_sv_extent":    {"pt": "VS ExtensÃ£o (kmÂ²)", "en": "SV Extent (kmÂ²)"},
    "col_sv_increment": {"pt": "VS Incremento (kmÂ²)", "en": "SV Increment (kmÂ²)"},
    "col_sv_area":      {"pt": "Ãrea VS (kmÂ²)",    "en": "SV Area (kmÂ²)"},
    "col_pct_total":    {"pt": "% do Total",        "en": "% of Total"},
    "col_yoy":          {"pt": "VariaÃ§Ã£o Anual (%)", "en": "Annual Change (%)"},
    # â”€â”€ Age class labels â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    "age_young":        {"pt": "Jovem (0â€“5 anos)",         "en": "Young (0â€“5 yr)"},
    "age_intermediate": {"pt": "IntermediÃ¡ria (5â€“15 anos)", "en": "Intermediate (5â€“15 yr)"},
    "age_mature":       {"pt": "Madura (â‰¥15 anos)",        "en": "Mature (â‰¥15 yr)"},
    # â”€â”€ Chart axis labels â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    "ax_year":       {"pt": "Ano",                         "en": "Year"},
    "ax_km2":        {"pt": "Ãrea (kmÂ²)",                  "en": "Area (kmÂ²)"},
    "ax_delta_km2":  {"pt": "Incremento LÃ­quido (kmÂ²)",    "en": "Net Increment (kmÂ²)"},
    "ax_cum_km2":    {"pt": "SupressÃ£o Acumulada (kmÂ²)",   "en": "Cumulative Suppression (kmÂ²)"},
    "ax_state":      {"pt": "Estado",                      "en": "State"},
    # â”€â”€ Source strings â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    "src_prodes": {
        "pt": "Fonte: INPE/PRODES. Calculado on-the-fly a partir dos dados GeoParquet.",
        "en": "Source: INPE/PRODES. Computed on-the-fly from primary GeoParquet files.",
    },
    "src_vs": {
        "pt": "Fonte: INPE/PRODES â€” VegetaÃ§Ã£o SecundÃ¡ria. Calculado a partir dos dados GeoParquet.",
        "en": "Source: INPE/PRODES â€” Secondary Vegetation. Computed from GeoParquet data.",
    },
    # â”€â”€ PPTX text â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    "pptx_lang_label": {"pt": "PT-BR", "en": "EN-US"},
    "pptx_subtitle": {
        "pt": "DinÃ¢mica EspaÃ§o-Temporal da SupressÃ£o de VegetaÃ§Ã£o Nativa\ne TendÃªncias de SucessÃ£o SecundÃ¡ria em Biomas Brasileiros",
        "en": "Spatial-Temporal Dynamics of Natural Vegetation Suppression\nand Secondary Succession Trends in Brazilian Biomes",
    },
    "pptx_credit": {
        "pt": "INPE/PRODES Â· MapBiomas Â· Imazon",
        "en": "INPE/PRODES Â· MapBiomas Â· Imazon",
    },
    "pptx_method_title": {
        "pt": "Notas MetodolÃ³gicas",
        "en": "Methodological Notes",
    },
    "pptx_method_body": {
        "pt": (
            "â€¢ SupressÃ£o: detecÃ§Ã£o por corte raso em imagem Ã³ptica (PRODES)\n"
            "â€¢ VegetaÃ§Ã£o secundÃ¡ria: mapeamento anual por sensoriamento remoto\n"
            "â€¢ Unidade espacial: municÃ­pio/estado (IBGE)\n"
            "â€¢ Unidade de Ã¡rea: kmÂ² (SIRGAS 2000)\n"
            "â€¢ PerÃ­odo: dados disponÃ­veis nos arquivos GeoParquet"
        ),
        "en": (
            "â€¢ Suppression: clear-cut detection via optical imagery (PRODES)\n"
            "â€¢ Secondary vegetation: annual mapping by remote sensing\n"
            "â€¢ Spatial unit: municipality/state (IBGE)\n"
            "â€¢ Area unit: kmÂ² (SIRGAS 2000)\n"
            "â€¢ Period: data available in GeoParquet files"
        ),
    },
}


def _t(key: str, lang: str) -> str:
    """Return the localized string for key/lang."""
    return T.get(key, {}).get(lang, key)

# ============================================================================
# EXCEL PALETTE  (The Economist Ã— Academic Journal)
# ============================================================================

_C_NAVY   = "1B3A4B"   # dark navy    â€” primary header
_C_FOREST = "2E7D32"   # forest green â€” secondary header
_C_WHITE  = "FFFFFF"
_C_ALT    = "F4F6F8"   # off-white alternating rows
_C_DARK   = "111111"
_C_MED    = "555555"
_C_RED    = "C0392B"
_C_GRN    = "2E7D32"
_C_BORDER = "CCCCCC"

_THIN   = Side(style="thin",   color=_C_BORDER)
_MEDIUM = Side(style="medium", color="888888")
_BORDER = Border(left=_THIN, right=_THIN, top=_THIN, bottom=_THIN)

def _fill(c: str) -> PatternFill:
    return PatternFill("solid", fgColor=c)

def _font(bold=False, size=10, color=_C_DARK, italic=False) -> Font:
    return Font(name="Calibri", bold=bold, size=size, color=color, italic=italic)

def _align(h="left", v="center", wrap=False) -> Alignment:
    return Alignment(horizontal=h, vertical=v, wrap_text=wrap)

def _mpl_color(c: str) -> str:
    return c if c.startswith("#") else f"#{c}"

def _hdr(ws, row: int, cols: list[str], col0: int = 1, color: str = _C_NAVY) -> None:
    for j, v in enumerate(cols, col0):
        c = ws.cell(row=row, column=j, value=v)
        c.font      = _font(bold=True, color=_C_WHITE, size=10)
        c.fill      = _fill(color)
        c.alignment = _align("center")
        c.border    = Border(bottom=Side(style="medium", color=_C_WHITE),
                             right=Side(style="thin", color=_C_WHITE))

def _row(ws, row: int, vals: list, col0: int = 1) -> None:
    bg = _C_ALT if row % 2 == 0 else _C_WHITE
    for j, v in enumerate(vals, col0):
        c = ws.cell(row=row, column=j, value=v)
        c.fill      = _fill(bg)
        c.border    = _BORDER
        c.alignment = _align("right" if isinstance(v, (int, float)) else "left")
        c.font      = _font()

def _src(ws, row: int, text: str) -> None:
    c = ws.cell(row=row, column=1, value=text)
    c.font = _font(italic=True, size=8, color=_C_MED)

def _title(ws, row: int, text: str) -> None:
    c = ws.cell(row=row, column=1, value=text)
    c.font = _font(bold=True, size=12, color=_C_NAVY)

def _widths(ws, w: dict[int, float]) -> None:
    for idx, width in w.items():
        ws.column_dimensions[get_column_letter(idx)].width = width

# ============================================================================
# MATPLOTLIB STYLE  (Publication / Journal Quality)
# ============================================================================

_CMAP_UNI  = "viridis"     # perceptually uniform, colorblind-safe
_CMAP_DIV  = "cividis"     # diverging variant
_FIG_DPI   = int(CONFIG["chart_dpi"])


def _pub_style() -> None:
    """Apply strict academic publication rcParams."""
    plt.rcParams.update({
        "font.family":        "sans-serif",
        "font.sans-serif":    ["Arial", "Helvetica Neue", "Liberation Sans", "DejaVu Sans"],
        "font.size":          10,
        "axes.titlesize":     11,
        "axes.labelsize":     10,
        "xtick.labelsize":    9,
        "ytick.labelsize":    9,
        "legend.fontsize":    9,
        "figure.facecolor":   "white",
        "axes.facecolor":     "white",
        "axes.grid":          True,
        "axes.grid.axis":     "y",
        "grid.color":         "#F0F0F0",
        "grid.linewidth":     0.8,
        "axes.axisbelow":     True,
    })


def _clean_ax(ax) -> None:
    """Remove top and right spines per academic convention."""
    ax.spines[["top", "right"]].set_visible(False)
    ax.spines["left"].set_color("#CCCCCC")
    ax.spines["bottom"].set_color("#CCCCCC")


def _save_chart(fig, name: str) -> Path:
    CHART_DIR.mkdir(parents=True, exist_ok=True)
    date_str = datetime.now().strftime("%Y-%m-%d")
    path = CHART_DIR / f"{name}_{date_str}.{CONFIG['chart_fmt']}"
    fig.savefig(str(path), dpi=_FIG_DPI, bbox_inches="tight",
                facecolor="white", format=CONFIG["chart_fmt"])
    plt.close(fig)
    return path


def _fig_to_buf(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=_FIG_DPI,
                bbox_inches="tight", facecolor="white")
    buf.seek(0)
    plt.close(fig)
    return buf

# ============================================================================
# SCHEMA PROBING  (flexible column detection)
# ============================================================================

#: Ordered candidates for each logical field (first found wins)
_AREA_CANDIDATES   = ("areakm", "area_km", "area_km2", "areakm2", "area_ha", "area")
_YEAR_CANDIDATES   = ("year", "yr_image", "yr", "ano", "ano_imagem", "data_year")
_STATE_CANDIDATES  = ("estado", "state", "uf", "sigla_uf", "nm_uf")
_MUNI_CANDIDATES   = ("municipio", "municipality", "nm_municipio", "nome_municipio")
_CLASS_CANDIDATES  = ("classname", "class", "categoria", "class_name", "land_use")
_AGE_CANDIDATES    = ("age", "idade", "age_yr", "anos")


def _probe(files: list[Path], candidates: tuple[str, ...]) -> str | None:
    """Return the first column from candidates found in any of the files."""
    for f in files[:5]:
        try:
            names_low = {n.lower(): n for n in pq.read_schema(str(f)).names}
            match = next((names_low[c] for c in candidates if c in names_low), None)
            if match:
                return match
        except Exception:
            pass
    return None


def _infer_km2_factor(files: list[Path], area_col: str) -> float:
    """
    Detect area unit from value magnitude:
      >500,000 â†’ mÂ²  (Ã·1e6)
      >5,000   â†’ ha  (Ã·100)
      else     â†’ kmÂ²  (Ã—1)
    """
    sample: list[float] = []
    for f in files[:3]:
        try:
            col = pq.read_table(str(f), columns=[area_col]).column(area_col).to_pylist()
            sample += [v for v in col if v and v > 0][:20]
        except Exception:
            pass
    if not sample:
        return 1.0
    med = sorted(sample)[len(sample) // 2]
    return 1 / 1_000_000 if med > 500_000 else (1 / 100 if med > 5_000 else 1.0)

# ============================================================================
# FILE DISCOVERY
# ============================================================================

_BIOME_NAMES  = {"Amazon Biome", "Legal Amazon", "Cerrado", "Caatinga",
                 "Pantanal", "Mata Atlantica", "Pampa"}
_DEFOR_KW     = ("deforestation", "desmatamento", "desmat")
_VS_KW        = ("vs_", "vegetacao_secundaria", "secondary_vegetation", "vegsec")
_AUX_SKIP     = ("border", "boundary", "hydrography", "indigenous",
                 "conservation_units", "settlement")


def _discover_suppression(gpq_dir: Path) -> dict[str, list[Path]]:
    """
    Discover natural vegetation suppression parquets, grouped by biome.
    Matches files containing deforestation keywords in their path,
    excluding known auxiliary layers.
    """
    biome_files: dict[str, list[Path]] = {}
    for pf in sorted(gpq_dir.rglob("*.parquet")):
        try:
            parts    = pf.relative_to(gpq_dir).parts
            path_low = "/".join(p.lower() for p in parts)
            if not any(k in path_low for k in _DEFOR_KW):
                continue
            if any(k in path_low for k in _AUX_SKIP):
                continue
            biome = next((p for p in parts if p in _BIOME_NAMES), None)
            if biome:
                biome_files.setdefault(biome, []).append(pf)
        except (ValueError, IndexError):
            pass
    return biome_files


def _discover_vs(gpq_dir: Path) -> list[Path]:
    """
    Discover secondary vegetation parquets (VS files).
    Matches files with 'vs_' prefix or 'vegetacao_secundaria' in path.
    """
    files: list[Path] = []
    for pf in sorted(gpq_dir.rglob("*.parquet")):
        try:
            name_low = pf.name.lower()
            path_low = str(pf).lower().replace("\\", "/")
            if any(k in name_low or k in path_low for k in _VS_KW):
                files.append(pf)
        except Exception:
            pass
    return files

# ============================================================================
# DUCKDB QUERY LAYER  (P1 â€“ P9)
# ============================================================================

class SchemaMap(NamedTuple):
    area:  str
    year:  str | None
    state: str | None
    muni:  str | None
    cls:   str | None
    age:   str | None
    factor: float


def _schema(files: list[Path]) -> SchemaMap | None:
    """Probe schema and return a SchemaMap, or None if area column not found."""
    area   = _probe(files, _AREA_CANDIDATES)
    if not area:
        return None
    factor = _infer_km2_factor(files, area)
    return SchemaMap(
        area   = area,
        year   = _probe(files, _YEAR_CANDIDATES),
        state  = _probe(files, _STATE_CANDIDATES),
        muni   = _probe(files, _MUNI_CANDIDATES),
        cls    = _probe(files, _CLASS_CANDIDATES),
        age    = _probe(files, _AGE_CANDIDATES),
        factor = factor,
    )


def _sql_paths(files: list[Path]) -> str:
    """Format file list for DuckDB read_parquet()."""
    return repr([str(f).replace("\\", "/") for f in files])


def _run(sql: str, label: str = "") -> list:
    """Execute a DuckDB query with comprehensive error handling."""
    try:
        return _DUCKDB.execute(sql).fetchall()
    except Exception as exc:
        if label:
            print(f"  [WARN] {label}: {exc}")
        return []


def _tuple_sort_key(row: tuple) -> tuple:
    """Sort rows with mixed None/string dimension values without changing data."""
    return tuple((v is None, "" if v is None else v) for v in row)


# â”€â”€ P1  Annual Suppression Series â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def p1_suppression_series(
    files: list[Path], sm: SchemaMap
) -> list[tuple]:
    """
    P1: S(t, s, m) = Î£_{p: year=t, state=s, muni=m} area_p

    Returns rows of (year, state, municipality, suppression_km2).
    If administrative columns absent, returns (year, suppression_km2).
    """
    paths = _sql_paths(files)
    sel   = [f'CAST("{sm.year}" AS INTEGER) AS year'] if sm.year else []
    grp   = ["year"] if sm.year else []

    if sm.state:
        sel.append(f'"{sm.state}" AS state'); grp.append("state")
    if sm.muni:
        sel.append(f'"{sm.muni}" AS municipality'); grp.append("municipality")

    sel.append(f'ROUND(SUM(CAST("{sm.area}" AS DOUBLE)) * {sm.factor}, 2) AS suppression_km2')

    where = f'WHERE "{sm.area}" IS NOT NULL AND CAST("{sm.area}" AS DOUBLE) > 0'
    if sm.year:
        where += f' AND CAST("{sm.year}" AS INTEGER) BETWEEN 2000 AND 2030'

    sql = f"""
        SELECT {", ".join(sel)}
        FROM   read_parquet({paths})
        {where}
        {("GROUP BY " + ", ".join(grp)) if grp else ""}
        ORDER  BY {", ".join(grp) if grp else "1"}
    """
    return _run(sql, "P1 suppression_series")


# â”€â”€ P2  Cumulative Suppression â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def p2_cumulative(p1_rows: list[tuple]) -> list[tuple]:
    """
    P2: C(t,s,m) = Î£_{Ï„â‰¤t} S(Ï„,s,m)

    Built from P1 results using a running sum window.
    Returns rows of (year, [state, [muni,]] suppression_km2, cumulative_km2).
    """
    if not p1_rows:
        return []

    # Determine row structure from width
    ncols = len(p1_rows[0])
    # Group by everything except last col (suppression_km2)
    groups: dict[tuple, list[tuple[int, float]]] = defaultdict(list)
    for row in p1_rows:
        key  = row[:-1]            # (year, [state, [muni]])
        year = int(key[0]) if key else 0
        km2  = float(row[-1])
        groups[key[1:]].append((year, km2))  # key[1:] = (state, muni) or ()

    result = []
    for group_key, year_vals in groups.items():
        year_vals.sort()
        cum = 0.0
        for year, km2 in year_vals:
            cum += km2
            result.append((*((year,) + group_key), round(km2, 2), round(cum, 2)))

    return sorted(result, key=_tuple_sort_key)


# â”€â”€ P3  Natural Vegetation Parameter A (stock estimate) â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def p3_nv_remaining(p2_rows: list[tuple]) -> list[tuple]:
    """
    P3: NV_A(t) = Ã‚â‚€ âˆ’ C(t)

    Ã‚â‚€ is estimated as the maximum cumulative value across all units
    (proxy for total suppressed forest at end of record) plus a 10%
    residual buffer representing the undetected suppression prior to
    the monitoring period.

    Returns rows of (year, [state,] nv_remaining_km2, nv_pct).
    """
    if not p2_rows:
        return []
    # Cumulative is always the last column; suppression is second-to-last
    max_cum = max(float(r[-1]) for r in p2_rows)
    a_hat   = max_cum * 1.10   # 10% buffer for pre-monitoring period

    result = []
    for r in p2_rows:
        cum       = float(r[-1])
        remaining = max(a_hat - cum, 0.0)
        pct       = round(remaining / a_hat * 100, 2) if a_hat else 0.0
        result.append((*r[:-2], round(remaining, 2), pct))

    return sorted(result, key=_tuple_sort_key)


# â”€â”€ P4  Natural Vegetation Parameter B (class partition) â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def p4_nv_by_class(
    files: list[Path], sm: SchemaMap
) -> list[tuple]:
    """
    P4: NV_B(t, c) = Î£_{p: year=t, class=c} area_p

    Returns rows of (year, class, suppression_km2, pct_of_year_total).
    """
    if not sm.cls or not sm.year:
        return []

    paths = _sql_paths(files)
    sql = f"""
        WITH annual_class AS (
            SELECT
                CAST("{sm.year}" AS INTEGER) AS year,
                COALESCE("{sm.cls}", 'Unknown') AS vegetation_class,
                SUM(CAST("{sm.area}" AS DOUBLE)) * {sm.factor} AS km2
            FROM read_parquet({paths})
            WHERE "{sm.area}" IS NOT NULL
              AND CAST("{sm.area}" AS DOUBLE) > 0
              AND CAST("{sm.year}" AS INTEGER) BETWEEN 2000 AND 2030
            GROUP BY year, vegetation_class
        ),
        totals AS (
            SELECT year, SUM(km2) AS year_total FROM annual_class GROUP BY year
        )
        SELECT
            ac.year,
            ac.vegetation_class,
            ROUND(ac.km2, 2) AS suppression_km2,
            ROUND(ac.km2 / NULLIF(t.year_total, 0) * 100, 2) AS pct_of_year
        FROM annual_class ac
        JOIN totals t USING (year)
        ORDER BY ac.year, ac.km2 DESC
    """
    return _run(sql, "P4 nv_by_class")


# â”€â”€ P5  Secondary Vegetation Annual Extent â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def p5_sv_extent(
    files: list[Path], sm: SchemaMap
) -> list[tuple]:
    """
    P5: VS(t, s) = Î£_{v: year=t, state=s} area_v

    Returns rows of (year, [state,] sv_extent_km2).
    """
    if not files:
        return []
    paths = _sql_paths(files)

    sel = [f'CAST("{sm.year}" AS INTEGER) AS year'] if sm.year else []
    grp = ["year"] if sm.year else []
    if sm.state:
        sel.append(f'"{sm.state}" AS state'); grp.append("state")
    sel.append(f'ROUND(SUM(CAST("{sm.area}" AS DOUBLE)) * {sm.factor}, 2) AS sv_extent_km2')

    where = f'WHERE "{sm.area}" IS NOT NULL AND CAST("{sm.area}" AS DOUBLE) > 0'
    if sm.year:
        where += f' AND CAST("{sm.year}" AS INTEGER) BETWEEN 2000 AND 2030'

    sql = f"""
        SELECT {", ".join(sel)}
        FROM   read_parquet({paths})
        {where}
        {("GROUP BY " + ", ".join(grp)) if grp else ""}
        ORDER  BY {", ".join(grp) if grp else "1"}
    """
    return _run(sql, "P5 sv_extent")


# â”€â”€ P6  Annual Net Increment of Secondary Vegetation â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def p6_sv_increment(p5_rows: list[tuple]) -> list[tuple]:
    """
    P6: Î”VS(t, s) = VS(t, s) âˆ’ VS(tâˆ’1, s)

    Computed via LAG applied to P5 results.
    Returns rows with additional column: net_increment_km2.
    """
    if not p5_rows:
        return []

    groups: dict[tuple, list[tuple[int, float]]] = defaultdict(list)
    for row in p5_rows:
        year   = int(row[0])
        extent = float(row[-1])
        key    = row[1:-1]   # state (or empty tuple)
        groups[key].append((year, extent))

    result = []
    for key, series in groups.items():
        series.sort()
        for i, (year, extent) in enumerate(series):
            prev      = series[i - 1][1] if i > 0 else extent
            increment = round(extent - prev, 2) if i > 0 else 0.0
            result.append((*((year,) + key), round(extent, 2), increment))

    return sorted(result, key=_tuple_sort_key)


# â”€â”€ P7  SV Parameter A â€” Age-Class Partition â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def p7_sv_by_age_class(
    files: list[Path], sm: SchemaMap, lang: str = "en"
) -> list[tuple]:
    """
    P7: VS_A(t, Î´) = Î£_{v: year=t, Î´_v=Î´} area_v
        Î´ âˆˆ {Young: 0â€“5yr, Intermediate: 5â€“15yr, Mature: â‰¥15yr}

    If an explicit age column is present, uses it; otherwise falls back
    to classname-based heuristic parsing.
    """
    if not files:
        return []
    paths = _sql_paths(files)

    if sm.age and sm.year:
        young_lbl = _t("age_young", lang)
        inter_lbl = _t("age_intermediate", lang)
        mat_lbl   = _t("age_mature", lang)
        sql = f"""
            SELECT
                CAST("{sm.year}" AS INTEGER) AS year,
                CASE
                    WHEN CAST("{sm.age}" AS DOUBLE) < 5  THEN '{young_lbl}'
                    WHEN CAST("{sm.age}" AS DOUBLE) < 15 THEN '{inter_lbl}'
                    ELSE '{mat_lbl}'
                END AS age_class,
                ROUND(SUM(CAST("{sm.area}" AS DOUBLE)) * {sm.factor}, 2) AS sv_area_km2
            FROM read_parquet({paths})
            WHERE "{sm.area}" IS NOT NULL
              AND CAST("{sm.area}" AS DOUBLE) > 0
              AND "{sm.age}" IS NOT NULL
              AND CAST("{sm.year}" AS INTEGER) BETWEEN 2000 AND 2030
            GROUP BY year, age_class
            ORDER BY year, age_class
        """
    elif sm.cls and sm.year:
        # Heuristic: parse numeric portion from classname (e.g., "sec_veg_3yr")
        young_lbl = _t("age_young", lang)
        inter_lbl = _t("age_intermediate", lang)
        mat_lbl   = _t("age_mature", lang)
        sql = f"""
            SELECT
                CAST("{sm.year}" AS INTEGER) AS year,
                CASE
                    WHEN TRY_CAST(REGEXP_EXTRACT("{sm.cls}", '\\d+') AS INTEGER) < 5  THEN '{young_lbl}'
                    WHEN TRY_CAST(REGEXP_EXTRACT("{sm.cls}", '\\d+') AS INTEGER) < 15 THEN '{inter_lbl}'
                    ELSE '{mat_lbl}'
                END AS age_class,
                ROUND(SUM(CAST("{sm.area}" AS DOUBLE)) * {sm.factor}, 2) AS sv_area_km2
            FROM read_parquet({paths})
            WHERE "{sm.area}" IS NOT NULL
              AND CAST("{sm.area}" AS DOUBLE) > 0
              AND CAST("{sm.year}" AS INTEGER) BETWEEN 2000 AND 2030
            GROUP BY year, age_class
            ORDER BY year, age_class
        """
    else:
        return []

    return _run(sql, "P7 sv_age_class")


# â”€â”€ P8  SV Parameter B â€” Land-Use History Partition â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def p8_sv_by_land_use(
    files: list[Path], sm: SchemaMap
) -> list[tuple]:
    """
    P8: VS_B(t, h) = Î£_{v: year=t, history=h} area_v

    Uses classname or class field as land-use history proxy.
    Returns rows of (year, land_use_class, sv_area_km2, pct_of_year).
    """
    if not sm.cls or not sm.year or not files:
        return []
    paths = _sql_paths(files)

    sql = f"""
        WITH by_class AS (
            SELECT
                CAST("{sm.year}" AS INTEGER) AS year,
                COALESCE(CAST("{sm.cls}" AS VARCHAR), 'Unknown') AS land_use_class,
                SUM(CAST("{sm.area}" AS DOUBLE)) * {sm.factor} AS km2
            FROM read_parquet({paths})
            WHERE "{sm.area}" IS NOT NULL
              AND CAST("{sm.area}" AS DOUBLE) > 0
              AND CAST("{sm.year}" AS INTEGER) BETWEEN 2000 AND 2030
            GROUP BY year, land_use_class
        ),
        totals AS (SELECT year, SUM(km2) AS yr_total FROM by_class GROUP BY year)
        SELECT
            bc.year,
            bc.land_use_class,
            ROUND(bc.km2, 2) AS sv_area_km2,
            ROUND(bc.km2 / NULLIF(t.yr_total, 0) * 100, 2) AS pct_of_year
        FROM by_class bc JOIN totals t USING (year)
        ORDER BY bc.year, bc.km2 DESC
    """
    return _run(sql, "P8 sv_land_use")


# â”€â”€ P9  Municipality Ã— State Cross-Tabulation â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

def p9_muni_state_matrix(
    files: list[Path], sm: SchemaMap, param_label: str
) -> list[tuple]:
    """
    P9: X_k(t, s, m) for parameter k

    Computes all parameters at municipalityâ€“state granularity.
    Returns rows of (year, state, municipality, suppression_km2,
                     cumulative_km2, pct_state, pct_national).
    """
    if not sm.year or not sm.state or not sm.muni or not files:
        return []
    paths = _sql_paths(files)

    sql = f"""
        WITH base AS (
            SELECT
                CAST("{sm.year}" AS INTEGER) AS year,
                "{sm.state}" AS state,
                "{sm.muni}" AS municipality,
                SUM(CAST("{sm.area}" AS DOUBLE)) * {sm.factor} AS km2
            FROM read_parquet({paths})
            WHERE "{sm.area}" IS NOT NULL
              AND CAST("{sm.area}" AS DOUBLE) > 0
              AND CAST("{sm.year}" AS INTEGER) BETWEEN 2000 AND 2030
              AND "{sm.state}" IS NOT NULL
              AND "{sm.muni}" IS NOT NULL
            GROUP BY year, state, municipality
        ),
        cumulative AS (
            SELECT *,
                SUM(km2) OVER (
                    PARTITION BY state, municipality ORDER BY year
                    ROWS BETWEEN UNBOUNDED PRECEDING AND CURRENT ROW
                ) AS cum_km2
            FROM base
        ),
        state_totals AS (
            SELECT year, state, SUM(km2) AS state_km2 FROM base GROUP BY year, state
        ),
        nat_totals AS (
            SELECT year, SUM(km2) AS nat_km2 FROM base GROUP BY year
        )
        SELECT
            c.year,
            c.state,
            c.municipality,
            ROUND(c.km2, 2)     AS suppression_km2,
            ROUND(c.cum_km2, 2) AS cumulative_km2,
            ROUND(c.km2 / NULLIF(s.state_km2, 0) * 100, 2) AS pct_state,
            ROUND(c.km2 / NULLIF(n.nat_km2,   0) * 100, 2) AS pct_national
        FROM      cumulative c
        JOIN state_totals s USING (year, state)
        JOIN nat_totals   n USING (year)
        ORDER BY c.year, c.state, c.km2 DESC
    """
    return _run(sql, f"P9 muni_state [{param_label}]")

# ============================================================================
# EXCEL SHEET BUILDERS
# ============================================================================

def _ws_p1(wb: Workbook, lang: str, rows: list[tuple]) -> None:
    ws = wb.create_sheet(_t("title_suppression", lang)[:31])
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "B3"
    _title(ws, 1, _t("title_suppression", lang))

    ncols = len(rows[0]) if rows else 0
    if ncols >= 4:
        headers = [_t("col_year", lang), _t("col_state", lang),
                   _t("col_muni", lang), _t("col_suppression", lang)]
    elif ncols == 3:
        headers = [_t("col_year", lang), _t("col_state", lang), _t("col_suppression", lang)]
    else:
        headers = [_t("col_year", lang), _t("col_suppression", lang)]

    _hdr(ws, 2, headers)
    for i, r in enumerate(rows, 3):
        _row(ws, i, list(r))
    _src(ws, len(rows) + 4, _t("src_prodes", lang))
    _widths(ws, {j + 1: 18 for j in range(len(headers))})


def _ws_p2(wb: Workbook, lang: str, rows: list[tuple]) -> None:
    ws = wb.create_sheet(_t("title_cumulative", lang)[:31])
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "B3"
    _title(ws, 1, _t("title_cumulative", lang))

    ncols = len(rows[0]) if rows else 0
    if ncols >= 5:
        headers = [_t("col_year", lang), _t("col_state", lang),
                   _t("col_muni", lang), _t("col_suppression", lang), _t("col_cumulative", lang)]
    elif ncols == 4:
        headers = [_t("col_year", lang), _t("col_state", lang),
                   _t("col_suppression", lang), _t("col_cumulative", lang)]
    else:
        headers = [_t("col_year", lang), _t("col_suppression", lang), _t("col_cumulative", lang)]

    _hdr(ws, 2, headers)
    for i, r in enumerate(rows, 3):
        _row(ws, i, list(r))
    _src(ws, len(rows) + 4, _t("src_prodes", lang))
    _widths(ws, {j + 1: 18 for j in range(len(headers))})


def _ws_p3(wb: Workbook, lang: str, rows: list[tuple]) -> None:
    ws = wb.create_sheet(_t("title_nv_a", lang)[:31])
    ws.sheet_view.showGridLines = False
    _title(ws, 1, _t("title_nv_a", lang))

    ncols = len(rows[0]) if rows else 2
    base_h = [_t("col_year", lang), _t("col_state", lang), _t("col_muni", lang)][:ncols - 2]
    headers = base_h + [_t("col_nv_remaining", lang), _t("col_nv_pct", lang)]
    _hdr(ws, 2, headers)
    for i, r in enumerate(rows, 3):
        _row(ws, i, list(r))
        pct = r[-1]
        color = _C_RED if isinstance(pct, float) and pct < 50 else _C_GRN
        ws.cell(row=i, column=len(headers)).font = _font(bold=True, color=color)
    _src(ws, len(rows) + 4, _t("src_prodes", lang))
    _widths(ws, {j + 1: 18 for j in range(len(headers))})


def _ws_p4(wb: Workbook, lang: str, rows: list[tuple]) -> None:
    ws = wb.create_sheet(_t("title_nv_b", lang)[:31])
    ws.sheet_view.showGridLines = False
    _title(ws, 1, _t("title_nv_b", lang))
    headers = [_t("col_year", lang), _t("col_class", lang),
               _t("col_suppression", lang), _t("col_pct_total", lang)]
    _hdr(ws, 2, headers, color=_C_FOREST)
    for i, r in enumerate(rows, 3):
        _row(ws, i, list(r))
    _src(ws, len(rows) + 4, _t("src_prodes", lang))
    _widths(ws, {1: 8, 2: 35, 3: 20, 4: 16})


def _ws_p5(wb: Workbook, lang: str, rows: list[tuple]) -> None:
    ws = wb.create_sheet(_t("title_sv_extent", lang)[:31])
    ws.sheet_view.showGridLines = False
    _title(ws, 1, _t("title_sv_extent", lang))
    ncols = len(rows[0]) if rows else 2
    base_h = ([_t("col_year", lang), _t("col_state", lang)] if ncols > 2
              else [_t("col_year", lang)])
    headers = base_h + [_t("col_sv_extent", lang)]
    _hdr(ws, 2, headers, color=_C_FOREST)
    for i, r in enumerate(rows, 3):
        _row(ws, i, list(r))
    _src(ws, len(rows) + 4, _t("src_vs", lang))
    _widths(ws, {j + 1: 18 for j in range(len(headers))})


def _ws_p6(wb: Workbook, lang: str, rows: list[tuple]) -> None:
    ws = wb.create_sheet(_t("title_sv_increment", lang)[:31])
    ws.sheet_view.showGridLines = False
    _title(ws, 1, _t("title_sv_increment", lang))
    ncols = len(rows[0]) if rows else 3
    base_h = ([_t("col_year", lang), _t("col_state", lang)] if ncols > 3
              else [_t("col_year", lang)])
    headers = base_h + [_t("col_sv_extent", lang), _t("col_sv_increment", lang)]
    _hdr(ws, 2, headers, color=_C_FOREST)
    for i, r in enumerate(rows, 3):
        _row(ws, i, list(r))
        inc = r[-1]
        if isinstance(inc, (int, float)):
            color = _C_GRN if inc >= 0 else _C_RED
            ws.cell(row=i, column=len(headers)).font = _font(bold=True, color=color)
    _src(ws, len(rows) + 4, _t("src_vs", lang))
    _widths(ws, {j + 1: 18 for j in range(len(headers))})


def _ws_p7(wb: Workbook, lang: str, rows: list[tuple]) -> None:
    ws = wb.create_sheet(_t("title_sv_a", lang)[:31])
    ws.sheet_view.showGridLines = False
    _title(ws, 1, _t("title_sv_a", lang))
    headers = [_t("col_year", lang), _t("col_age_class", lang), _t("col_sv_area", lang)]
    _hdr(ws, 2, headers, color=_C_FOREST)
    for i, r in enumerate(rows, 3):
        _row(ws, i, list(r))
    _src(ws, len(rows) + 4, _t("src_vs", lang))
    _widths(ws, {1: 8, 2: 30, 3: 20})


def _ws_p8(wb: Workbook, lang: str, rows: list[tuple]) -> None:
    ws = wb.create_sheet(_t("title_sv_b", lang)[:31])
    ws.sheet_view.showGridLines = False
    _title(ws, 1, _t("title_sv_b", lang))
    headers = [_t("col_year", lang), "Land-Use Class / Classe de Uso",
               _t("col_sv_area", lang), _t("col_pct_total", lang)]
    _hdr(ws, 2, headers, color=_C_FOREST)
    for i, r in enumerate(rows, 3):
        _row(ws, i, list(r))
    _src(ws, len(rows) + 4, _t("src_vs", lang))
    _widths(ws, {1: 8, 2: 40, 3: 20, 4: 16})


def _ws_p9(wb: Workbook, lang: str, rows: list[tuple]) -> None:
    ws = wb.create_sheet(_t("title_muni", lang)[:31])
    ws.sheet_view.showGridLines = False
    ws.freeze_panes = "D3"
    _title(ws, 1, _t("title_muni", lang))
    headers = [_t("col_year", lang), _t("col_state", lang), _t("col_muni", lang),
               _t("col_suppression", lang), _t("col_cumulative", lang),
               "% Estado / State", "% Nacional / National"]
    _hdr(ws, 2, headers)
    for i, r in enumerate(rows[:5000], 3):   # cap at 5000 rows for xlsx perf
        _row(ws, i, list(r))
    if len(rows) > 5000:
        ws.cell(row=5004, column=1,
                value=f"[Truncado para 5.000 linhas de {len(rows)} / "
                      f"Truncated to 5,000 of {len(rows)} rows]").font = _font(italic=True)
    _src(ws, min(len(rows), 5000) + 4, _t("src_prodes", lang))
    _widths(ws, {1: 8, 2: 14, 3: 30, 4: 20, 5: 20, 6: 14, 7: 16})


def _ws_methodology(wb: Workbook, lang: str) -> None:
    ws = wb.create_sheet(_t("title_methodology", lang)[:31])
    ws.sheet_view.showGridLines = False
    _title(ws, 1, _t("title_methodology", lang))

    notes = {
        "pt": [
            ("P1 â€” SÃ©rie de SupressÃ£o Anual",
             "S(t,s,m) = Î£ area dos polÃ­gonos de supressÃ£o para o ano t, "
             "estado s e municÃ­pio m. Unidade: kmÂ². Fonte: INPE/PRODES."),
            ("P2 â€” SupressÃ£o Acumulada",
             "C(t,s,m) = Î£_{Ï„â‰¤t} S(Ï„,s,m). Calculada via funÃ§Ã£o de janela SQL "
             "(SUM OVER PARTITION BY estado, municÃ­pio ORDER BY ano)."),
            ("P3 â€” VegetaÃ§Ã£o Nativa Remanescente (A)",
             "NV_A(t) = Ã‚â‚€ âˆ’ C(t), onde Ã‚â‚€ = C(t_max)Ã—1,10 Ã© o estimador "
             "do estoque florestal no inÃ­cio do perÃ­odo monitorado."),
            ("P4 â€” VegetaÃ§Ã£o Nativa por Classe (B)",
             "PartiÃ§Ã£o de S(t) pela classe de uso/cobertura detectada."),
            ("P5 â€” ExtensÃ£o Anual de VegetaÃ§Ã£o SecundÃ¡ria",
             "VS(t,s) = Î£ Ã¡rea dos polÃ­gonos de VS para o ano t e estado s."),
            ("P6 â€” Incremento LÃ­quido Anual de VS",
             "Î”VS(t,s) = VS(t,s) âˆ’ VS(tâˆ’1,s). Positivo = recuperaÃ§Ã£o lÃ­quida."),
            ("P7 â€” VS por Classe de Idade (A)",
             "PartiÃ§Ã£o por idade: Jovem (<5yr), IntermediÃ¡ria (5â€“15yr), Madura (â‰¥15yr)."),
            ("P8 â€” VS por HistÃ³rico de Uso (B)",
             "PartiÃ§Ã£o por classe/histÃ³rico de uso disponÃ­vel no atributo 'classname'."),
            ("P9 â€” Matriz MunicÃ­pio Ã— Estado",
             "Todos os parÃ¢metros calculados na escala municipal e "
             "agregados ao nÃ­vel estadual. Inclui % do total estadual e nacional."),
            ("Dados",
             f"Computado on-the-fly a partir dos GeoParquets primÃ¡rios. "
             f"Gerado em: {datetime.now().strftime('%Y-%m-%d %H:%M')}"),
        ],
        "en": [
            ("P1 â€” Annual Suppression Series",
             "S(t,s,m) = Î£ area of suppression polygons for year t, state s, "
             "municipality m. Unit: kmÂ². Source: INPE/PRODES."),
            ("P2 â€” Cumulative Suppression",
             "C(t,s,m) = Î£_{Ï„â‰¤t} S(Ï„,s,m). Computed via SQL window function "
             "(SUM OVER PARTITION BY state, municipality ORDER BY year)."),
            ("P3 â€” Remaining Natural Vegetation (A)",
             "NV_A(t) = Ã‚â‚€ âˆ’ C(t), where Ã‚â‚€ = C(t_max)Ã—1.10 is the "
             "estimated forest stock at the start of the monitored period."),
            ("P4 â€” Natural Vegetation by Class (B)",
             "Partition of S(t) by detected land-cover/use class."),
            ("P5 â€” Annual Secondary Vegetation Extent",
             "VS(t,s) = Î£ area of SV polygons for year t and state s."),
            ("P6 â€” Annual Net Increment of SV",
             "Î”VS(t,s) = VS(t,s) âˆ’ VS(tâˆ’1,s). Positive = net recovery."),
            ("P7 â€” SV by Age Class (A)",
             "Partition by age: Young (<5yr), Intermediate (5â€“15yr), Mature (â‰¥15yr)."),
            ("P8 â€” SV by Land-Use History (B)",
             "Partition by class/land-use history from the 'classname' attribute."),
            ("P9 â€” Municipality Ã— State Matrix",
             "All parameters at municipal scale, aggregated to state level. "
             "Includes % of state total and % of national total."),
            ("Data",
             f"Computed on-the-fly from primary GeoParquet files. "
             f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"),
        ],
    }[lang]

    row = 3
    for title, body in notes:
        ws.cell(row=row, column=1, value=title).font = _font(bold=True, color=_C_NAVY, size=11)
        ws.cell(row=row + 1, column=1, value=body).font = _font(color=_C_DARK, size=10)
        ws.cell(row=row + 1, column=1).alignment = _align(wrap=True)
        ws.row_dimensions[row + 1].height = 40
        row += 3
    ws.column_dimensions["A"].width = 95

# ============================================================================
# MATPLOTLIB CHART FUNCTIONS  (publication quality, journal ready)
# ============================================================================

def chart_suppression_trend(p1_national: list[tuple], lang: str) -> Path:
    """Bar chart of annual national suppression with viridis color scale."""
    _pub_style()
    years = [int(r[0]) for r in p1_national if len(r) == 2]
    km2   = [float(r[1]) for r in p1_national if len(r) == 2]
    if not years:
        return None

    cmap   = plt.get_cmap(_CMAP_UNI, len(years))
    colors = [cmap(i / max(len(years) - 1, 1)) for i in range(len(years))]

    fig, ax = plt.subplots(figsize=(10, 5))
    bars = ax.bar(years, km2, color=colors, width=0.75, zorder=3, linewidth=0)
    _clean_ax(ax)

    ax.set_xlabel(_t("ax_year", lang), labelpad=8)
    ax.set_ylabel(_t("ax_km2", lang), labelpad=8)
    ax.set_title(_t("title_suppression", lang), fontsize=12, fontweight="bold",
                 loc="left", pad=10)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:,.0f}"))

    # Annotate peak and most recent
    if km2:
        peak_i = km2.index(max(km2))
        ax.annotate(f"{km2[peak_i]:,.0f}",
                    xy=(years[peak_i], km2[peak_i]),
                    xytext=(0, 6), textcoords="offset points",
                    ha="center", fontsize=8, color="#C0392B", fontweight="bold")

    ax.text(0.01, -0.12, _t("src_prodes", lang),
            transform=ax.transAxes, fontsize=7, color="#888888", style="italic")

    plt.tight_layout()
    return _save_chart(fig, f"suppression_trend_{lang}")


def chart_cumulative_suppression(p2_national: list[tuple], lang: str) -> Path:
    """Cumulative suppression line chart with shaded area."""
    _pub_style()
    # Expect rows: (year, suppression, cumulative)  OR  (year, cumulative)
    rows_2col = [r for r in p2_national if len(r) == 2]
    rows_3col = [r for r in p2_national if len(r) == 3]
    if rows_3col:
        years = [int(r[0]) for r in rows_3col]
        cum   = [float(r[2]) for r in rows_3col]
    elif rows_2col:
        years = [int(r[0]) for r in rows_2col]
        cum   = list(np.cumsum([float(r[1]) for r in rows_2col]))
    else:
        return None

    fig, ax = plt.subplots(figsize=(10, 5))
    ax.fill_between(years, cum, alpha=0.15, color="#C0392B")
    ax.plot(years, cum, color="#C0392B", lw=2.5, zorder=4)
    ax.scatter(years, cum, color="#C0392B", s=30, zorder=5)
    _clean_ax(ax)

    ax.set_xlabel(_t("ax_year", lang), labelpad=8)
    ax.set_ylabel(_t("ax_cum_km2", lang), labelpad=8)
    ax.set_title(_t("title_cumulative", lang), fontsize=12, fontweight="bold",
                 loc="left", pad=10)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:,.0f}"))
    ax.text(0.01, -0.12, _t("src_prodes", lang),
            transform=ax.transAxes, fontsize=7, color="#888888", style="italic")

    plt.tight_layout()
    return _save_chart(fig, f"cumulative_suppression_{lang}")


def chart_sv_dynamics(p5_national: list[tuple], lang: str) -> Path:
    """Secondary vegetation extent â€” bar chart with cividis color scale."""
    _pub_style()
    years = [int(r[0]) for r in p5_national if len(r) == 2]
    km2   = [float(r[1]) for r in p5_national if len(r) == 2]
    if not years:
        return None

    cmap   = plt.get_cmap(_CMAP_DIV, len(years))
    colors = [cmap(i / max(len(years) - 1, 1)) for i in range(len(years))]

    fig, ax = plt.subplots(figsize=(10, 5))
    ax.bar(years, km2, color=colors, width=0.75, zorder=3, linewidth=0)
    _clean_ax(ax)

    ax.set_xlabel(_t("ax_year", lang), labelpad=8)
    ax.set_ylabel(_t("ax_km2", lang), labelpad=8)
    ax.set_title(_t("title_sv_extent", lang), fontsize=12, fontweight="bold",
                 loc="left", pad=10)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:,.0f}"))
    ax.text(0.01, -0.12, _t("src_vs", lang),
            transform=ax.transAxes, fontsize=7, color="#888888", style="italic")

    plt.tight_layout()
    return _save_chart(fig, f"sv_dynamics_{lang}")


def chart_sv_increment(p6_national: list[tuple], lang: str) -> Path:
    """Net SV increment â€” diverging bar chart (green=positive, red=negative)."""
    _pub_style()
    years = [int(r[0]) for r in p6_national if len(r) >= 3]
    inc   = [float(r[-1]) for r in p6_national if len(r) >= 3]
    if not years:
        return None

    colors = [_mpl_color(_C_GRN if v >= 0 else _C_RED) for v in inc]
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.bar(years, inc, color=colors, width=0.75, zorder=3, linewidth=0)
    ax.axhline(0, color="#888888", lw=0.8, zorder=2)
    _clean_ax(ax)

    ax.set_xlabel(_t("ax_year", lang), labelpad=8)
    ax.set_ylabel(_t("ax_delta_km2", lang), labelpad=8)
    ax.set_title(_t("title_sv_increment", lang), fontsize=12, fontweight="bold",
                 loc="left", pad=10)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:+,.0f}"))
    ax.text(0.01, -0.12, _t("src_vs", lang),
            transform=ax.transAxes, fontsize=7, color="#888888", style="italic")

    plt.tight_layout()
    return _save_chart(fig, f"sv_increment_{lang}")


def chart_state_ranking(p1_rows: list[tuple], lang: str, top_n: int = 15) -> Path:
    """Horizontal bar â€” top states by cumulative suppression (most recent year)."""
    _pub_style()
    # Filter for rows with state column (ncols==3: year, state, km2)
    state_rows = [r for r in p1_rows if len(r) == 3]
    if not state_rows:
        return None

    max_year = max(int(r[0]) for r in state_rows)
    year_rows = [(str(r[1]), float(r[2])) for r in state_rows if int(r[0]) == max_year]
    year_rows.sort(key=lambda x: x[1], reverse=True)
    year_rows = year_rows[:top_n]
    states, km2 = zip(*year_rows) if year_rows else ([], [])

    cmap   = plt.get_cmap(_CMAP_UNI, len(states))
    colors = [cmap(i / max(len(states) - 1, 1)) for i in range(len(states))]

    fig, ax = plt.subplots(figsize=(9, max(4, len(states) * 0.45)))
    bars = ax.barh(range(len(states)), km2, color=colors, height=0.7, zorder=3)
    ax.set_yticks(range(len(states)))
    ax.set_yticklabels(states, fontsize=9)
    ax.invert_yaxis()
    _clean_ax(ax)
    ax.spines["left"].set_visible(False)
    ax.set_xlabel(_t("ax_km2", lang), labelpad=8)
    ax.set_title(f"{_t('title_suppression', lang)} â€” {max_year}",
                 fontsize=11, fontweight="bold", loc="left", pad=10)
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:,.0f}"))

    for i, v in enumerate(km2):
        ax.text(v + max(km2) * 0.01, i, f"{v:,.0f}",
                va="center", fontsize=8, color=_mpl_color(_C_MED))

    ax.text(0.01, -0.08, _t("src_prodes", lang),
            transform=ax.transAxes, fontsize=7, color="#888888", style="italic")
    plt.tight_layout()
    return _save_chart(fig, f"state_ranking_{lang}")


def chart_sv_subclass(p7_rows: list[tuple], lang: str) -> Path:
    """Stacked bar â€” SV by age class over time (viridis)."""
    _pub_style()
    if not p7_rows:
        return None

    from collections import defaultdict
    by_year_class: dict[int, dict[str, float]] = defaultdict(lambda: defaultdict(float))
    all_classes: set[str] = set()
    for r in p7_rows:
        if len(r) < 3:
            continue
        yr, cls, km2 = int(r[0]), str(r[1]), float(r[2])
        by_year_class[yr][cls] += km2
        all_classes.add(cls)

    years   = sorted(by_year_class)
    classes = sorted(all_classes)
    cmap    = plt.get_cmap(_CMAP_UNI, max(len(classes), 3))
    colors  = [cmap(i / max(len(classes) - 1, 1)) for i in range(len(classes))]

    fig, ax = plt.subplots(figsize=(11, 5))
    bottom  = np.zeros(len(years))
    for ci, cls in enumerate(classes):
        vals = np.array([by_year_class[y].get(cls, 0.0) for y in years])
        ax.bar(years, vals, bottom=bottom, color=colors[ci],
               label=cls, width=0.75, zorder=3, linewidth=0)
        bottom += vals

    _clean_ax(ax)
    ax.set_xlabel(_t("ax_year", lang), labelpad=8)
    ax.set_ylabel(_t("ax_km2", lang), labelpad=8)
    ax.set_title(_t("title_sv_a", lang), fontsize=12, fontweight="bold",
                 loc="left", pad=10)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:,.0f}"))
    ax.legend(fontsize=8, loc="upper left", frameon=False)
    ax.text(0.01, -0.12, _t("src_vs", lang),
            transform=ax.transAxes, fontsize=7, color="#888888", style="italic")

    plt.tight_layout()
    return _save_chart(fig, f"sv_subclass_{lang}")

# ============================================================================
# POWERPOINT SLIDE BUILDERS  (bilingual, academic layout)
# ============================================================================

_PPT_W = 10.0   # inches
_PPT_H = 5.625  # 16:9


def _in(x: float) -> int:
    return int(x * 914_400)


def _rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _pptx_tb(slide, text, l, t, w, h, *, size=11, bold=False,
             color=_C_DARK, italic=False, align=1, wrap=True):
    from pptx.util import Pt
    box = slide.shapes.add_textbox(_in(l), _in(t), _in(w), _in(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(str(text).split("\n")):
        p   = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text          = line
        run.font.size     = Pt(size)
        run.font.bold     = bold
        run.font.italic   = italic
        run.font.color.rgb = _rgb(color)
    return box


def _pptx_rect(slide, l, t, w, h, fill_hex):
    sh = slide.shapes.add_shape(1, _in(l), _in(t), _in(w), _in(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = _rgb(fill_hex)
    sh.line.fill.background()
    return sh


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _lang_colors(lang: str) -> tuple[str, str]:
    """Return (primary, accent) hex colors for the language block."""
    return ("1B5E20", "E8F5E9") if lang == "pt" else ("0D47A1", "E3F2FD")


def _slide_cover(prs: Presentation, lang: str) -> None:
    sl     = _blank_slide(prs)
    pri, _ = _lang_colors(lang)
    _pptx_rect(sl, 0, 0, _PPT_W, 0.07, pri)
    _pptx_rect(sl, 0, _PPT_H - 0.07, _PPT_W, 0.07, pri)
    _pptx_tb(sl, _t("pptx_lang_label", lang), 0.35, 0.12, 2, 0.28,
             size=7, bold=True, color=pri)
    _pptx_tb(sl, "PRODES Â· Imazon", 0.35, 0.45, 8.5, 0.5,
             size=22, bold=True, color=_C_NAVY)
    _pptx_tb(sl, _t("pptx_subtitle", lang), 0.35, 1.05, 7.8, 1.4,
             size=13, color=_C_MED)
    _pptx_rect(sl, 0.35, 2.6, 2.5, 0.05, pri)
    _pptx_tb(sl, _t("pptx_credit", lang), 0.35, 2.75, 8.5, 0.4,
             size=9, italic=True, color=_C_MED)
    year = datetime.now().year
    _pptx_tb(sl, str(year), 0.35, 3.15, 8.5, 0.35, size=9, color="#AAAAAA")


def _slide_chart(prs: Presentation, lang: str,
                 chart_path: Path | None, title: str,
                 subtitle: str, source: str) -> None:
    """Chart slide with academic sidebar for methodology notes."""
    sl     = _blank_slide(prs)
    pri, acc = _lang_colors(lang)
    _pptx_rect(sl, 0, 0, _PPT_W, 0.06, pri)
    _pptx_rect(sl, 0, _PPT_H - 0.06, _PPT_W, 0.06, pri)

    # Sidebar
    _pptx_rect(sl, 7.6, 0.06, 2.4, _PPT_H - 0.12, acc)
    _pptx_tb(sl, _t("pptx_method_title", lang), 7.7, 0.2, 2.2, 0.35,
             size=8, bold=True, color=pri)
    _pptx_tb(sl, _t("pptx_method_body", lang), 7.7, 0.55, 2.2, 3.5,
             size=7, color=_C_MED, wrap=True)

    # Main area
    _pptx_tb(sl, title, 0.25, 0.1, 7.2, 0.4, size=12, bold=True, color=_C_DARK)
    _pptx_tb(sl, subtitle, 0.25, 0.5, 7.2, 0.28, size=8, color=_C_MED)

    if chart_path and chart_path.exists():
        try:
            sl.shapes.add_picture(str(chart_path), _in(0.2), _in(0.82),
                                  width=_in(7.2))
        except Exception:
            _pptx_tb(sl, "[Chart unavailable]", 0.2, 0.82, 7.2, 3.5,
                     size=9, color="#AAAAAA", italic=True)
    else:
        _pptx_tb(sl, "[No data available for this parameter]", 0.2, 2.5, 7.2, 0.5,
                 size=9, color="#AAAAAA", italic=True)

    _pptx_tb(sl, source, 0.25, 5.3, 7.2, 0.25, size=6.5, italic=True, color="#AAAAAA")


def _slide_table_summary(prs: Presentation, lang: str,
                         rows: list[tuple], title: str,
                         headers: list[str], source: str) -> None:
    """Text-table summary slide (top 8 rows for readability)."""
    sl     = _blank_slide(prs)
    pri, acc = _lang_colors(lang)
    _pptx_rect(sl, 0, 0, _PPT_W, 0.06, pri)
    _pptx_rect(sl, 0, _PPT_H - 0.06, _PPT_W, 0.06, pri)
    _pptx_tb(sl, title, 0.25, 0.1, 9.5, 0.4, size=12, bold=True, color=_C_DARK)
    _pptx_tb(sl, source, 0.25, 5.3, 9.5, 0.25, size=6.5, italic=True, color="#AAAAAA")

    display_rows = rows[:8]
    if not display_rows:
        _pptx_tb(sl, "[No data]", 0.25, 1.0, 9.5, 0.4, size=10, color="#AAAAAA")
        return

    col_w = 9.5 / max(len(headers), 1)
    _pptx_rect(sl, 0.25, 0.58, 9.5, 0.38, pri)
    for j, h in enumerate(headers):
        _pptx_tb(sl, h, 0.25 + j * col_w, 0.6, col_w, 0.34,
                 size=8.5, bold=True, color=_C_WHITE, align=2)

    for i, row in enumerate(display_rows):
        y = 0.96 + i * 0.44
        bg = "F4F6F8" if i % 2 == 0 else _C_WHITE
        _pptx_rect(sl, 0.25, y, 9.5, 0.42, bg)
        for j, val in enumerate(row[:len(headers)]):
            _pptx_tb(sl, str(round(val, 2) if isinstance(val, float) else val),
                     0.25 + j * col_w, y + 0.04, col_w, 0.36,
                     size=8, color=_C_DARK,
                     align=3 if isinstance(val, (int, float)) else 1)


def build_pptx(lang: str, chart_paths: dict[str, Path | None],
               data: dict) -> Presentation:
    """Construct the full bilingual PPTX for one language block."""
    prs = Presentation()
    prs.slide_width  = _in(_PPT_W)
    prs.slide_height = _in(_PPT_H)

    _slide_cover(prs, lang)

    specs = [
        ("suppression_trend",     "title_suppression",   "src_prodes"),
        ("cumulative_suppression", "title_cumulative",    "src_prodes"),
        ("sv_dynamics",            "title_sv_extent",     "src_vs"),
        ("sv_increment",           "title_sv_increment",  "src_vs"),
        ("state_ranking",          "title_suppression",   "src_prodes"),
        ("sv_subclass",            "title_sv_a",          "src_vs"),
    ]
    for key, title_key, src_key in specs:
        _slide_chart(
            prs, lang,
            chart_paths.get(f"{key}_{lang}"),
            _t(title_key, lang),
            f"INPE/PRODES Â· Imazon  Â·  {datetime.now().year}",
            _t(src_key, lang),
        )

    # P9 summary table slide
    p9 = data.get("p9", [])[:8]
    if p9:
        headers = [_t("col_year", lang), _t("col_state", lang), _t("col_muni", lang),
                   _t("col_suppression", lang), _t("col_cumulative", lang)]
        _slide_table_summary(prs, lang, p9, _t("title_muni", lang),
                             headers, _t("src_prodes", lang))

    return prs

# ============================================================================
# MAIN ORCHESTRATOR
# ============================================================================

def main() -> None:
    ensure_pipeline_dirs()
    now = datetime.now().strftime("%Y-%m-%d %H:%M")
    print(f"\n{SEP}")
    print(f"  PRODES Analytics Pipeline  v{__version__}  |  {now}")
    print(f"{SEP}\n")

    gpq_dir = Path(str(CONFIG["geoparquet_dir"]))
    require_existing_dir(gpq_dir, "GeoParquet")

    TABLES_DIR.mkdir(parents=True, exist_ok=True)
    CHART_DIR.mkdir(parents=True, exist_ok=True)

    # â”€â”€ 1. Discover files â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    print("  [1/5] Discovering data files...")
    discover_timer = StageTimer("06_discover_inputs")
    supp_by_biome = _discover_suppression(gpq_dir)
    vs_files      = _discover_vs(gpq_dir)

    all_supp = [f for fs in supp_by_biome.values() for f in fs]
    amazon_files = [f for bd, fs in supp_by_biome.items()
                    if bd in {"Amazon Biome", "Legal Amazon"} for f in fs]

    print(f"     Suppression files  : {len(all_supp)} ({len(supp_by_biome)} biomes)")
    print(f"     Secondary veg files: {len(vs_files)}")
    all_input_files = sorted(set(all_supp + vs_files))
    input_quality = {
        "contract": to_jsonable(GEOPARQUET_CONTRACT),
        "inventory": file_inventory(all_input_files),
        "freshness": freshness_metrics(all_input_files, GEOPARQUET_CONTRACT.freshness),
        "parquet_profile": parquet_quality_profile(
            all_input_files,
            GEOPARQUET_CONTRACT,
        ),
    }
    OBS_LOG.emit(
        "stage_metrics",
        **to_jsonable(
            discover_timer.finish(
                "ok",
                output_row_count=len(all_input_files),
                anomalies={
                    "schema": input_quality["parquet_profile"].get(
                        "schema_anomalies", []
                    ),
                    "distribution": input_quality["parquet_profile"].get(
                        "distribution_anomalies", []
                    ),
                    "freshness": input_quality["freshness"].get("stale", []),
                },
            )
        ),
    )

    if not all_supp:
        sys.exit("[FATAL] No suppression parquet files found. Run scripts 02 and 05 first.")

    # â”€â”€ 2. Schema probing â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    print("\n  [2/5] Probing schemas...")
    sm_supp = _schema(all_supp or amazon_files)
    sm_vs   = _schema(vs_files) if vs_files else None

    if not sm_supp:
        sys.exit("[FATAL] Could not detect area column in suppression parquets.")

    print(f"     Suppression â€” area: {sm_supp.area}  year: {sm_supp.year}  "
          f"state: {sm_supp.state}  muni: {sm_supp.muni}  "
          f"class: {sm_supp.cls}  factor: {sm_supp.factor:.2e}")
    if sm_vs:
        print(f"     Secondary veg  â€” area: {sm_vs.area}  year: {sm_vs.year}  "
              f"class: {sm_vs.cls}  age: {sm_vs.age}")

    # â”€â”€ 3. Compute parameters â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    print("\n  [3/5] Computing parameters P1-P9 via DuckDB...")
    compute_timer = StageTimer("06_compute_parameters_p1_p9")

    p1 = p1_suppression_series(all_supp, sm_supp)
    print(f"     P1 rows: {len(p1)}")

    p2 = p2_cumulative(p1)
    print(f"     P2 rows: {len(p2)}")

    p3 = p3_nv_remaining(p2)
    print(f"     P3 rows: {len(p3)}")

    p4 = p4_nv_by_class(all_supp, sm_supp)
    print(f"     P4 rows: {len(p4)}")

    p5 = p5_sv_extent(vs_files, sm_vs) if sm_vs else []
    print(f"     P5 rows: {len(p5)}")

    p6 = p6_sv_increment(p5)
    print(f"     P6 rows: {len(p6)}")

    p7_pt = p7_sv_by_age_class(vs_files, sm_vs, "pt") if sm_vs else []
    p7_en = p7_sv_by_age_class(vs_files, sm_vs, "en") if sm_vs else []
    print(f"     P7 rows: {len(p7_pt)}")

    p8 = p8_sv_by_land_use(vs_files, sm_vs) if sm_vs else []
    print(f"     P8 rows: {len(p8)}")

    p9 = p9_muni_state_matrix(amazon_files or all_supp, sm_supp, "suppression")
    print(f"     P9 rows: {len(p9)}")
    row_counts = {
        "p1": len(p1),
        "p2": len(p2),
        "p3": len(p3),
        "p4": len(p4),
        "p5": len(p5),
        "p6": len(p6),
        "p7_pt": len(p7_pt),
        "p7_en": len(p7_en),
        "p8": len(p8),
        "p9": len(p9),
    }
    OBS_LOG.emit(
        "stage_metrics",
        **to_jsonable(
            compute_timer.finish(
                "ok",
                input_row_count=input_quality["parquet_profile"].get("row_count"),
                output_row_count=sum(row_counts.values()),
            )
        ),
    )

    # â”€â”€ 4. Generate charts â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    print("\n  [4/5] Generating publication-quality charts...")

    # P1 national aggregate (year, km2 only)
    p1_nat = [(r[0], r[-1]) for r in p1
              if len(r) == 2 or (len(r) == 3 and not sm_supp.state)
              or (len(r) >= 2)]
    # Deduplicate to national total per year
    nat_dict: dict[int, float] = {}
    for r in p1:
        yr = int(r[0]); km2 = float(r[-1])
        nat_dict[yr] = nat_dict.get(yr, 0) + km2
    p1_nat = [(y, v) for y, v in sorted(nat_dict.items())]

    p5_nat_dict: dict[int, float] = {}
    for r in p5:
        yr = int(r[0]); km2 = float(r[-1])
        p5_nat_dict[yr] = p5_nat_dict.get(yr, 0) + km2
    p5_nat = [(y, v) for y, v in sorted(p5_nat_dict.items())]

    p6_nat_dict: dict[int, float] = {}
    for r in p6:
        yr = int(r[0]); inc = float(r[-1])
        p6_nat_dict[yr] = p6_nat_dict.get(yr, 0) + inc
    p6_nat = [(y, v, p6_nat_dict.get(y, 0)) for y, v in sorted(p5_nat_dict.items())]

    chart_paths: dict[str, Path | None] = {}
    for lang in ("pt", "en"):
        p7_lang = p7_pt if lang == "pt" else p7_en
        chart_paths[f"suppression_trend_{lang}"]     = chart_suppression_trend(p1_nat, lang)
        chart_paths[f"cumulative_suppression_{lang}"] = chart_cumulative_suppression(p2, lang)
        chart_paths[f"sv_dynamics_{lang}"]            = chart_sv_dynamics(p5_nat, lang) if p5_nat else None
        chart_paths[f"sv_increment_{lang}"]           = chart_sv_increment(p6_nat, lang) if p6_nat else None
        chart_paths[f"state_ranking_{lang}"]          = chart_state_ranking(p1, lang) if sm_supp.state else None
        chart_paths[f"sv_subclass_{lang}"]            = chart_sv_subclass(p7_lang, lang) if p7_lang else None
        print(f"     [{lang.upper()}] charts generated")

    # â”€â”€ 5. Export workbooks + PPTX â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    print("\n  [5/5] Writing Excel workbooks and PowerPoint presentations...")
    date_str = datetime.now().strftime("%Y-%m-%d")
    output_artifacts: list[Path] = []

    for lang in ("pt", "en"):
        p7_lang = p7_pt if lang == "pt" else p7_en
        wb = Workbook()
        wb.remove(wb.active)

        _ws_p1(wb, lang, p1)
        _ws_p2(wb, lang, p2)
        _ws_p3(wb, lang, p3)
        _ws_p4(wb, lang, p4)
        _ws_p5(wb, lang, p5)
        _ws_p6(wb, lang, p6)
        _ws_p7(wb, lang, p7_lang)
        _ws_p8(wb, lang, p8)
        _ws_p9(wb, lang, p9)
        _ws_methodology(wb, lang)

        ll  = "PT" if lang == "pt" else "EN"
        out = TABLES_DIR / f"PRODES_Analytics_{ll}_{date_str}.xlsx"
        wb.save(str(out))
        output_artifacts.append(out)
        print(f"     [{ll}] Excel saved: {out}")

        data = {"p9": p9}
        prs  = build_pptx(lang, chart_paths, data)
        ppt_out = TABLES_DIR / f"PRODES_Analytics_{ll}_{date_str}.pptx"
        prs.save(str(ppt_out))
        output_artifacts.append(ppt_out)
        print(f"     [{ll}] PPTX  saved: {ppt_out}")

    artifacts = validate_nonempty_files(output_artifacts, "analytics export")
    report_path = write_run_report(
        REPORT_DIR,
        Path(__file__).name,
        {
            "status": "ok",
            "version": __version__,
            "geoparquet_dir": str(gpq_dir),
            "tables_dir": str(TABLES_DIR),
            "input_quality": input_quality,
            "input_counts": {
                "suppression_files": len(all_supp),
                "suppression_biomes": len(supp_by_biome),
                "secondary_vegetation_files": len(vs_files),
            },
            "row_counts": row_counts,
            "output_contract": to_jsonable(ANALYTICS_EXPORT_CONTRACT),
            "artifacts": artifacts,
            "lineage": LineageRecord(
                stage_name="06_export_tables",
                upstream_sources=[str(gpq_dir)],
                transformation="Compute P1-P9 analytic tables from GeoParquet with DuckDB and export bilingual Excel/PPTX artifacts.",
                downstream_outputs=[str(p) for p in output_artifacts],
                contracts=[GEOPARQUET_CONTRACT.name, ANALYTICS_EXPORT_CONTRACT.name],
            ),
        },
    )

    print(f"\n{DIV}")
    print(f"  Output folder : {TABLES_DIR}")
    print(f"  Charts folder : {CHART_DIR}")
    print(f"  Quality report: {report_path}")
    print(f"{SEP}\n")


if __name__ == "__main__":
    main()

